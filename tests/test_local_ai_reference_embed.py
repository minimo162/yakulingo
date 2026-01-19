from __future__ import annotations

import heapq
from pathlib import Path

from unittest.mock import Mock, patch

from yakulingo.services.translation_service import TranslationService
from yakulingo.ui.app import YakuLingoApp
from yakulingo.ui.state import TranslationBackend

from yakulingo.config.settings import AppSettings
from yakulingo.services.local_ai_prompt_builder import LocalPromptBuilder
from yakulingo.services.prompt_builder import PromptBuilder


def _make_builder() -> LocalPromptBuilder:
    repo_root = Path(__file__).resolve().parents[1]
    prompts_dir = repo_root / "prompts"
    return LocalPromptBuilder(
        prompts_dir,
        base_prompt_builder=PromptBuilder(prompts_dir),
        settings=AppSettings(),
    )


def _make_glossary_pairs(
    rows: list[tuple[str, str]],
) -> list[tuple[str, str, str, str, str, str]]:
    pairs: list[tuple[str, str, str, str, str, str]] = []
    for source, target in rows:
        source_folded = LocalPromptBuilder._normalize_for_glossary_match(source)
        target_folded = (
            LocalPromptBuilder._normalize_for_glossary_match(target) if target else ""
        )
        source_compact = LocalPromptBuilder._compact_for_glossary_match(source_folded)
        target_compact = LocalPromptBuilder._compact_for_glossary_match(target_folded)
        pairs.append(
            (
                source,
                target,
                source_folded,
                target_folded,
                source_compact,
                target_compact,
            )
        )
    return pairs


def _filter_glossary_pairs_linear(
    pairs: list[tuple[str, str, str, str, str, str]],
    input_text: str,
    *,
    max_lines: int,
) -> tuple[list[tuple[str, str]], bool]:
    text = (input_text or "").strip()
    if not text:
        return [], False

    text_folded = LocalPromptBuilder._normalize_for_glossary_match(text)
    text_compact = LocalPromptBuilder._compact_for_glossary_match(text_folded)

    seen: set[str] = set()
    heap: list[tuple[int, int, str, str]] = []
    matched_count = 0
    for idx, (
        source,
        target,
        source_folded,
        target_folded,
        source_compact,
        target_compact,
    ) in enumerate(pairs):
        source = (source or "").strip()
        if not source:
            continue
        if source in seen:
            continue

        matched = LocalPromptBuilder._matches_glossary_term(
            text_folded=text_folded,
            text_compact=text_compact,
            term_folded=source_folded,
            term_compact=source_compact,
        )
        if not matched and target:
            matched = LocalPromptBuilder._matches_glossary_term(
                text_folded=text_folded,
                text_compact=text_compact,
                term_folded=target_folded,
                term_compact=target_compact,
            )
        if not matched:
            continue

        seen.add(source)
        matched_count += 1
        key = max(len(source_folded or source), len(target_folded or target))
        item = (key, -idx, source, target)
        if len(heap) < max_lines:
            heapq.heappush(heap, item)
        else:
            if item > heap[0]:
                heapq.heapreplace(heap, item)

    if not heap:
        return [], False

    selected = sorted(heap, key=lambda x: (-x[0], -x[1]))
    return [(source, target) for _, _, source, target in selected], matched_count > max_lines


def test_local_reference_embed_supports_binary_formats(tmp_path: Path) -> None:
    builder = _make_builder()
    cases: list[tuple[Path, str]] = []

    from docx import Document

    docx_path = tmp_path / "ref.docx"
    doc = Document()
    doc.add_paragraph("Docx content")
    doc.save(docx_path)
    cases.append((docx_path, "Docx content"))

    import openpyxl

    xlsx_path = tmp_path / "ref.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Xlsx content"
    wb.save(xlsx_path)
    wb.close()
    cases.append((xlsx_path, "Xlsx content"))

    from pptx import Presentation
    from pptx.util import Inches

    pptx_path = tmp_path / "ref.pptx"
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    textbox.text_frame.text = "Pptx content"
    pres.save(pptx_path)
    cases.append((pptx_path, "Pptx content"))

    import fitz

    pdf_path = tmp_path / "ref.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "Pdf content")
    pdf.save(pdf_path)
    pdf.close()
    cases.append((pdf_path, "Pdf content"))

    for path, expected in cases:
        embedded = builder.build_reference_embed([path], input_text="sample")
        assert expected in embedded.text
        assert not any("未対応の参照ファイル" in w for w in embedded.warnings)


def test_local_prompt_includes_full_rules_for_short_text() -> None:
    builder = _make_builder()
    prompt = builder.build_text_to_en_single(
        "短文",
        style="minimal",
        reference_files=None,
        detected_language="日本語",
    )
    full_rules = builder._get_translation_rules("en").strip()
    assert full_rules
    assert full_rules in prompt


def test_local_prompt_includes_numeric_hints_for_oku() -> None:
    builder = _make_builder()
    prompt = builder.build_text_to_en_single(
        "売上高は2兆2,385億円(前年同期比1,554億円減)となりました。",
        style="minimal",
        reference_files=None,
        detected_language="日本語",
    )
    expected_rules = builder._get_translation_rules("en").strip()
    assert expected_rules
    assert expected_rules in prompt
    assert "数値変換ヒント" in prompt
    assert "2兆2,385億円 -> 22,385 oku yen" in prompt
    assert "1,554億円 -> 1,554 oku yen" in prompt


def test_local_prompt_includes_numeric_hints_for_oku_in_en_3style() -> None:
    builder = _make_builder()
    prompt = builder.build_text_to_en_single(
        "売上高は2兆2,385億円(前年同期比1,554億円減)となりました。",
        style="minimal",
        reference_files=None,
        detected_language="日本語",
    )
    expected_rules = builder._get_translation_rules("en").strip()
    assert expected_rules
    assert expected_rules in prompt
    assert "数値変換ヒント" in prompt
    assert "2兆2,385億円 -> 22,385 oku yen" in prompt
    assert "1,554億円 -> 1,554 oku yen" in prompt


def test_local_prompt_includes_numeric_hints_for_oku_in_en_missing_styles() -> None:
    builder = _make_builder()
    prompt = builder.build_text_to_en_single(
        "売上高は2兆2,385億円(前年同期比1,554億円減)となりました。",
        style="minimal",
        reference_files=None,
        detected_language="日本語",
    )
    expected_rules = builder._get_translation_rules("en").strip()
    assert expected_rules
    assert expected_rules in prompt
    assert "数値変換ヒント" in prompt
    assert "2兆2,385億円 -> 22,385 oku yen" in prompt
    assert "1,554億円 -> 1,554 oku yen" in prompt


def test_local_batch_prompt_includes_numeric_hints_for_oku() -> None:
    builder = _make_builder()
    prompt = builder.build_batch(
        [
            "売上高は2兆2,385億円となりました。",
            "前年同期比1,554億円減となりました。",
        ],
        output_language="en",
        translation_style="concise",
        reference_files=None,
    )
    expected_rules = builder._get_translation_rules("en").strip()
    assert expected_rules
    assert expected_rules in prompt
    assert "数値変換ヒント" in prompt
    assert "2兆2,385億円 -> 22,385 oku yen" in prompt
    assert "1,554億円 -> 1,554 oku yen" in prompt


def test_local_prompt_includes_full_rules_for_en_3style_short_text() -> None:
    builder = _make_builder()
    prompt = builder.build_text_to_en_single(
        "短文",
        style="minimal",
        reference_files=None,
        detected_language="日本語",
    )
    full_rules = builder._get_translation_rules("en").strip()
    assert full_rules
    assert full_rules in prompt


def test_local_prompt_includes_full_rules_for_en_to_jp_numeric_text() -> None:
    builder = _make_builder()
    text = "Revenue was 220k yen."
    prompt = builder.build_text_to_jp(
        text,
        reference_files=None,
        detected_language="英語",
    )
    expected_rules = builder._get_translation_rules("jp").strip()
    assert expected_rules
    assert expected_rules in prompt


def test_local_prompt_includes_full_rules_for_plain_en_to_jp_text() -> None:
    builder = _make_builder()
    prompt = builder.build_text_to_jp(
        "Hello world.",
        reference_files=None,
        detected_language="英語",
    )
    full_rules = builder._get_translation_rules("jp").strip()
    assert full_rules
    assert full_rules in prompt


def test_local_reference_embed_filters_bundled_glossary(tmp_path: Path) -> None:
    builder = _make_builder()
    glossary_path = tmp_path / "glossary.csv"
    glossary_path.write_text(
        "営業利益,Operating Profit\n売上高,Revenue\n", encoding="utf-8"
    )
    embedded = builder.build_reference_embed(
        [glossary_path], input_text="営業利益が増加"
    )
    assert "営業利益 翻译成 Operating Profit" in embedded.text
    assert "売上高 翻译成 Revenue" not in embedded.text


def test_local_reference_embed_matches_glossary_source_with_fullwidth_ascii(
    tmp_path: Path,
) -> None:
    builder = _make_builder()
    glossary_path = tmp_path / "glossary.csv"
    glossary_path.write_text("B/E台数,B/E Vol.\n", encoding="utf-8")
    embedded = builder.build_reference_embed(
        [glossary_path], input_text="Ｂ／Ｅ台数が増加"
    )
    assert "B/E台数 翻译成 B/E Vol." in embedded.text


def test_local_reference_embed_matches_glossary_target_with_hyphenated_phrase(
    tmp_path: Path,
) -> None:
    builder = _make_builder()
    glossary_path = tmp_path / "glossary.csv"
    glossary_path.write_text(
        "営業利益,Operating Profit\n売上高,Revenue\n", encoding="utf-8"
    )
    embedded = builder.build_reference_embed(
        [glossary_path], input_text="Operating-Profit improved."
    )
    assert "営業利益 翻译成 Operating Profit" in embedded.text
    assert "売上高 翻译成 Revenue" not in embedded.text


def test_local_reference_embed_truncates_bundled_glossary_to_max_lines(
    tmp_path: Path,
) -> None:
    builder = _make_builder()
    glossary_path = tmp_path / "glossary.csv"
    rows: list[str] = []
    tokens: list[str] = []
    for idx in range(100):
        source = f"TERM{idx:03d}"
        target = f"T{idx:03d}"
        rows.append(f"{source},{target}")
        tokens.append(source)
    glossary_path.write_text("\n".join(rows) + "\n", encoding="utf-8")

    embedded = builder.build_reference_embed(
        [glossary_path], input_text=" ".join(tokens)
    )
    assert embedded.truncated is True
    assert embedded.text.count(" 翻译成 ") == 80


def test_local_glossary_filter_indexed_matches_linear_for_same_input() -> None:
    pairs = _make_glossary_pairs(
        [
            ("foobar", "X"),
            ("bar", "BAR"),
            ("foobarish", "Y"),
            ("営業利益", "Operating Profit"),
            ("profit", "利益"),
            ("foobar", "X2"),
            ("abcd1234", "A"),
            ("abcd", "B"),
        ]
    )
    glossary = LocalPromptBuilder._build_glossary_index(pairs)

    cases: list[tuple[str, int]] = [
        ("foo_bar", 10),
        ("foobarish", 10),
        ("営業利益が増加", 10),
        ("profit improved", 10),
        ("profit利益", 10),
        ("abcd1234 abcd", 1),
        ("abcd1234 abcd", 2),
    ]
    for input_text, max_lines in cases:
        assert LocalPromptBuilder._filter_glossary_pairs(
            glossary, input_text, max_lines=max_lines
        ) == _filter_glossary_pairs_linear(pairs, input_text, max_lines=max_lines)


def test_local_glossary_filter_indexed_matches_linear_when_no_candidates() -> None:
    pairs = _make_glossary_pairs(
        [
            ("alpha", "A"),
            ("beta", "B"),
        ]
    )
    glossary = LocalPromptBuilder._build_glossary_index(pairs)

    input_text = "gamma"
    assert LocalPromptBuilder._filter_glossary_pairs(
        glossary, input_text, max_lines=10
    ) == _filter_glossary_pairs_linear(pairs, input_text, max_lines=10)


def _make_temp_builder(
    tmp_path: Path, *, use_bundled_glossary: bool = True
) -> LocalPromptBuilder:
    prompts_dir = tmp_path / "prompts"
    prompts_dir.mkdir(parents=True, exist_ok=True)

    (prompts_dir / "translation_rules.txt").write_text("RULES_MARKER", encoding="utf-8")
    (prompts_dir / "local_text_translate_to_en_single_json.txt").write_text(
        "{translation_rules}\n{reference_section}\n{input_text}\n", encoding="utf-8"
    )
    (prompts_dir / "local_batch_translate_to_en_json.txt").write_text(
        "{translation_rules}\n{reference_section}\n{items_json}\n", encoding="utf-8"
    )

    settings = AppSettings()
    settings.use_bundled_glossary = use_bundled_glossary
    return LocalPromptBuilder(
        prompts_dir,
        base_prompt_builder=PromptBuilder(prompts_dir),
        settings=settings,
    )


def test_local_reference_embed_auto_includes_bundled_glossary_when_enabled(
    tmp_path: Path,
) -> None:
    (tmp_path / "glossary.csv").write_text(
        "営業利益,Operating Profit\n売上高,Revenue\n", encoding="utf-8"
    )
    builder = _make_temp_builder(tmp_path, use_bundled_glossary=True)

    embedded = builder.build_reference_embed(None, input_text="営業利益が増加")
    assert "[REFERENCE:file=glossary.csv]" in embedded.text
    assert "営業利益 翻译成 Operating Profit" in embedded.text
    assert "売上高" not in embedded.text


def test_local_reference_embed_does_not_auto_include_bundled_glossary_when_disabled(
    tmp_path: Path,
) -> None:
    (tmp_path / "glossary.csv").write_text(
        "営業利益,Operating Profit\n売上高,Revenue\n", encoding="utf-8"
    )
    builder = _make_temp_builder(tmp_path, use_bundled_glossary=False)

    embedded = builder.build_reference_embed(None, input_text="営業利益が増加")
    assert embedded.text == ""


def test_local_text_prompt_includes_rules_and_bundled_glossary_when_enabled(
    tmp_path: Path,
) -> None:
    (tmp_path / "glossary.csv").write_text(
        "営業利益,Operating Profit\n売上高,Revenue\n", encoding="utf-8"
    )
    builder = _make_temp_builder(tmp_path, use_bundled_glossary=True)

    prompt = builder.build_text_to_en_single(
        "営業利益が増加",
        style="concise",
        reference_files=None,
        detected_language="日本語",
    )
    assert "RULES_MARKER" in prompt
    assert "営業利益 翻译成 Operating Profit" in prompt


def test_local_batch_prompt_includes_rules_and_bundled_glossary_when_enabled(
    tmp_path: Path,
) -> None:
    (tmp_path / "glossary.csv").write_text(
        "営業利益,Operating Profit\n売上高,Revenue\n", encoding="utf-8"
    )
    builder = _make_temp_builder(tmp_path, use_bundled_glossary=True)

    prompt = builder.build_batch(
        ["営業利益が増加"],
        output_language="en",
        translation_style="concise",
        reference_files=None,
    )
    assert "RULES_MARKER" in prompt
    assert "営業利益 翻译成 Operating Profit" in prompt


def test_local_followup_prompt_includes_rules_and_bundled_glossary_when_enabled(
    tmp_path: Path,
) -> None:
    (tmp_path / "glossary.csv").write_text(
        "営業利益,Operating Profit\n売上高,Revenue\n", encoding="utf-8"
    )
    builder = _make_temp_builder(tmp_path, use_bundled_glossary=True)

    prompt = builder.build_text_to_en_single(
        "営業利益が増加",
        style="minimal",
        reference_files=None,
        detected_language="日本語",
    )
    assert "RULES_MARKER" in prompt
    assert "営業利益 翻译成 Operating Profit" in prompt


def test_local_reference_embed_truncates_large_file(tmp_path: Path) -> None:
    builder = _make_builder()
    ref_path = tmp_path / "ref.txt"
    ref_path.write_text("A" * 2100, encoding="utf-8")
    embedded = builder.build_reference_embed([ref_path], input_text="sample")
    assert embedded.truncated is True
    assert any("上限 2000 文字" in w for w in embedded.warnings)


def test_local_reference_embed_truncates_total_limit(tmp_path: Path) -> None:
    builder = _make_builder()
    path_a = tmp_path / "a.txt"
    path_b = tmp_path / "b.txt"
    path_c = tmp_path / "c.txt"
    path_a.write_text("A" * 2000, encoding="utf-8")
    path_b.write_text("B" * 2000, encoding="utf-8")
    path_c.write_text("C" * 10, encoding="utf-8")
    embedded = builder.build_reference_embed(
        [path_a, path_b, path_c], input_text="sample"
    )
    assert embedded.truncated is True
    assert any("合計上限 4000 文字" in w for w in embedded.warnings)


def test_local_reference_embed_uses_file_cache_for_text(tmp_path: Path) -> None:
    builder = _make_builder()
    ref_path = tmp_path / "ref.txt"
    ref_path.write_text("Reference content", encoding="utf-8")

    original_read_text = Path.read_text
    with patch.object(Path, "read_text", autospec=True) as mock_read:
        mock_read.side_effect = lambda self, *args, **kwargs: original_read_text(
            self, *args, **kwargs
        )
        first = builder.build_reference_embed([ref_path], input_text="alpha")
        second = builder.build_reference_embed([ref_path], input_text="beta")

    assert "Reference content" in first.text
    assert "Reference content" in second.text
    assert mock_read.call_count == 2


def test_local_followup_reference_embed_includes_local_reference(
    tmp_path: Path,
) -> None:
    settings = AppSettings()
    settings.translation_backend = "local"
    app = YakuLingoApp()
    app.translation_service = TranslationService(
        Mock(),
        settings,
        prompts_dir=Path(__file__).resolve().parents[1] / "prompts",
    )
    app.state.translation_backend = TranslationBackend.LOCAL

    ref_path = tmp_path / "ref.txt"
    ref_path.write_text("Reference content", encoding="utf-8")

    prompt = app._build_follow_up_prompt(
        "review",
        source_text="source",
        translation="translation",
        reference_files=[ref_path],
    )
    assert prompt is not None
    assert "[REFERENCE:file=ref.txt]" in prompt
    assert "Reference content" in prompt


def test_local_prompt_includes_json_guard_block() -> None:
    builder = _make_builder()
    prompt = builder.build_text_to_en_single(
        "sample",
        style="concise",
        reference_files=None,
        detected_language="Japanese",
    )
    assert '{"translation":"..."}' in prompt


def test_local_batch_embed_reference_even_when_flag_false(tmp_path: Path) -> None:
    builder = _make_builder()
    ref_path = tmp_path / "ref.txt"
    ref_path.write_text("Reference content", encoding="utf-8")
    prompt = builder.build_batch(
        ["sample"],
        has_reference_files=False,
        output_language="en",
        translation_style="concise",
        reference_files=[ref_path],
    )
    assert "[REFERENCE:file=ref.txt]" in prompt
