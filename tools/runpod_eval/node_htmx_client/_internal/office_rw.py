#!/usr/bin/env python3
from __future__ import annotations

import json
import os
import posixpath
import re
import sys
import zipfile
import xml.etree.ElementTree as ET
from typing import Any

XML_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "pr": "http://schemas.openxmlformats.org/package/2006/relationships",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "s": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
}
INVALID_SHEET_CHARS_RE = re.compile(r"[:\\/?*\[\]]")


def _emit(payload: dict[str, Any]) -> None:
    sys.stdout.write(json.dumps(payload, ensure_ascii=False))


def _fail(message: str, exit_code: int = 1) -> None:
    _emit({"ok": False, "error": str(message)})
    raise SystemExit(exit_code)


def _clamp_int(raw: str, default: int, minimum: int, maximum: int) -> int:
    try:
        value = int(str(raw).strip())
    except Exception:
        return default
    return max(minimum, min(maximum, value))


def _safe_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value
    if isinstance(value, (int, float, bool)):
        return str(value)
    return json.dumps(value, ensure_ascii=False)


def _truncate_text(text: str, max_chars: int) -> tuple[str, bool]:
    if len(text) <= max_chars:
        return text, False
    return text[:max_chars], True


def _extract_docx_text(file_path: str, max_chars: int, max_items: int) -> dict[str, Any]:
    with zipfile.ZipFile(file_path) as zf:
        if "word/document.xml" not in zf.namelist():
            raise RuntimeError("word/document.xml not found.")
        root = ET.fromstring(zf.read("word/document.xml"))
        paragraphs: list[str] = []
        item_count = 0
        item_limit_reached = False
        for node in root.findall(".//w:p", XML_NS):
            text = "".join(t.text or "" for t in node.findall(".//w:t", XML_NS)).strip()
            if not text:
                continue
            paragraphs.append(text)
            item_count += 1
            if item_count >= max_items:
                item_limit_reached = True
                break

    joined = "\n".join(paragraphs).strip()
    truncated_text, text_limit_reached = _truncate_text(joined, max_chars)
    return {
        "format": "docx",
        "parser": "openxml-zip",
        "paragraphs": item_count,
        "truncated": text_limit_reached or item_limit_reached,
        "text": truncated_text,
    }


def _slide_sort_key(name: str) -> tuple[int, str]:
    match = re.search(r"slide(\d+)\.xml$", name)
    if match:
        return int(match.group(1)), name
    return 10**9, name


def _extract_pptx_text(file_path: str, max_chars: int, max_items: int) -> dict[str, Any]:
    with zipfile.ZipFile(file_path) as zf:
        slide_files = [
            name for name in zf.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]
        slide_files.sort(key=_slide_sort_key)
        lines: list[str] = []
        text_runs = 0
        item_limit_reached = False
        for slide_index, slide_path in enumerate(slide_files, start=1):
            root = ET.fromstring(zf.read(slide_path))
            chunks = [
                (node.text or "").strip()
                for node in root.findall(".//a:t", XML_NS)
                if (node.text or "").strip()
            ]
            if not chunks:
                continue
            lines.append(f"--- slide {slide_index} ---")
            for chunk in chunks:
                lines.append(chunk)
                text_runs += 1
                if text_runs >= max_items:
                    item_limit_reached = True
                    break
            if item_limit_reached:
                break

    joined = "\n".join(lines).strip()
    truncated_text, text_limit_reached = _truncate_text(joined, max_chars)
    return {
        "format": "pptx",
        "parser": "openxml-zip",
        "slides": len(slide_files),
        "text_runs": text_runs,
        "truncated": text_limit_reached or item_limit_reached,
        "text": truncated_text,
    }


def _extract_shared_strings(zf: zipfile.ZipFile) -> list[str]:
    path = "xl/sharedStrings.xml"
    if path not in zf.namelist():
        return []
    root = ET.fromstring(zf.read(path))
    output: list[str] = []
    for node in root.findall(".//s:si", XML_NS):
        output.append("".join(node.itertext()))
    return output


def _load_sheet_targets(zf: zipfile.ZipFile) -> list[tuple[str, str]]:
    rel_map: dict[str, str] = {}
    rels_path = "xl/_rels/workbook.xml.rels"
    if rels_path in zf.namelist():
        rel_root = ET.fromstring(zf.read(rels_path))
        for rel in rel_root.findall(".//pr:Relationship", XML_NS):
            rel_id = rel.attrib.get("Id", "").strip()
            target = rel.attrib.get("Target", "").strip()
            if not rel_id or not target:
                continue
            normalized = posixpath.normpath(posixpath.join("xl", target.lstrip("/")))
            rel_map[rel_id] = normalized

    workbook_path = "xl/workbook.xml"
    if workbook_path not in zf.namelist():
        return []
    wb_root = ET.fromstring(zf.read(workbook_path))
    results: list[tuple[str, str]] = []
    for index, sheet in enumerate(wb_root.findall(".//s:sheet", XML_NS), start=1):
        name = sheet.attrib.get("name", "").strip() or f"Sheet{index}"
        rel_id = sheet.attrib.get(f"{{{XML_NS['r']}}}id", "").strip()
        target = rel_map.get(rel_id, "")
        if target and target in zf.namelist():
            results.append((name, target))
    return results


def _sheet_file_sort_key(name: str) -> tuple[int, str]:
    match = re.search(r"sheet(\d+)\.xml$", name)
    if match:
        return int(match.group(1)), name
    return 10**9, name


def _xlsx_cell_value(cell: ET.Element, shared_strings: list[str]) -> str:
    cell_type = cell.attrib.get("t", "").strip()
    if cell_type == "inlineStr":
        inline = cell.find("s:is", XML_NS)
        return "".join(inline.itertext()).strip() if inline is not None else ""

    value_node = cell.find("s:v", XML_NS)
    raw = (value_node.text or "").strip() if value_node is not None else ""
    if cell_type == "s" and raw:
        try:
            index = int(raw)
            if 0 <= index < len(shared_strings):
                return shared_strings[index]
        except Exception:
            pass
    if not raw:
        formula_node = cell.find("s:f", XML_NS)
        if formula_node is not None and (formula_node.text or "").strip():
            return f"={(formula_node.text or '').strip()}"
    if cell_type == "b":
        if raw == "1":
            return "TRUE"
        if raw == "0":
            return "FALSE"
    return raw


def _extract_xlsx_text(file_path: str, max_chars: int, max_items: int) -> dict[str, Any]:
    with zipfile.ZipFile(file_path) as zf:
        shared_strings = _extract_shared_strings(zf)
        sheets = _load_sheet_targets(zf)
        if not sheets:
            sheet_files = [
                name for name in zf.namelist()
                if name.startswith("xl/worksheets/sheet") and name.endswith(".xml")
            ]
            sheet_files.sort(key=_sheet_file_sort_key)
            sheets = [(f"Sheet{i + 1}", name) for i, name in enumerate(sheet_files)]

        lines: list[str] = []
        filled_cells = 0
        item_limit_reached = False

        for sheet_index, (sheet_name, sheet_path) in enumerate(sheets, start=1):
            if sheet_path not in zf.namelist():
                continue
            root = ET.fromstring(zf.read(sheet_path))
            lines.append(f"[Sheet] {sheet_name}")
            has_values = False
            for cell in root.findall(".//s:sheetData/s:row/s:c", XML_NS):
                value = _xlsx_cell_value(cell, shared_strings).strip()
                if not value:
                    continue
                has_values = True
                ref = cell.attrib.get("r", "").strip() or f"R{sheet_index}C{filled_cells + 1}"
                lines.append(f"{ref}\t{value}")
                filled_cells += 1
                if filled_cells >= max_items:
                    item_limit_reached = True
                    break
            if not has_values:
                lines.append("(No non-empty cells)")
            if item_limit_reached:
                break
            lines.append("")

    joined = "\n".join(lines).strip()
    truncated_text, text_limit_reached = _truncate_text(joined, max_chars)
    return {
        "format": "xlsx",
        "parser": "openxml-zip",
        "sheets": len(sheets),
        "cells": filled_cells,
        "truncated": text_limit_reached or item_limit_reached,
        "text": truncated_text,
    }


def _parse_json_if_possible(raw_text: str) -> Any | None:
    stripped = raw_text.strip()
    if not stripped:
        return None
    if stripped[0] not in "{[":
        return None
    try:
        return json.loads(stripped)
    except Exception:
        return None


def _sanitize_sheet_name(raw_name: Any, index: int, used: set[str]) -> str:
    base = _safe_text(raw_name).strip() or f"Sheet{index}"
    base = INVALID_SHEET_CHARS_RE.sub("_", base).strip()
    if not base:
        base = f"Sheet{index}"
    base = base[:31]
    name = base
    suffix_index = 2
    while name in used:
        suffix = f"_{suffix_index}"
        name = f"{base[: max(1, 31 - len(suffix))]}{suffix}"
        suffix_index += 1
    used.add(name)
    return name


def _normalize_xlsx_sheet_specs(parsed: Any, raw_text: str) -> list[dict[str, Any]]:
    if isinstance(parsed, dict):
        sheets = parsed.get("sheets")
        if isinstance(sheets, list) and sheets:
            output = []
            for i, item in enumerate(sheets, start=1):
                if isinstance(item, dict):
                    output.append({
                        "name": item.get("name") or item.get("sheet") or f"Sheet{i}",
                        "rows": item.get("rows") if isinstance(item.get("rows"), list) else [],
                        "cells": item.get("cells") if isinstance(item.get("cells"), list) else [],
                    })
                else:
                    output.append({
                        "name": f"Sheet{i}",
                        "rows": [[item]],
                        "cells": [],
                    })
            return output
        if isinstance(parsed.get("rows"), list) or isinstance(parsed.get("cells"), list):
            return [{
                "name": parsed.get("name") or parsed.get("sheet") or "Sheet1",
                "rows": parsed.get("rows") if isinstance(parsed.get("rows"), list) else [],
                "cells": parsed.get("cells") if isinstance(parsed.get("cells"), list) else [],
            }]
        if isinstance(parsed.get("text"), str):
            rows = [[line] for line in parsed.get("text", "").splitlines()]
            return [{"name": parsed.get("name") or "Sheet1", "rows": rows, "cells": []}]

    if isinstance(parsed, list) and parsed:
        if any(isinstance(item, dict) and ("rows" in item or "cells" in item or "name" in item) for item in parsed):
            output = []
            for i, item in enumerate(parsed, start=1):
                if isinstance(item, dict):
                    output.append({
                        "name": item.get("name") or item.get("sheet") or f"Sheet{i}",
                        "rows": item.get("rows") if isinstance(item.get("rows"), list) else [],
                        "cells": item.get("cells") if isinstance(item.get("cells"), list) else [],
                    })
                else:
                    output.append({"name": f"Sheet{i}", "rows": [[item]], "cells": []})
            return output
        rows = []
        for row in parsed:
            if isinstance(row, list):
                rows.append(row)
            else:
                rows.append([row])
        return [{"name": "Sheet1", "rows": rows, "cells": []}]

    plain_lines = raw_text.splitlines()
    if plain_lines:
        rows = [[line] for line in plain_lines]
    elif raw_text:
        rows = [[raw_text]]
    else:
        rows = []
    return [{"name": "Sheet1", "rows": rows, "cells": []}]


def _xlsx_value(value: Any) -> Any:
    if value is None:
        return ""
    if isinstance(value, (int, float, bool, str)):
        return value
    return _safe_text(value)


def _write_xlsx(file_path: str, parsed: Any, raw_text: str) -> dict[str, Any]:
    try:
        from openpyxl import Workbook
    except Exception as exc:  # pragma: no cover - runtime dependency error path
        raise RuntimeError(f"openpyxl import failed: {exc}") from exc

    specs = _normalize_xlsx_sheet_specs(parsed, raw_text)
    if not specs:
        specs = [{"name": "Sheet1", "rows": [], "cells": []}]

    wb = Workbook()
    default_sheet = wb.active
    used_names: set[str] = set()
    written_cells = 0

    for index, spec in enumerate(specs, start=1):
        name = _sanitize_sheet_name(spec.get("name"), index, used_names)
        if index == 1:
            ws = default_sheet
            ws.title = name
        else:
            ws = wb.create_sheet(title=name)

        rows = spec.get("rows")
        if not isinstance(rows, list):
            rows = []
        for row_index, row in enumerate(rows, start=1):
            values: list[Any]
            if isinstance(row, dict):
                if "cell" in row:
                    cell_ref = _safe_text(row.get("cell")).strip()
                    if cell_ref:
                        ws[cell_ref] = _xlsx_value(row.get("value"))
                        written_cells += 1
                    continue
                values = list(row.values())
            elif isinstance(row, (list, tuple)):
                values = list(row)
            else:
                values = [row]
            for col_index, value in enumerate(values, start=1):
                ws.cell(row=row_index, column=col_index, value=_xlsx_value(value))
                written_cells += 1

        cells = spec.get("cells")
        if isinstance(cells, list):
            for cell in cells:
                if not isinstance(cell, dict):
                    continue
                cell_ref = _safe_text(cell.get("cell")).strip()
                if not cell_ref:
                    continue
                ws[cell_ref] = _xlsx_value(cell.get("value"))
                written_cells += 1

    wb.save(file_path)
    return {
        "format": "xlsx",
        "sheets": len(specs),
        "cells": written_cells,
    }


def _normalize_docx_paragraphs(parsed: Any, raw_text: str) -> list[str]:
    if isinstance(parsed, dict):
        paragraphs = parsed.get("paragraphs")
        if isinstance(paragraphs, list):
            return [_safe_text(item) for item in paragraphs]
        if isinstance(parsed.get("text"), str):
            return parsed["text"].splitlines()
    if isinstance(parsed, list):
        return [_safe_text(item) for item in parsed]
    lines = raw_text.splitlines()
    if lines:
        return lines
    if raw_text:
        return [raw_text]
    return []


def _write_docx(file_path: str, parsed: Any, raw_text: str) -> dict[str, Any]:
    try:
        from docx import Document
    except Exception as exc:  # pragma: no cover - runtime dependency error path
        raise RuntimeError(f"python-docx import failed: {exc}") from exc

    paragraphs = _normalize_docx_paragraphs(parsed, raw_text)
    doc = Document()
    for paragraph in paragraphs:
        doc.add_paragraph(paragraph)
    doc.save(file_path)
    return {
        "format": "docx",
        "paragraphs": len(paragraphs),
    }


def _slide_from_text(text: Any, default_title: str) -> dict[str, Any]:
    lines = [line.strip() for line in _safe_text(text).splitlines()]
    lines = [line for line in lines if line]
    if not lines:
        return {"title": default_title, "lines": []}
    return {
        "title": lines[0],
        "lines": lines[1:],
    }


def _normalize_slide_specs(parsed: Any, raw_text: str) -> list[dict[str, Any]]:
    if isinstance(parsed, dict):
        slides = parsed.get("slides")
        if isinstance(slides, list) and slides:
            output = []
            for i, item in enumerate(slides, start=1):
                default_title = f"Slide {i}"
                if isinstance(item, dict):
                    title = _safe_text(item.get("title")).strip() or default_title
                    lines_raw = (
                        item.get("lines")
                        if isinstance(item.get("lines"), list)
                        else (
                            item.get("bullets")
                            if isinstance(item.get("bullets"), list)
                            else (
                                item.get("text").splitlines()
                                if isinstance(item.get("text"), str)
                                else []
                            )
                        )
                    )
                    output.append({
                        "title": title,
                        "lines": [_safe_text(x) for x in lines_raw],
                    })
                else:
                    output.append(_slide_from_text(item, default_title))
            return output
        if isinstance(parsed.get("text"), str):
            return _normalize_slide_specs(None, parsed["text"])

    if isinstance(parsed, list) and parsed:
        output = []
        for i, item in enumerate(parsed, start=1):
            output.append(_slide_from_text(item, f"Slide {i}"))
        return output

    source = raw_text or ""
    blocks = [source]
    if "\f" in source:
        blocks = [block for block in source.split("\f")]
    elif "\n\n" in source:
        blocks = re.split(r"\n\s*\n+", source)

    output = []
    for i, block in enumerate(blocks, start=1):
        if not block.strip() and len(blocks) > 1:
            continue
        output.append(_slide_from_text(block, f"Slide {i}"))

    if not output:
        output = [{"title": "Slide 1", "lines": []}]
    return output


def _write_pptx(file_path: str, parsed: Any, raw_text: str) -> dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:  # pragma: no cover - runtime dependency error path
        raise RuntimeError(f"python-pptx import failed: {exc}") from exc

    slide_specs = _normalize_slide_specs(parsed, raw_text)
    prs = Presentation()

    for index, spec in enumerate(slide_specs, start=1):
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        title_text = _safe_text(spec.get("title")).strip() or f"Slide {index}"
        if slide.shapes.title is not None:
            slide.shapes.title.text = title_text

        body_frame = None
        for placeholder in slide.placeholders:
            if not getattr(placeholder, "has_text_frame", False):
                continue
            if slide.shapes.title is not None and placeholder == slide.shapes.title:
                continue
            body_frame = placeholder.text_frame
            break

        if body_frame is None:
            textbox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(8.0), Inches(4.5))
            body_frame = textbox.text_frame

        lines = spec.get("lines")
        if not isinstance(lines, list):
            lines = []
        body_frame.clear()
        if lines:
            body_frame.text = _safe_text(lines[0])
            for line in lines[1:]:
                paragraph = body_frame.add_paragraph()
                paragraph.text = _safe_text(line)

    prs.save(file_path)
    return {
        "format": "pptx",
        "slides": len(slide_specs),
        "text_lines": sum(len(spec.get("lines") or []) for spec in slide_specs),
    }


def _read_dispatch(file_path: str, max_chars: int, max_items: int) -> dict[str, Any]:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".xlsx":
        return _extract_xlsx_text(file_path, max_chars=max_chars, max_items=max_items)
    if ext == ".docx":
        return _extract_docx_text(file_path, max_chars=max_chars, max_items=max_items)
    if ext == ".pptx":
        return _extract_pptx_text(file_path, max_chars=max_chars, max_items=max_items)
    raise RuntimeError(f"Unsupported office extension: {ext}")


def _write_dispatch(file_path: str, raw_text: str) -> dict[str, Any]:
    ext = os.path.splitext(file_path)[1].lower()
    parsed = _parse_json_if_possible(raw_text)
    if ext == ".xlsx":
        return _write_xlsx(file_path, parsed=parsed, raw_text=raw_text)
    if ext == ".docx":
        return _write_docx(file_path, parsed=parsed, raw_text=raw_text)
    if ext == ".pptx":
        return _write_pptx(file_path, parsed=parsed, raw_text=raw_text)
    raise RuntimeError(f"Unsupported office extension: {ext}")


def main() -> None:
    if len(sys.argv) < 3:
        _fail("Usage: office_rw.py <read|write> <file_path> [max_chars] [max_items]")

    mode = str(sys.argv[1]).strip().lower()
    file_path = os.path.abspath(str(sys.argv[2]).strip())
    if not file_path:
        _fail("File path is empty.")

    try:
        if mode == "read":
            max_chars = _clamp_int(sys.argv[3] if len(sys.argv) >= 4 else "80000", 80000, 1000, 2_000_000)
            max_items = _clamp_int(sys.argv[4] if len(sys.argv) >= 5 else "4000", 4000, 1, 200_000)
            payload = _read_dispatch(file_path, max_chars=max_chars, max_items=max_items)
            payload["ok"] = True
            _emit(payload)
            return

        if mode == "write":
            raw_text = sys.stdin.read()
            os.makedirs(os.path.dirname(file_path) or ".", exist_ok=True)
            payload = _write_dispatch(file_path, raw_text=raw_text)
            payload["ok"] = True
            payload["bytes"] = os.path.getsize(file_path) if os.path.exists(file_path) else 0
            _emit(payload)
            return

        _fail(f"Unknown mode: {mode}")
    except SystemExit:
        raise
    except Exception as exc:
        message = str(exc)
        if message:
            print(message, file=sys.stderr)
        _fail(message or "office helper failed")


if __name__ == "__main__":
    main()
