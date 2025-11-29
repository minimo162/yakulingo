# ecm_translate/processors/pptx_processor.py
"""
Processor for PowerPoint files (.pptx, .ppt).
"""

from pathlib import Path
from typing import Iterator
from pptx import Presentation
from pptx.util import Pt

from .base import FileProcessor
from .translators import CellTranslator, ParagraphTranslator
from .font_manager import FontManager, FontTypeDetector
from ecm_translate.models.types import TextBlock, FileInfo, FileType


class PptxProcessor(FileProcessor):
    """
    Processor for PowerPoint files (.pptx, .ppt).

    Translation targets:
    - Shape text frames (ParagraphTranslator)
    - Table cells (CellTranslator - Excel-compatible)
    - Speaker notes (ParagraphTranslator)

    Preserved:
    - Slide layouts
    - Animations
    - Transitions
    - Images
    - Charts
    """

    def __init__(self):
        self.cell_translator = CellTranslator()
        self.para_translator = ParagraphTranslator()
        self.font_type_detector = FontTypeDetector()

    @property
    def file_type(self) -> FileType:
        return FileType.POWERPOINT

    @property
    def supported_extensions(self) -> list[str]:
        return ['.pptx', '.ppt']

    def get_file_info(self, file_path: Path) -> FileInfo:
        """Get PowerPoint file info"""
        prs = Presentation(file_path)

        slide_count = len(prs.slides)
        text_count = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                # Text shapes
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text and self.para_translator.should_translate(para.text):
                            text_count += 1

                # Tables (Excel-compatible)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = cell.text_frame.text if cell.text_frame else ""
                            if cell_text and self.cell_translator.should_translate(cell_text):
                                text_count += 1

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                for para in slide.notes_slide.notes_text_frame.paragraphs:
                    if para.text and self.para_translator.should_translate(para.text):
                        text_count += 1

        return FileInfo(
            path=file_path,
            file_type=FileType.POWERPOINT,
            size_bytes=file_path.stat().st_size,
            slide_count=slide_count,
            text_block_count=text_count,
        )

    def extract_text_blocks(self, file_path: Path) -> Iterator[TextBlock]:
        """Extract text from slides, shapes, tables, notes"""
        prs = Presentation(file_path)

        for slide_idx, slide in enumerate(prs.slides):
            shape_counter = 0
            table_counter = 0

            for shape in slide.shapes:
                # === Text Shapes ===
                if shape.has_text_frame:
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        if para.text and self.para_translator.should_translate(para.text):
                            # Get font info from first run
                            font_name = None
                            font_size = 18.0  # default for PPT
                            if para.runs:
                                first_run = para.runs[0]
                                if first_run.font.name:
                                    font_name = first_run.font.name
                                if first_run.font.size:
                                    font_size = first_run.font.size.pt

                            # Get dominant font if multiple runs
                            if len(para.runs) > 1:
                                font_names = [r.font.name for r in para.runs if r.font.name]
                                if font_names:
                                    font_name = self.font_type_detector.get_dominant_font(font_names)

                            yield TextBlock(
                                id=f"s{slide_idx}_sh{shape_counter}_p{para_idx}",
                                text=para.text,
                                location=f"Slide {slide_idx + 1}, Shape {shape_counter + 1}",
                                metadata={
                                    'type': 'shape',
                                    'slide': slide_idx,
                                    'shape': shape_counter,
                                    'para': para_idx,
                                    'font_name': font_name,
                                    'font_size': font_size,
                                }
                            )
                    shape_counter += 1

                # === Tables (Excel-compatible) ===
                if shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for cell_idx, cell in enumerate(row.cells):
                            cell_text = cell.text_frame.text if cell.text_frame else ""
                            if cell_text and self.cell_translator.should_translate(cell_text):
                                # Get font info
                                font_name = None
                                font_size = 14.0
                                if cell.text_frame and cell.text_frame.paragraphs:
                                    first_para = cell.text_frame.paragraphs[0]
                                    if first_para.runs:
                                        first_run = first_para.runs[0]
                                        if first_run.font.name:
                                            font_name = first_run.font.name
                                        if first_run.font.size:
                                            font_size = first_run.font.size.pt

                                yield TextBlock(
                                    id=f"s{slide_idx}_tbl{table_counter}_r{row_idx}_c{cell_idx}",
                                    text=cell_text,
                                    location=f"Slide {slide_idx + 1}, Table {table_counter + 1}, Row {row_idx + 1}, Cell {cell_idx + 1}",
                                    metadata={
                                        'type': 'table_cell',
                                        'slide': slide_idx,
                                        'table': table_counter,
                                        'row': row_idx,
                                        'col': cell_idx,
                                        'font_name': font_name,
                                        'font_size': font_size,
                                    }
                                )
                    table_counter += 1

            # === Speaker Notes ===
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_frame = slide.notes_slide.notes_text_frame
                for para_idx, para in enumerate(notes_frame.paragraphs):
                    if para.text and self.para_translator.should_translate(para.text):
                        yield TextBlock(
                            id=f"s{slide_idx}_notes_{para_idx}",
                            text=para.text,
                            location=f"Slide {slide_idx + 1}, Notes",
                            metadata={
                                'type': 'notes',
                                'slide': slide_idx,
                                'para': para_idx,
                            }
                        )

    def apply_translations(
        self,
        input_path: Path,
        output_path: Path,
        translations: dict[str, str],
        direction: str = "jp_to_en",
    ) -> None:
        """Apply translations to PowerPoint"""
        prs = Presentation(input_path)
        font_manager = FontManager(direction)

        for slide_idx, slide in enumerate(prs.slides):
            shape_counter = 0
            table_counter = 0

            for shape in slide.shapes:
                # === Apply to text shapes ===
                if shape.has_text_frame:
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        block_id = f"s{slide_idx}_sh{shape_counter}_p{para_idx}"
                        if block_id in translations:
                            self._apply_to_paragraph(para, translations[block_id], font_manager)
                    shape_counter += 1

                # === Apply to tables ===
                if shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for cell_idx, cell in enumerate(row.cells):
                            block_id = f"s{slide_idx}_tbl{table_counter}_r{row_idx}_c{cell_idx}"
                            if block_id in translations:
                                if cell.text_frame and cell.text_frame.paragraphs:
                                    self._apply_to_paragraph(
                                        cell.text_frame.paragraphs[0],
                                        translations[block_id],
                                        font_manager
                                    )
                                    # Clear remaining paragraphs
                                    for para in cell.text_frame.paragraphs[1:]:
                                        for run in para.runs:
                                            run.text = ""
                    table_counter += 1

            # === Apply to speaker notes ===
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_frame = slide.notes_slide.notes_text_frame
                for para_idx, para in enumerate(notes_frame.paragraphs):
                    block_id = f"s{slide_idx}_notes_{para_idx}"
                    if block_id in translations:
                        self._apply_to_paragraph(para, translations[block_id], font_manager)

        prs.save(output_path)

    def _apply_to_paragraph(self, para, translated_text: str, font_manager: FontManager) -> None:
        """
        Apply translation to paragraph, preserving paragraph style.

        Strategy:
        - Keep first run's formatting (font, size, color)
        - Set translated text to first run
        - Clear remaining runs
        """
        if para.runs:
            first_run = para.runs[0]

            # Get original font info
            original_font_name = first_run.font.name
            original_font_size = first_run.font.size.pt if first_run.font.size else 18.0

            # Get new font settings
            new_font_name, new_font_size = font_manager.select_font(
                original_font_name,
                original_font_size
            )

            # Apply translation
            first_run.text = translated_text

            # Apply new font
            first_run.font.name = new_font_name
            first_run.font.size = Pt(new_font_size)

            # Clear remaining runs
            for run in para.runs[1:]:
                run.text = ""
        else:
            # No runs - add text directly
            para.text = translated_text
