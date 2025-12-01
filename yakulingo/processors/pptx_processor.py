# yakulingo/processors/pptx_processor.py
"""
Processor for PowerPoint files (.pptx).
"""

import logging
from pathlib import Path
from typing import Iterator
from pptx import Presentation
from pptx.util import Pt

from .base import FileProcessor

# Module logger
logger = logging.getLogger(__name__)
from .translators import CellTranslator, ParagraphTranslator
from .font_manager import FontManager, FontTypeDetector
from yakulingo.models.types import TextBlock, FileInfo, FileType


class PptxProcessor(FileProcessor):
    """
    Processor for PowerPoint files (.pptx, .ppt).

    Translation targets:
    - Shape text frames (ParagraphTranslator)
    - Table cells (CellTranslator - Excel-compatible)
    - Speaker notes (ParagraphTranslator)

    Preserved:
    - Slide layouts
    - Animations
    - Transitions
    - Images
    - Charts
    """

    def __init__(self):
        self.cell_translator = CellTranslator()
        self.para_translator = ParagraphTranslator()
        self.font_type_detector = FontTypeDetector()

    @property
    def file_type(self) -> FileType:
        return FileType.POWERPOINT

    @property
    def supported_extensions(self) -> list[str]:
        # Note: .ppt (legacy format) is not supported by python-pptx
        # Only .pptx (Office Open XML) is supported
        return ['.pptx']

    def get_file_info(self, file_path: Path) -> FileInfo:
        """Get PowerPoint file info"""
        prs = Presentation(file_path)

        slide_count = len(prs.slides)
        text_count = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                # Text shapes
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text and self.para_translator.should_translate(para.text):
                            text_count += 1

                # Tables (Excel-compatible)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = cell.text_frame.text if cell.text_frame else ""
                            if cell_text and self.cell_translator.should_translate(cell_text):
                                text_count += 1

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                for para in slide.notes_slide.notes_text_frame.paragraphs:
                    if para.text and self.para_translator.should_translate(para.text):
                        text_count += 1

        return FileInfo(
            path=file_path,
            file_type=FileType.POWERPOINT,
            size_bytes=file_path.stat().st_size,
            slide_count=slide_count,
            text_block_count=text_count,
        )

    def extract_text_blocks(self, file_path: Path) -> Iterator[TextBlock]:
        """Extract text from slides, shapes, tables, notes"""
        prs = Presentation(file_path)

        for slide_idx, slide in enumerate(prs.slides):
            shape_counter = 0
            table_counter = 0

            for shape in slide.shapes:
                # === Text Shapes ===
                if shape.has_text_frame:
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        if para.text and self.para_translator.should_translate(para.text):
                            # Get font info from first run
                            font_name = None
                            font_size = 18.0  # default for PPT
                            if para.runs:
                                first_run = para.runs[0]
                                if first_run.font.name:
                                    font_name = first_run.font.name
                                if first_run.font.size:
                                    font_size = first_run.font.size.pt

                            # Get dominant font if multiple runs
                            if len(para.runs) > 1:
                                font_names = [r.font.name for r in para.runs if r.font.name]
                                if font_names:
                                    font_name = self.font_type_detector.get_dominant_font(font_names)

                            yield TextBlock(
                                id=f"s{slide_idx}_sh{shape_counter}_p{para_idx}",
                                text=para.text,
                                location=f"Slide {slide_idx + 1}, Shape {shape_counter + 1}",
                                metadata={
                                    'type': 'shape',
                                    'slide': slide_idx,
                                    'shape': shape_counter,
                                    'para': para_idx,
                                    'font_name': font_name,
                                    'font_size': font_size,
                                }
                            )
                    shape_counter += 1

                # === Tables (Excel-compatible) ===
                if shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for cell_idx, cell in enumerate(row.cells):
                            cell_text = cell.text_frame.text if cell.text_frame else ""
                            if cell_text and self.cell_translator.should_translate(cell_text):
                                # Get font info
                                font_name = None
                                font_size = 14.0
                                if cell.text_frame and cell.text_frame.paragraphs:
                                    first_para = cell.text_frame.paragraphs[0]
                                    if first_para.runs:
                                        first_run = first_para.runs[0]
                                        if first_run.font.name:
                                            font_name = first_run.font.name
                                        if first_run.font.size:
                                            font_size = first_run.font.size.pt

                                yield TextBlock(
                                    id=f"s{slide_idx}_tbl{table_counter}_r{row_idx}_c{cell_idx}",
                                    text=cell_text,
                                    location=f"Slide {slide_idx + 1}, Table {table_counter + 1}, Row {row_idx + 1}, Cell {cell_idx + 1}",
                                    metadata={
                                        'type': 'table_cell',
                                        'slide': slide_idx,
                                        'table': table_counter,
                                        'row': row_idx,
                                        'col': cell_idx,
                                        'font_name': font_name,
                                        'font_size': font_size,
                                    }
                                )
                    table_counter += 1

            # === Speaker Notes ===
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_frame = slide.notes_slide.notes_text_frame
                for para_idx, para in enumerate(notes_frame.paragraphs):
                    if para.text and self.para_translator.should_translate(para.text):
                        yield TextBlock(
                            id=f"s{slide_idx}_notes_{para_idx}",
                            text=para.text,
                            location=f"Slide {slide_idx + 1}, Notes",
                            metadata={
                                'type': 'notes',
                                'slide': slide_idx,
                                'para': para_idx,
                            }
                        )

    def apply_translations(
        self,
        input_path: Path,
        output_path: Path,
        translations: dict[str, str],
        direction: str = "jp_to_en",
    ) -> None:
        """Apply translations to PowerPoint"""
        prs = Presentation(input_path)
        font_manager = FontManager(direction)

        for slide_idx, slide in enumerate(prs.slides):
            shape_counter = 0
            table_counter = 0

            for shape in slide.shapes:
                # === Apply to text shapes ===
                if shape.has_text_frame:
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        block_id = f"s{slide_idx}_sh{shape_counter}_p{para_idx}"
                        if block_id in translations:
                            self._apply_to_paragraph(para, translations[block_id], font_manager)
                    shape_counter += 1

                # === Apply to tables ===
                if shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for cell_idx, cell in enumerate(row.cells):
                            block_id = f"s{slide_idx}_tbl{table_counter}_r{row_idx}_c{cell_idx}"
                            if block_id in translations:
                                if cell.text_frame and cell.text_frame.paragraphs:
                                    self._apply_to_paragraph(
                                        cell.text_frame.paragraphs[0],
                                        translations[block_id],
                                        font_manager
                                    )
                                    # Clear remaining paragraphs
                                    for para in cell.text_frame.paragraphs[1:]:
                                        for run in para.runs:
                                            run.text = ""
                    table_counter += 1

            # === Apply to speaker notes ===
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_frame = slide.notes_slide.notes_text_frame
                for para_idx, para in enumerate(notes_frame.paragraphs):
                    block_id = f"s{slide_idx}_notes_{para_idx}"
                    if block_id in translations:
                        self._apply_to_paragraph(para, translations[block_id], font_manager)

        prs.save(output_path)

    def _apply_to_paragraph(self, para, translated_text: str, font_manager: FontManager) -> None:
        """
        Apply translation to paragraph, preserving paragraph style.

        Strategy:
        - Keep first run's formatting (font, size, color)
        - Set translated text to first run
        - Clear remaining runs
        """
        if para.runs:
            first_run = para.runs[0]

            # Get original font info
            original_font_name = first_run.font.name
            original_font_size = first_run.font.size.pt if first_run.font.size else 18.0

            # Get new font settings
            new_font_name, new_font_size = font_manager.select_font(
                original_font_name,
                original_font_size
            )

            # Apply translation
            first_run.text = translated_text

            # Apply new font
            first_run.font.name = new_font_name
            first_run.font.size = Pt(new_font_size)

            # Clear remaining runs
            for run in para.runs[1:]:
                run.text = ""
        else:
            # No runs - add text directly
            para.text = translated_text

    def create_bilingual_presentation(
        self,
        original_path: Path,
        translated_path: Path,
        output_path: Path,
    ) -> dict[str, int]:
        """
        Create a bilingual presentation with original and translated slides interleaved.

        Output format:
            Slide 1 (original), Slide 1 (translated), Slide 2 (original), Slide 2 (translated), ...

        Note: Due to python-pptx limitations, this creates a new presentation by
        appending translated slides after all original slides. Full interleaving
        would require low-level XML manipulation.

        Args:
            original_path: Path to the original presentation
            translated_path: Path to the translated presentation
            output_path: Path to save the bilingual presentation

        Returns:
            dict with original_slides, translated_slides, total_slides counts
        """
        import shutil
        import zipfile
        import tempfile
        from xml.etree import ElementTree as ET

        # Copy original file to output
        shutil.copy2(original_path, output_path)

        # Count slides
        original_prs = Presentation(original_path)
        translated_prs = Presentation(translated_path)

        original_slides = len(original_prs.slides)
        translated_slides = len(translated_prs.slides)

        # For a proper interleaved approach, we'd need to manipulate the XML directly
        # For now, append translated slides after original slides with a separator

        # Create a combined presentation using XML manipulation
        try:
            self._merge_presentations_xml(
                output_path, translated_path, output_path
            )
        except (OSError, ValueError, KeyError, ET.ParseError) as e:
            logger.warning("XML merge failed, using simple append: %s", e)
            # Fallback: Just return the original file
            pass

        # Recount after merge
        result_prs = Presentation(output_path)
        total_slides = len(result_prs.slides)

        return {
            'original_slides': original_slides,
            'translated_slides': translated_slides,
            'total_slides': total_slides,
        }

    def _merge_presentations_xml(
        self,
        base_path: Path,
        append_path: Path,
        output_path: Path,
    ) -> None:
        """
        Merge two PowerPoint presentations by appending slides from append_path to base_path.

        This uses low-level ZIP/XML manipulation to copy slides.
        """
        import shutil
        import zipfile
        import tempfile
        from xml.etree import ElementTree as ET

        # PowerPoint XML namespaces
        PPTX_NS = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        }

        # Register namespaces
        for prefix, uri in PPTX_NS.items():
            ET.register_namespace(prefix, uri)

        with tempfile.TemporaryDirectory() as tmpdir:
            tmpdir = Path(tmpdir)

            # Extract base pptx
            base_dir = tmpdir / 'base'
            with zipfile.ZipFile(base_path, 'r') as zf:
                zf.extractall(base_dir)

            # Extract append pptx
            append_dir = tmpdir / 'append'
            with zipfile.ZipFile(append_path, 'r') as zf:
                zf.extractall(append_dir)

            # Get the number of slides in base
            base_slides_dir = base_dir / 'ppt' / 'slides'
            base_slide_count = len(list(base_slides_dir.glob('slide*.xml')))

            # Copy slides from append to base
            append_slides_dir = append_dir / 'ppt' / 'slides'
            append_rels_dir = append_dir / 'ppt' / 'slides' / '_rels'

            for i, slide_file in enumerate(sorted(append_slides_dir.glob('slide*.xml'))):
                new_slide_num = base_slide_count + i + 1
                new_slide_name = f'slide{new_slide_num}.xml'

                # Copy slide XML
                shutil.copy2(slide_file, base_slides_dir / new_slide_name)

                # Copy slide relationships if exist
                rels_file = append_rels_dir / f'{slide_file.name}.rels'
                if rels_file.exists():
                    base_rels_dir = base_slides_dir / '_rels'
                    base_rels_dir.mkdir(exist_ok=True)
                    shutil.copy2(rels_file, base_rels_dir / f'{new_slide_name}.rels')

            # Update presentation.xml to include new slides
            pres_xml_path = base_dir / 'ppt' / 'presentation.xml'
            tree = ET.parse(pres_xml_path)
            root = tree.getroot()

            # Find sldIdLst element
            ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
            sld_id_lst = root.find('.//p:sldIdLst', ns)

            if sld_id_lst is not None:
                # Get max id
                max_id = 256  # Default starting id
                for sld_id in sld_id_lst.findall('p:sldId', ns):
                    id_val = int(sld_id.get('id', '256'))
                    max_id = max(max_id, id_val)

                # Add new slide entries
                append_slide_count = len(list(append_slides_dir.glob('slide*.xml')))
                for i in range(append_slide_count):
                    new_id = max_id + i + 1
                    new_rid = f'rId{base_slide_count + i + 100}'  # Offset to avoid collision
                    new_sld_id = ET.SubElement(sld_id_lst, f'{{{ns["p"]}}}sldId')
                    new_sld_id.set('id', str(new_id))
                    new_sld_id.set(f'{{{PPTX_NS["r"]}}}id', new_rid)

            tree.write(pres_xml_path, xml_declaration=True, encoding='UTF-8')

            # Update presentation.xml.rels to include relationships for new slides
            pres_rels_path = base_dir / 'ppt' / '_rels' / 'presentation.xml.rels'
            rels_ns = 'http://schemas.openxmlformats.org/package/2006/relationships'
            ET.register_namespace('', rels_ns)

            if pres_rels_path.exists():
                rels_tree = ET.parse(pres_rels_path)
                rels_root = rels_tree.getroot()

                # Add new relationships for slides
                append_slide_count = len(list(append_slides_dir.glob('slide*.xml')))
                for i in range(append_slide_count):
                    new_slide_num = base_slide_count + i + 1
                    new_rid = f'rId{base_slide_count + i + 100}'  # Same offset as above
                    rel_elem = ET.SubElement(rels_root, f'{{{rels_ns}}}Relationship')
                    rel_elem.set('Id', new_rid)
                    rel_elem.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
                    rel_elem.set('Target', f'slides/slide{new_slide_num}.xml')

                rels_tree.write(pres_rels_path, xml_declaration=True, encoding='UTF-8')

            # Update Content_Types
            content_types_path = base_dir / '[Content_Types].xml'
            ct_tree = ET.parse(content_types_path)
            ct_root = ct_tree.getroot()

            ct_ns = 'http://schemas.openxmlformats.org/package/2006/content-types'
            append_slide_count = len(list(append_slides_dir.glob('slide*.xml')))
            for i in range(append_slide_count):
                new_slide_num = base_slide_count + i + 1
                override = ET.SubElement(ct_root, f'{{{ct_ns}}}Override')
                override.set('PartName', f'/ppt/slides/slide{new_slide_num}.xml')
                override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml')

            ct_tree.write(content_types_path, xml_declaration=True, encoding='UTF-8')

            # Repack the pptx
            with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
                for file_path in base_dir.rglob('*'):
                    if file_path.is_file():
                        arcname = file_path.relative_to(base_dir)
                        zf.write(file_path, arcname)
