# tests/test_translation_service.py
"""Tests for yakulingo.services.translation_service"""

import pytest
import tempfile
from pathlib import Path
from unittest.mock import Mock, MagicMock, patch

from yakulingo.models.types import (
    TranslationStatus,
    TextBlock,
    FileType,
)
from yakulingo.config.settings import AppSettings
from yakulingo.processors.pdf_processor import TranslationCell
from yakulingo.services.translation_service import (
    BatchTranslator,
    TranslationCache,
    TranslationService,
)
from yakulingo.services.prompt_builder import PromptBuilder


class TestBatchTranslatorCreateBatches:
    """Tests for BatchTranslator._create_batches()"""

    @pytest.fixture
    def batch_translator(self):
        """Create BatchTranslator with real PromptBuilder (copilot still mocked)"""
        mock_copilot = Mock()
        prompt_builder = PromptBuilder()  # Use real PromptBuilder
        return BatchTranslator(mock_copilot, prompt_builder)

    def test_empty_blocks(self, batch_translator):
        """Empty block list returns empty batches"""
        batches = batch_translator._create_batches([])
        assert batches == []

    def test_single_block(self, batch_translator):
        """Single block creates single batch"""
        blocks = [TextBlock(id="1", text="Hello", location="A1")]
        batches = batch_translator._create_batches(blocks)

        assert len(batches) == 1
        assert len(batches[0]) == 1
        assert batches[0][0].id == "1"

    def test_multiple_blocks_single_batch(self, batch_translator):
        """Multiple small blocks fit in one batch"""
        blocks = [
            TextBlock(id="1", text="Hello", location="A1"),
            TextBlock(id="2", text="World", location="A2"),
            TextBlock(id="3", text="Test", location="A3"),
        ]
        batches = batch_translator._create_batches(blocks)

        assert len(batches) == 1
        assert len(batches[0]) == 3

    def test_respects_max_chars_limit(self, batch_translator):
        """Batches respect MAX_CHARS_PER_BATCH limit"""
        # Create blocks with large text (6000 chars each)
        large_text = "x" * 6000
        blocks = [
            TextBlock(id="1", text=large_text, location="A1"),
            TextBlock(id="2", text=large_text, location="A2"),
            TextBlock(id="3", text=large_text, location="A3"),
        ]
        batches = batch_translator._create_batches(blocks)

        # 6000 + 6000 > 10000, so should split after first block
        assert len(batches) == 3

    def test_preserves_block_order(self, batch_translator):
        """Block order is preserved across batches when split by character limit"""
        # Create blocks with large text that will force splits
        # Each block is 4000 chars, so 2 blocks = 8000 > 7000 limit
        large_text = "x" * 4000
        blocks = [
            TextBlock(id=str(i), text=large_text, location=f"A{i}")
            for i in range(4)
        ]
        batches = batch_translator._create_batches(blocks)

        # Should create 4 batches (one per block due to char limit)
        assert len(batches) == 4

        # Order should be preserved
        assert batches[0][0].id == "0"
        assert batches[1][0].id == "1"
        assert batches[2][0].id == "2"
        assert batches[3][0].id == "3"


class TestFilterBlocksBySection:
    """Section filtering for partial translation."""

    @pytest.fixture
    def service(self):
        return TranslationService(Mock(), AppSettings())

    def test_empty_selection_skips_all_blocks(self, service):
        """No sections selected should result in no blocks being processed."""

        blocks = [
            TextBlock(id="sheet1_A1", text="Hello", location="Sheet1, A1", metadata={"sheet_idx": 0}),
            TextBlock(id="sheet2_A1", text="World", location="Sheet2, A1", metadata={"sheet_idx": 1}),
        ]

        filtered = service._filter_blocks_by_section(blocks, [])

        assert filtered == []

    def test_selected_section_included(self, service):
        """Only blocks from selected sections are kept."""

        blocks = [
            TextBlock(id="sheet1_A1", text="Hello", location="Sheet1, A1", metadata={"sheet_idx": 0}),
            TextBlock(id="sheet2_A1", text="World", location="Sheet2, A1", metadata={"sheet_idx": 1}),
        ]

        filtered = service._filter_blocks_by_section(blocks, [1])

        assert [block.id for block in filtered] == ["sheet2_A1"]


class TestTranslationServiceInit:
    """Tests for TranslationService initialization"""

    @pytest.fixture
    def mock_copilot(self):
        return Mock()

    @pytest.fixture
    def settings(self):
        return AppSettings()

    def test_registers_excel_processors(self, mock_copilot, settings):
        """Excel extensions are registered"""
        service = TranslationService(mock_copilot, settings)
        assert '.xlsx' in service.processors
        assert '.xls' in service.processors

    def test_registers_word_processors(self, mock_copilot, settings):
        """Word extensions are registered"""
        service = TranslationService(mock_copilot, settings)
        assert '.docx' in service.processors
        # .doc (legacy format) is not supported by python-docx
        assert '.doc' not in service.processors

    def test_registers_powerpoint_processors(self, mock_copilot, settings):
        """PowerPoint extensions are registered"""
        service = TranslationService(mock_copilot, settings)
        assert '.pptx' in service.processors
        # .ppt (legacy format) is not supported by python-pptx
        assert '.ppt' not in service.processors

    def test_registers_pdf_processor(self, mock_copilot, settings):
        """PDF extension is registered"""
        service = TranslationService(mock_copilot, settings)
        assert '.pdf' in service.processors


class TestTranslationServiceSupportedFiles:
    """Tests for TranslationService file support methods"""

    @pytest.fixture
    def service(self):
        return TranslationService(Mock(), AppSettings())

    def test_is_supported_file_xlsx(self, service):
        assert service.is_supported_file(Path("test.xlsx")) is True

    def test_is_supported_file_xls(self, service):
        assert service.is_supported_file(Path("test.xls")) is True

    def test_is_supported_file_docx(self, service):
        assert service.is_supported_file(Path("test.docx")) is True

    def test_is_supported_file_pptx(self, service):
        assert service.is_supported_file(Path("test.pptx")) is True

    def test_is_supported_file_pdf(self, service):
        assert service.is_supported_file(Path("test.pdf")) is True

    def test_is_supported_file_txt(self, service):
        assert service.is_supported_file(Path("test.txt")) is True

    def test_is_supported_file_unsupported(self, service):
        assert service.is_supported_file(Path("test.csv")) is False
        assert service.is_supported_file(Path("test.jpg")) is False
        assert service.is_supported_file(Path("test.xyz")) is False

    def test_is_supported_file_case_insensitive(self, service):
        assert service.is_supported_file(Path("test.XLSX")) is True
        assert service.is_supported_file(Path("test.Docx")) is True
        assert service.is_supported_file(Path("test.PDF")) is True

    def test_get_supported_extensions(self, service):
        extensions = service.get_supported_extensions()
        assert '.xlsx' in extensions
        assert '.docx' in extensions
        assert '.pptx' in extensions
        assert '.pdf' in extensions


class TestTranslationServiceGenerateOutputPath:
    """Tests for TranslationService._generate_output_path() - bidirectional"""

    @pytest.fixture
    def service(self):
        return TranslationService(Mock(), AppSettings())

    def test_adds_translated_suffix(self, service):
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "report.xlsx"
            input_path.touch()

            output = service._generate_output_path(input_path)

            assert output.name == "report_translated.xlsx"

    def test_preserves_extension(self, service):
        with tempfile.TemporaryDirectory() as tmpdir:
            for ext in ['.xlsx', '.docx', '.pptx', '.pdf']:
                input_path = Path(tmpdir) / f"file{ext}"

                output = service._generate_output_path(input_path)

                assert output.suffix == ext

    def test_adds_number_if_exists(self, service):
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "report.xlsx"
            input_path.touch()

            # Create the _translated file
            existing = Path(tmpdir) / "report_translated.xlsx"
            existing.touch()

            output = service._generate_output_path(input_path)

            assert output.name == "report_translated_2.xlsx"

    def test_increments_number_until_unique(self, service):
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "report.xlsx"
            input_path.touch()

            # Create multiple _translated files
            (Path(tmpdir) / "report_translated.xlsx").touch()
            (Path(tmpdir) / "report_translated_2.xlsx").touch()
            (Path(tmpdir) / "report_translated_3.xlsx").touch()

            output = service._generate_output_path(input_path)

            assert output.name == "report_translated_4.xlsx"

    def test_output_in_same_directory_by_default(self, service):
        with tempfile.TemporaryDirectory() as tmpdir:
            input_path = Path(tmpdir) / "report.xlsx"
            input_path.touch()

            output = service._generate_output_path(input_path)

            assert output.parent == input_path.parent

    def test_output_in_custom_directory(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output_dir = Path(tmpdir) / "output"
            output_dir.mkdir()

            settings = AppSettings(output_directory=str(output_dir))
            service = TranslationService(Mock(), settings)

            input_path = Path(tmpdir) / "report.xlsx"

            output = service._generate_output_path(input_path)

            assert output.parent == output_dir

    def test_sanitizes_forbidden_chars_in_filename(self, service):
        with tempfile.TemporaryDirectory() as tmpdir:
            # Windows-forbidden characters (:/?* etc.) should be replaced
            input_path = Path(tmpdir) / "レポート:2024?.xlsx"

            output = service._generate_output_path(input_path)

            assert output.name == "レポート_2024__translated.xlsx"

    def test_preserves_full_width_characters(self, service):
        with tempfile.TemporaryDirectory() as tmpdir:
            # Full-width characters should remain intact after sanitization
            input_path = Path(tmpdir) / "データ２０２４.xlsx"

            output = service._generate_output_path(input_path)

            assert output.name == "データ２０２４_translated.xlsx"


class TestTranslationServiceGetProcessor:
    """Tests for TranslationService._get_processor()"""

    @pytest.fixture
    def service(self):
        return TranslationService(Mock(), AppSettings())

    def test_get_processor_xlsx(self, service):
        processor = service._get_processor(Path("test.xlsx"))
        assert processor is not None

    def test_get_processor_txt(self, service):
        from yakulingo.processors.txt_processor import TxtProcessor
        processor = service._get_processor(Path("test.txt"))
        assert isinstance(processor, TxtProcessor)

    def test_get_processor_unsupported_raises(self, service):
        with pytest.raises(ValueError) as exc:
            service._get_processor(Path("test.xyz"))
        assert "Unsupported file type" in str(exc.value)

    def test_get_processor_case_insensitive(self, service):
        # Should not raise
        processor = service._get_processor(Path("test.XLSX"))
        assert processor is not None


class TestTranslationServiceCancel:
    """Tests for TranslationService.cancel()"""

    @pytest.fixture
    def service(self):
        return TranslationService(Mock(), AppSettings())

    def test_cancel_sets_flag(self, service):
        assert service._cancel_event.is_set() is False
        service.cancel()
        assert service._cancel_event.is_set() is True


class TestBatchTranslatorTranslateBlocks:
    """Tests for BatchTranslator.translate_blocks() with real PromptBuilder"""

    def test_translate_single_batch(self):
        """Test translation of blocks that fit in single batch"""
        mock_copilot = Mock()
        mock_copilot.translate_sync.return_value = ["Hello", "World"]

        prompt_builder = PromptBuilder()  # Use real PromptBuilder

        translator = BatchTranslator(mock_copilot, prompt_builder)

        blocks = [
            TextBlock(id="1", text="こんにちは", location="A1"),
            TextBlock(id="2", text="世界", location="A2"),
        ]

        results = translator.translate_blocks(blocks)

        assert results["1"] == "Hello"
        assert results["2"] == "World"
        assert mock_copilot.translate_sync.call_count == 1

    def test_translate_multiple_batches(self):
        """Test translation spanning multiple batches (due to char limit)"""
        mock_copilot = Mock()
        # Return different results for each batch (2 blocks per batch)
        mock_copilot.translate_sync.side_effect = [
            ["Trans0", "Trans1"],
            ["Trans2", "Trans3"],
        ]

        prompt_builder = PromptBuilder()  # Use real PromptBuilder

        # Create translator with small char limit to force multiple batches
        translator = BatchTranslator(mock_copilot, prompt_builder, max_chars_per_batch=1000)

        # Create 4 blocks with DIFFERENT texts (~400 chars each, 800 chars per 2 blocks < 1000 limit)
        # Using different texts to avoid deduplication
        blocks = [
            TextBlock(id=str(i), text=f"text_{i}_" + "x" * 390, location=f"A{i}")
            for i in range(4)
        ]

        results = translator.translate_blocks(blocks)

        assert len(results) == 4
        assert results["0"] == "Trans0"
        assert results["3"] == "Trans3"
        assert mock_copilot.translate_sync.call_count == 2

    def test_progress_callback_called(self):
        """Test that progress callback is invoked"""
        mock_copilot = Mock()
        mock_copilot.translate_sync.return_value = ["Hello"]

        prompt_builder = PromptBuilder()  # Use real PromptBuilder

        translator = BatchTranslator(mock_copilot, prompt_builder)

        blocks = [TextBlock(id="1", text="Test", location="A1")]

        progress_calls = []
        def on_progress(progress):
            progress_calls.append(progress)

        translator.translate_blocks(blocks, on_progress=on_progress)

        assert len(progress_calls) == 1
        assert progress_calls[0].current == 0
        assert progress_calls[0].total == 1

    def test_translate_duplicate_texts_in_batch(self):
        """Test that duplicate texts within a batch are deduplicated before sending to Copilot"""
        mock_copilot = Mock()
        # Only 2 unique texts should be sent, so only 2 translations returned
        mock_copilot.translate_sync.return_value = ["Hello", "World"]

        prompt_builder = PromptBuilder()
        translator = BatchTranslator(mock_copilot, prompt_builder)

        # 5 blocks with only 2 unique texts
        blocks = [
            TextBlock(id="1", text="こんにちは", location="A1"),
            TextBlock(id="2", text="世界", location="A2"),
            TextBlock(id="3", text="こんにちは", location="A3"),  # duplicate
            TextBlock(id="4", text="こんにちは", location="A4"),  # duplicate
            TextBlock(id="5", text="世界", location="A5"),  # duplicate
        ]

        results = translator.translate_blocks(blocks)

        # All 5 blocks should have translations
        assert len(results) == 5
        assert results["1"] == "Hello"
        assert results["2"] == "World"
        assert results["3"] == "Hello"  # same as block 1
        assert results["4"] == "Hello"  # same as block 1
        assert results["5"] == "World"  # same as block 2

        # Copilot should only be called once with 2 unique texts
        assert mock_copilot.translate_sync.call_count == 1
        call_args = mock_copilot.translate_sync.call_args
        texts_sent = call_args[0][0]  # First positional argument
        assert len(texts_sent) == 2
        assert texts_sent[0] == "こんにちは"
        assert texts_sent[1] == "世界"

    def test_translate_all_identical_texts(self):
        """Test batch where all texts are identical"""
        mock_copilot = Mock()
        # Only 1 unique text, so only 1 translation returned
        mock_copilot.translate_sync.return_value = ["Hello"]

        prompt_builder = PromptBuilder()
        translator = BatchTranslator(mock_copilot, prompt_builder)

        # 10 identical blocks
        blocks = [
            TextBlock(id=str(i), text="こんにちは", location=f"A{i}")
            for i in range(10)
        ]

        results = translator.translate_blocks(blocks)

        # All 10 blocks should have the same translation
        assert len(results) == 10
        for i in range(10):
            assert results[str(i)] == "Hello"

        # Copilot should only be called once with 1 unique text
        assert mock_copilot.translate_sync.call_count == 1
        call_args = mock_copilot.translate_sync.call_args
        texts_sent = call_args[0][0]
        assert len(texts_sent) == 1
        assert texts_sent[0] == "こんにちは"


# --- Tests: translate_text() ---

class TestTranslationServiceTranslateText:
    """Tests for TranslationService.translate_text() - bidirectional"""

    @pytest.fixture
    def mock_copilot(self):
        mock = Mock()
        mock.translate_single.return_value = "Translated text"
        return mock

    @pytest.fixture
    def service(self, mock_copilot):
        return TranslationService(mock_copilot, AppSettings())

    def test_translate_text_japanese(self, service, mock_copilot):
        """Translate Japanese text (should detect and translate to English)"""
        result = service.translate_text("こんにちは")

        assert result.status == TranslationStatus.COMPLETED
        assert result.output_text == "Translated text"
        assert result.blocks_translated == 1
        assert result.blocks_total == 1
        assert result.duration_seconds >= 0

        # Verify copilot was called
        mock_copilot.translate_single.assert_called_once()

    def test_translate_text_english(self, service, mock_copilot):
        """Translate English text (should detect and translate to Japanese)"""
        result = service.translate_text("Hello World")

        assert result.status == TranslationStatus.COMPLETED
        assert result.output_text == "Translated text"

    def test_translate_text_with_reference_files(self, service, mock_copilot, tmp_path):
        """Translation with reference files passes them to copilot"""
        glossary = tmp_path / "glossary.csv"
        glossary.write_text("JP,EN\nテスト,Test\n")

        result = service.translate_text(
            "テスト文章",
            reference_files=[glossary],
        )

        assert result.status == TranslationStatus.COMPLETED

        # Check reference files were passed
        call_args = mock_copilot.translate_single.call_args
        assert call_args[0][2] == [glossary]

    def test_translate_text_error_returns_failed(self, mock_copilot):
        """Error during translation returns FAILED status"""
        mock_copilot.translate_single.side_effect = RuntimeError("Translation error")

        service = TranslationService(mock_copilot, AppSettings())

        result = service.translate_text("テスト")

        assert result.status == TranslationStatus.FAILED
        assert result.error_message == "Translation error"
        assert result.duration_seconds >= 0

    def test_translate_text_records_duration(self, service, mock_copilot):
        """Translation records duration"""
        result = service.translate_text("テスト")

        assert result.duration_seconds is not None
        assert result.duration_seconds >= 0


# --- Tests: translate_file() ---

class TestTranslationServiceTranslateFile:
    """Tests for TranslationService.translate_file() - bidirectional"""

    @pytest.fixture
    def mock_copilot(self):
        mock = Mock()
        mock.translate_sync.return_value = ["Translated 1", "Translated 2"]
        return mock

    @pytest.fixture
    def sample_xlsx(self, tmp_path):
        """Create sample Excel file with translatable content"""
        import openpyxl
        file_path = tmp_path / "sample.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"] = "テスト1"
        ws["A2"] = "テスト2"
        wb.save(file_path)
        return file_path

    @pytest.fixture
    def empty_xlsx(self, tmp_path):
        """Create empty Excel file"""
        import openpyxl
        file_path = tmp_path / "empty.xlsx"
        wb = openpyxl.Workbook()
        wb.save(file_path)
        return file_path

    def test_translate_file_basic(self, mock_copilot, sample_xlsx):
        """Basic file translation"""
        service = TranslationService(mock_copilot, AppSettings())

        result = service.translate_file(sample_xlsx)

        assert result.status == TranslationStatus.COMPLETED
        assert result.output_path is not None
        assert result.output_path.exists()
        assert "_translated" in result.output_path.name
        assert result.blocks_translated == 2
        assert result.blocks_total == 2

    def test_translate_file_creates_output(self, mock_copilot, sample_xlsx):
        """Output file is created"""
        service = TranslationService(mock_copilot, AppSettings())

        result = service.translate_file(sample_xlsx)

        assert result.output_path.exists()
        assert result.output_path.suffix == ".xlsx"

    def test_translate_file_progress_callback(self, mock_copilot, sample_xlsx):
        """Progress callback is called during translation"""
        service = TranslationService(mock_copilot, AppSettings())

        progress_updates = []

        def on_progress(progress):
            progress_updates.append(progress)

        result = service.translate_file(
            sample_xlsx,
            on_progress=on_progress,
        )

        assert result.status == TranslationStatus.COMPLETED
        assert len(progress_updates) > 0
        # Final progress should be 100
        assert progress_updates[-1].current == 100

    def test_translate_file_empty_returns_warning(self, mock_copilot, empty_xlsx):
        """Empty file returns completed with warning"""
        service = TranslationService(mock_copilot, AppSettings())

        result = service.translate_file(empty_xlsx)

        assert result.status == TranslationStatus.COMPLETED
        assert result.blocks_total == 0
        assert result.warnings is not None
        assert any("No translatable" in w for w in result.warnings)

    def test_translate_file_with_reference(self, mock_copilot, sample_xlsx, tmp_path):
        """File translation with reference files"""
        glossary = tmp_path / "glossary.csv"
        glossary.write_text("JP,EN\n")

        service = TranslationService(mock_copilot, AppSettings())

        result = service.translate_file(
            sample_xlsx,
            reference_files=[glossary],
        )

        assert result.status == TranslationStatus.COMPLETED

    def test_translate_file_cancellation(self, mock_copilot, sample_xlsx):
        """Cancellation during translation stops the process"""
        # Note: _cancel_requested is reset at start of translate_file()
        # So cancellation only works if triggered during translation
        service = TranslationService(mock_copilot, AppSettings())

        # Simulate cancellation during batch translation
        def cancel_during_translate(*args, **kwargs):
            service.cancel()
            return ["Translated 1", "Translated 2"]

        mock_copilot.translate_sync.side_effect = cancel_during_translate

        result = service.translate_file(sample_xlsx)

        # Cancellation is checked after batch translation completes
        # Result is CANCELLED because flag was set during translation
        assert result.status == TranslationStatus.CANCELLED

    def test_translate_file_error(self, sample_xlsx):
        """Error during translation returns FAILED"""
        mock_copilot = Mock()
        mock_copilot.translate_sync.side_effect = RuntimeError("API Error")

        service = TranslationService(mock_copilot, AppSettings())

        result = service.translate_file(sample_xlsx)

        assert result.status == TranslationStatus.FAILED
        assert "API Error" in result.error_message

    def test_translate_pdf_missing_dependency(self, tmp_path):
        """Missing PDF dependencies are reported as failures instead of crashing"""
        pdf_path = tmp_path / "missing_dep.pdf"
        pdf_path.write_bytes(b"%PDF-1.4 dummy")

        service = TranslationService(Mock(), AppSettings())

        mock_processor = Mock()
        mock_processor.get_page_count.side_effect = ImportError("pymupdf not installed")

        with patch.object(service, '_get_processor', return_value=mock_processor):
            result = service.translate_file(pdf_path)

        assert result.status == TranslationStatus.FAILED
        assert "pymupdf not installed" in result.error_message

    def test_pdf_translation_respects_selected_sections(self, tmp_path):
        """PDF translation applies only to chosen pages"""
        settings = AppSettings(output_directory=str(tmp_path))
        service = TranslationService(Mock(), settings)

        processor = Mock()
        processor.get_page_count.return_value = 3
        processor.failed_pages = []

        page1_block = TextBlock(
            id="page_0_block_0", text="Page1", location="Page 1", metadata={'page_idx': 0}
        )
        page2_block = TextBlock(
            id="page_1_block_0", text="Page2", location="Page 2", metadata={'page_idx': 1}
        )

        # PDFMathTranslate compliant: page_cells is always None
        processor.extract_text_blocks_streaming.return_value = [
            ([page1_block], None),
            ([page2_block], None),
        ]

        translations = {page2_block.id: "Translated page 2"}
        service.batch_translator = Mock()
        service.batch_translator.translate_blocks.return_value = translations

        input_path = tmp_path / "sectioned.pdf"
        input_path.write_bytes(b"%PDF-1.4\n%%EOF")

        with patch.object(service, '_get_processor', return_value=processor):
            result = service.translate_file(
                input_path,
                selected_sections=[1],
            )

        # Translation should target only the second page (index 1)
        translate_blocks_args = service.batch_translator.translate_blocks.call_args.args[0]
        assert translate_blocks_args == [page2_block]

        # PDFMathTranslate compliant: apply_translations is called with text_blocks
        apply_kwargs = processor.apply_translations.call_args.kwargs
        assert apply_kwargs["pages"] == [2]
        assert apply_kwargs["text_blocks"] == [page2_block]

        assert result.status == TranslationStatus.COMPLETED
        assert result.blocks_total == 1

    def test_pdf_glossary_export_uses_text_blocks(self, tmp_path):
        """Glossary export for PDFs uses TextBlock metadata (PDFMathTranslate compliant)"""
        settings = AppSettings(export_glossary=True, output_directory=str(tmp_path))
        service = TranslationService(Mock(), settings)

        processor = Mock()
        processor.get_page_count.return_value = 1
        processor.failed_pages = []

        # PDFMathTranslate compliant: page_cells is always None
        block = TextBlock(id="page_0_block_0", text="原文", location="Page 1")
        processor.extract_text_blocks_streaming.return_value = [([block], None)]

        # Create translated PDF output
        processor.apply_translations.side_effect = (
            lambda _input_path, output_path, *_args, **_kwargs: output_path.write_bytes(b"pdf")
        )
        processor.create_bilingual_pdf = Mock()

        service.batch_translator = Mock()
        service.batch_translator.translate_blocks.return_value = {block.id: "訳文"}

        input_path = tmp_path / "input.pdf"
        input_path.write_bytes(b"%PDF-1.4\n%%EOF")

        # Avoid loading real processors (openpyxl dependency) during test
        service._processors = {'.pdf': processor}

        with patch.object(service, '_get_processor', return_value=processor):
            result = service.translate_file(input_path)

        glossary_path = tmp_path / "input_glossary.csv"

        # PDFMathTranslate compliant: glossary export uses _export_glossary_csv
        # which is called with TextBlocks (not processor.export_glossary_csv with cells)
        assert result.glossary_path == glossary_path
        assert glossary_path.exists()
        assert glossary_path.read_text(encoding="utf-8-sig").startswith("original,translated")

    def test_processor_warning_reasons_are_surface(self):
        """Failed page reasons should appear in user warnings."""
        settings = AppSettings()
        service = TranslationService(Mock(), settings)

        processor = Mock()
        processor.failed_pages = [2, 5]
        processor.failed_page_reasons = {
            2: "No embedded text detected",
            5: "Layout analysis failed",
        }
        # Ensure _layout_fallback_used is False to test only failed pages warning
        processor._layout_fallback_used = False

        warnings = service._collect_processor_warnings(processor)

        assert warnings == [
            "Pages skipped: 2 (No embedded text detected), 5 (Layout analysis failed)"
        ]

    def test_processor_warning_layout_fallback(self):
        """Layout fallback warning should appear when PP-DocLayout-L is unavailable."""
        settings = AppSettings()
        service = TranslationService(Mock(), settings)

        processor = Mock()
        processor.failed_pages = []
        processor._layout_fallback_used = True

        warnings = service._collect_processor_warnings(processor)

        assert len(warnings) == 1
        assert "PP-DocLayout-L" in warnings[0]
        assert "レイアウト解析" in warnings[0]

    def test_translate_file_custom_output_dir(self, mock_copilot, sample_xlsx, tmp_path):
        """Output goes to custom directory when configured"""
        output_dir = tmp_path / "custom_output"
        output_dir.mkdir()

        settings = AppSettings(output_directory=str(output_dir))
        service = TranslationService(mock_copilot, settings)

        result = service.translate_file(sample_xlsx)

        assert result.status == TranslationStatus.COMPLETED
        assert result.output_path.parent == output_dir


# --- Tests: get_file_info() ---

class TestTranslationServiceGetFileInfo:
    """Tests for TranslationService.get_file_info()"""

    @pytest.fixture
    def service(self):
        return TranslationService(Mock(), AppSettings())

    @pytest.fixture
    def sample_xlsx(self, tmp_path):
        import openpyxl
        file_path = tmp_path / "info_test.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"] = "テスト"
        ws["A2"] = "テスト2"
        wb.save(file_path)
        return file_path

    def test_get_file_info_returns_info(self, service, sample_xlsx):
        """get_file_info returns FileInfo object"""
        from yakulingo.models.types import FileInfo, FileType

        info = service.get_file_info(sample_xlsx)

        assert isinstance(info, FileInfo)
        assert info.path == sample_xlsx
        assert info.file_type == FileType.EXCEL
        assert info.size_bytes > 0

    def test_get_file_info_delegates_to_processor(self, service, sample_xlsx):
        """get_file_info uses correct processor"""
        info = service.get_file_info(sample_xlsx)

        # Excel processor returns correct file type
        assert info.file_type == FileType.EXCEL

    def test_get_file_info_txt(self, service, tmp_path):
        """get_file_info works for txt file"""
        txt_file = tmp_path / "test.txt"
        txt_file.write_text("content")

        info = service.get_file_info(txt_file)
        assert info.file_type == FileType.TEXT
        assert info.size_bytes > 0

    def test_get_file_info_unsupported_raises(self, service, tmp_path):
        """get_file_info raises for unsupported file type"""
        xyz_file = tmp_path / "test.xyz"
        xyz_file.write_text("content")

        with pytest.raises(ValueError):
            service.get_file_info(xyz_file)


# --- Tests: is_japanese_text() ---

class TestIsJapaneseText:
    """Tests for is_japanese_text() function"""

    def test_hiragana_detected(self):
        """Hiragana text is detected as Japanese"""
        from yakulingo.services.translation_service import is_japanese_text
        assert is_japanese_text("こんにちは") is True
        assert is_japanese_text("ひらがなテスト") is True

    def test_katakana_detected(self):
        """Katakana text is detected as Japanese"""
        from yakulingo.services.translation_service import is_japanese_text
        assert is_japanese_text("カタカナ") is True
        assert is_japanese_text("テスト") is True

    def test_kanji_detected(self):
        """Kanji text is detected as Japanese"""
        from yakulingo.services.translation_service import is_japanese_text
        assert is_japanese_text("日本語") is True
        assert is_japanese_text("漢字") is True

    def test_mixed_japanese(self):
        """Mixed Japanese content is detected"""
        from yakulingo.services.translation_service import is_japanese_text
        assert is_japanese_text("東京タワー") is True
        assert is_japanese_text("これはテストです") is True

    def test_english_not_detected(self):
        """English text is not detected as Japanese"""
        from yakulingo.services.translation_service import is_japanese_text
        assert is_japanese_text("Hello World") is False
        assert is_japanese_text("This is a test") is False

    def test_numbers_not_detected(self):
        """Numbers are not detected as Japanese"""
        from yakulingo.services.translation_service import is_japanese_text
        assert is_japanese_text("12345") is False
        assert is_japanese_text("100.50") is False

    def test_empty_string(self):
        """Empty string returns False"""
        from yakulingo.services.translation_service import is_japanese_text
        assert is_japanese_text("") is False

    def test_whitespace_only(self):
        """Whitespace only returns False"""
        from yakulingo.services.translation_service import is_japanese_text
        assert is_japanese_text("   ") is False
        assert is_japanese_text("\t\n") is False

    def test_mixed_english_japanese_above_threshold(self):
        """Mixed content above threshold is Japanese"""
        from yakulingo.services.translation_service import is_japanese_text
        # More Japanese than threshold (default 0.3)
        assert is_japanese_text("日本語 test") is True  # 3 JP / 7 total > 0.3

    def test_mixed_english_japanese_below_threshold(self):
        """Mixed content below threshold is not Japanese"""
        from yakulingo.services.translation_service import is_japanese_text
        # Less Japanese than threshold (default 0.3)
        # "a 日" has 1 JP / 2 alphanumeric = 0.5, so it's actually above threshold
        # Need more English characters to get below threshold
        assert is_japanese_text("This is English text 日") is False  # 1 JP / ~18 total < 0.3
        assert is_japanese_text("abcdefghij日") is False  # 1 JP / 11 total ≈ 0.09 < 0.3

    def test_custom_threshold(self):
        """Custom threshold works correctly"""
        from yakulingo.services.translation_service import is_japanese_text
        text = "日本 English"  # 2 JP / 9 total = ~0.22
        assert is_japanese_text(text, threshold=0.2) is True
        assert is_japanese_text(text, threshold=0.5) is False

    def test_halfwidth_katakana(self):
        """Halfwidth katakana is detected"""
        from yakulingo.services.translation_service import is_japanese_text
        # Halfwidth katakana: ｱｲｳｴｵ (U+FF65-U+FF9F)
        assert is_japanese_text("ｱｲｳｴｵ") is True

    def test_punctuation_ignored(self):
        """Punctuation is not counted in detection"""
        from yakulingo.services.translation_service import is_japanese_text
        assert is_japanese_text("。！？、") is False  # Only punctuation
        assert is_japanese_text("こんにちは。") is True  # Japanese + punctuation


# --- Tests: translate_text_with_options() ---

class TestTranslateTextWithOptions:
    """Tests for TranslationService.translate_text_with_options()"""

    @pytest.fixture
    def mock_copilot(self):
        mock = Mock()
        mock.translate_single.return_value = """[1]
訳文: Short translation
解説: Brief and concise

[2]
訳文: Medium length translation
解説: Standard translation

[3]
訳文: A longer, more detailed translation
解説: More verbose option"""
        return mock

    @pytest.fixture
    def service(self, mock_copilot, tmp_path):
        # Create prompts directory with test templates
        prompts_dir = tmp_path / "prompts"
        prompts_dir.mkdir()

        # Create text translation templates
        to_en = prompts_dir / "text_translate_to_en.txt"
        to_en.write_text("Translate to EN: {input_text}")

        to_jp = prompts_dir / "text_translate_to_jp.txt"
        to_jp.write_text("Translate to JP: {input_text}")

        return TranslationService(mock_copilot, AppSettings(), prompts_dir=prompts_dir)

    def test_japanese_input_returns_english_options(self, mock_copilot):
        """Japanese input returns English translation options"""
        # First call: language detection, second call: translation
        mock_copilot.translate_single.side_effect = [
            "日本語",  # detect_language result
            """訳文: Hello
解説: Greeting translation""",  # translation result
        ]
        service = TranslationService(mock_copilot, AppSettings())

        result = service.translate_text_with_options("こんにちは")

        assert result.output_language == "en"
        assert result.detected_language == "日本語"
        assert result.source_text == "こんにちは"
        assert result.source_char_count == 5

    def test_english_input_returns_japanese_option(self, mock_copilot):
        """English input returns Japanese translation"""
        # First call: language detection, second call: translation
        mock_copilot.translate_single.side_effect = [
            "英語",  # detect_language result
            """訳文: こんにちは
解説: 挨拶の翻訳です""",  # translation result
        ]

        service = TranslationService(mock_copilot, AppSettings())

        result = service.translate_text_with_options("Hello")

        assert result.output_language == "jp"
        assert result.detected_language == "英語"
        assert result.source_text == "Hello"

    def test_error_returns_error_message(self):
        """Error during translation returns error in result"""
        mock_copilot = Mock()
        mock_copilot.translate_single.side_effect = RuntimeError("API Error")

        service = TranslationService(mock_copilot, AppSettings())

        result = service.translate_text_with_options("テスト")

        assert result.error_message is not None
        assert "API Error" in result.error_message

    def test_fallback_when_no_prompt_file(self, mock_copilot):
        """Falls back to basic translation when prompt file missing"""
        mock_copilot.translate_single.return_value = "Translated text"

        service = TranslationService(mock_copilot, AppSettings())

        result = service.translate_text_with_options("テスト")

        # Should still return a result
        assert result.source_text == "テスト"
        # May have options or error depending on parsing

    def test_streaming_callback_passed_to_translation(self, mock_copilot, service):
        """Streaming callback is forwarded to Copilot translation call"""
        mock_copilot.translate_single.side_effect = [
            "日本語",  # detect_language result
            """訳文: Hello\n解説: Greeting""",
        ]

        chunks: list[str] = []
        def on_chunk(text: str):
            chunks.append(text)

        service.translate_text_with_options("こんにちは", on_chunk=on_chunk)

        # The translation call should receive the streaming callback
        assert mock_copilot.translate_single.call_args_list
        assert mock_copilot.translate_single.call_args_list[-1].args[3] is on_chunk

    def test_streaming_callback_passed_in_fallback(self, mock_copilot, service):
        """Fallback path also forwards streaming callback to translation"""
        # Force fallback by returning None from get_text_template
        service.prompt_builder.get_text_template = Mock(return_value=None)
        mock_copilot.translate_single.side_effect = [
            "日本語",  # detect_language result
            "訳文: Hello\n解説: Greeting",  # translation result
        ]

        chunks: list[str] = []
        def on_chunk(text: str):
            chunks.append(text)

        service.translate_text_with_options("こんにちは", on_chunk=on_chunk)

        assert mock_copilot.translate_single.call_args_list
        assert mock_copilot.translate_single.call_args_list[-1].args[3] is on_chunk


# --- Tests: detect_language_local() ---

class TestDetectLanguageLocal:
    """Tests for detect_language_local() function"""

    def test_hiragana_detected_as_japanese(self):
        """Hiragana text is detected as Japanese"""
        from yakulingo.services.translation_service import detect_language_local
        assert detect_language_local("こんにちは") == "日本語"

    def test_katakana_detected_as_japanese(self):
        """Katakana text is detected as Japanese"""
        from yakulingo.services.translation_service import detect_language_local
        assert detect_language_local("テスト") == "日本語"

    def test_mixed_japanese_detected(self):
        """Mixed hiragana/katakana/kanji is detected as Japanese"""
        from yakulingo.services.translation_service import detect_language_local
        assert detect_language_local("日本語テスト") == "日本語"
        assert detect_language_local("今日はいい天気ですね") == "日本語"

    def test_latin_detected_as_english(self):
        """Latin alphabet text is detected as English"""
        from yakulingo.services.translation_service import detect_language_local
        assert detect_language_local("Hello world") == "英語"
        assert detect_language_local("This is a test") == "英語"

    def test_hangul_detected_as_korean(self):
        """Hangul text is detected as Korean"""
        from yakulingo.services.translation_service import detect_language_local
        assert detect_language_local("안녕하세요") == "韓国語"

    def test_cjk_only_returns_none(self):
        """CJK-only text (no kana) returns None for Copilot detection"""
        from yakulingo.services.translation_service import detect_language_local
        # These could be Chinese or Japanese, need Copilot
        assert detect_language_local("東京") is None
        assert detect_language_local("中国") is None
        assert detect_language_local("日本人") is None

    def test_empty_text_returns_none(self):
        """Empty text returns None"""
        from yakulingo.services.translation_service import detect_language_local
        assert detect_language_local("") is None
        assert detect_language_local("   ") is None


# --- Tests: detect_language() (hybrid) ---

class TestDetectLanguage:
    """Tests for TranslationService.detect_language() with hybrid approach"""

    @pytest.fixture
    def mock_copilot(self):
        return Mock()

    def test_detect_language_japanese_local(self, mock_copilot):
        """Japanese text (with kana) is detected locally without Copilot"""
        service = TranslationService(mock_copilot, AppSettings())

        result = service.detect_language("こんにちは")

        assert result == "日本語"
        # Copilot should NOT be called (local detection)
        mock_copilot.translate_single.assert_not_called()

    def test_detect_language_english_local(self, mock_copilot):
        """English text is detected locally without Copilot"""
        service = TranslationService(mock_copilot, AppSettings())

        result = service.detect_language("Hello world")

        assert result == "英語"
        # Copilot should NOT be called (local detection)
        mock_copilot.translate_single.assert_not_called()

    def test_detect_language_korean_local(self, mock_copilot):
        """Korean text is detected locally without Copilot"""
        service = TranslationService(mock_copilot, AppSettings())

        result = service.detect_language("안녕하세요")

        assert result == "韓国語"
        # Copilot should NOT be called (local detection)
        mock_copilot.translate_single.assert_not_called()

    def test_detect_language_cjk_uses_copilot(self, mock_copilot):
        """CJK-only text (ambiguous) uses Copilot for detection"""
        mock_copilot.translate_single.return_value = "中国語"
        service = TranslationService(mock_copilot, AppSettings())

        result = service.detect_language("东京")  # Simplified Chinese chars

        assert result == "中国語"
        # Copilot SHOULD be called (CJK-only is ambiguous)
        mock_copilot.translate_single.assert_called_once()

    def test_detect_language_normalizes_english_variations(self, mock_copilot):
        """detect_language normalizes English language names to Japanese"""
        mock_copilot.translate_single.return_value = "Japanese"
        service = TranslationService(mock_copilot, AppSettings())

        # Use CJK-only text to trigger Copilot
        result = service.detect_language("東京")

        assert result == "日本語"

    def test_detect_language_fallback_on_empty_response(self, mock_copilot):
        """detect_language falls back to local detection on empty Copilot response"""
        mock_copilot.translate_single.return_value = ""
        service = TranslationService(mock_copilot, AppSettings())

        # CJK-only text with empty Copilot response
        result = service.detect_language("東京")
        # Fallback uses is_japanese_text which counts CJK as Japanese
        assert result == "日本語"

    def test_detect_language_fallback_on_long_response(self, mock_copilot):
        """detect_language falls back to local detection on invalid long response"""
        # Simulate Copilot error response
        mock_copilot.translate_single.return_value = "申し訳ございません。これについてチャットできません。"
        service = TranslationService(mock_copilot, AppSettings())

        # CJK-only text should still be detected via fallback
        result = service.detect_language("東京")
        assert result == "日本語"


    def test_detect_language_fallback_on_error_response(self, mock_copilot):
        """detect_language falls back when Copilot returns error message"""
        # This error message is longer than 20 chars, triggering fallback
        mock_copilot.translate_single.return_value = "I can't help with that request. Please try again."
        service = TranslationService(mock_copilot, AppSettings())

        # CJK-only text should be detected via fallback
        result = service.detect_language("東京")
        assert result == "日本語"


class TestExtractDetectionSample:
    """Tests for TranslationService.extract_detection_sample"""

    def test_retries_with_alternate_direction_for_english_docx(self, tmp_path):
        """English-only files still return a sample for detection."""
        from docx import Document

        doc_path = tmp_path / "english.docx"
        doc = Document()
        doc.add_paragraph("Hello world")
        doc.save(doc_path)

        service = TranslationService(Mock(), AppSettings())

        sample = service.extract_detection_sample(doc_path)

        assert sample is not None
        assert "Hello world" in sample

    def test_pdf_uses_fast_extraction_path(self, tmp_path):
        """PDF files use fast extraction without PP-DocLayout-L."""
        from unittest.mock import patch, MagicMock

        pdf_path = tmp_path / "test.pdf"
        pdf_path.write_bytes(b"%PDF-1.4 dummy")

        service = TranslationService(Mock(), AppSettings())

        # Mock the PDF processor's fast extraction method
        mock_processor = MagicMock()
        mock_processor.file_type = FileType.PDF
        mock_processor.extract_sample_text_fast.return_value = "テスト日本語テキスト"

        with patch.object(service, '_get_processor', return_value=mock_processor):
            sample = service.extract_detection_sample(pdf_path)

        # Verify fast extraction was called
        mock_processor.extract_sample_text_fast.assert_called_once_with(pdf_path)
        # Verify standard extraction was NOT called
        mock_processor.extract_text_blocks.assert_not_called()

        assert sample is not None
        assert "テスト日本語テキスト" in sample

    def test_pdf_falls_back_to_standard_extraction(self, tmp_path):
        """PDF falls back to standard extraction if fast path returns None."""
        from unittest.mock import patch, MagicMock
        from yakulingo.models.types import TextBlock

        pdf_path = tmp_path / "test.pdf"
        pdf_path.write_bytes(b"%PDF-1.4 dummy")

        service = TranslationService(Mock(), AppSettings())

        # Mock the PDF processor - fast extraction returns None
        mock_processor = MagicMock()
        mock_processor.file_type = FileType.PDF
        mock_processor.extract_sample_text_fast.return_value = None

        # Standard extraction returns blocks
        mock_block = TextBlock(id="1", text="Fallback text", location="Page 1", metadata={})
        mock_processor.extract_text_blocks.return_value = iter([mock_block])

        with patch.object(service, '_get_processor', return_value=mock_processor):
            sample = service.extract_detection_sample(pdf_path)

        # Verify both methods were called
        mock_processor.extract_sample_text_fast.assert_called_once()
        mock_processor.extract_text_blocks.assert_called()

        assert sample is not None
        assert "Fallback text" in sample


# --- Tests: adjust_translation() ---

class TestAdjustTranslation:
    """Tests for TranslationService.adjust_translation()"""

    @pytest.fixture
    def mock_copilot(self):
        mock = Mock()
        mock.translate_single.return_value = """訳文: Adjusted translation
解説: This is the adjusted version"""
        return mock

    @pytest.fixture
    def service(self, mock_copilot):
        return TranslationService(mock_copilot, AppSettings())

    def test_shorter_adjustment(self, service, mock_copilot):
        """Shorter adjustment returns shorter text"""
        result = service.adjust_translation("Original long translation", "shorter")

        assert result is not None
        mock_copilot.translate_single.assert_called_once()

    def test_longer_adjustment(self, service, mock_copilot):
        """Longer adjustment returns longer text"""
        result = service.adjust_translation("Short text", "longer")

        assert result is not None
        mock_copilot.translate_single.assert_called_once()

    def test_custom_adjustment(self, service, mock_copilot):
        """Custom adjustment instruction works"""
        result = service.adjust_translation("Text", "make it more formal")

        assert result is not None

    def test_error_returns_none(self):
        """Error during adjustment returns None"""
        mock_copilot = Mock()
        mock_copilot.translate_single.side_effect = RuntimeError("API Error")

        service = TranslationService(mock_copilot, AppSettings())

        result = service.adjust_translation("Text", "shorter")

        assert result is None


# --- Tests: Parsing methods ---

class TestParseMultiOptionResult:
    """Tests for TranslationService._parse_multi_option_result()"""

    @pytest.fixture
    def service(self):
        return TranslationService(Mock(), AppSettings())

    def test_parse_standard_format(self, service):
        """Parse standard multi-option format"""
        raw = """[1]
訳文: First translation
解説: First explanation

[2]
訳文: Second translation
解説: Second explanation

[3]
訳文: Third translation
解説: Third explanation"""

        options = service._parse_multi_option_result(raw)

        assert len(options) == 3
        assert options[0].text == "First translation"
        assert options[0].explanation == "First explanation"
        assert options[2].text == "Third translation"

    def test_parse_empty_result(self, service):
        """Parse empty result returns empty list"""
        options = service._parse_multi_option_result("")
        assert options == []

    def test_parse_malformed_result(self, service):
        """Parse malformed result handles gracefully"""
        raw = "Just some text without proper format"
        options = service._parse_multi_option_result(raw)
        # Should return empty or handle gracefully
        assert isinstance(options, list)


class TestParseSingleTranslationResult:
    """Tests for TranslationService._parse_single_translation_result()"""

    @pytest.fixture
    def service(self):
        return TranslationService(Mock(), AppSettings())

    def test_parse_standard_format(self, service):
        """Parse standard single translation format"""
        raw = """訳文: 翻訳されたテキスト
解説: これは翻訳の説明です"""

        options = service._parse_single_translation_result(raw)

        assert len(options) == 1
        assert options[0].text == "翻訳されたテキスト"
        assert options[0].explanation == "これは翻訳の説明です"

    def test_parse_without_explanation(self, service):
        """Parse result without explanation"""
        raw = """訳文: 翻訳されたテキスト"""

        options = service._parse_single_translation_result(raw)

        assert len(options) == 1
        assert options[0].text == "翻訳されたテキスト"

    def test_parse_fallback_format(self, service):
        """Parse fallback format (first line as text)"""
        raw = """Translation text
Some additional info"""

        options = service._parse_single_translation_result(raw)

        assert len(options) == 1
        assert options[0].text == "Translation text"

    def test_parse_empty_result(self, service):
        """Parse empty result returns empty list"""
        options = service._parse_single_translation_result("")
        assert options == []

    def test_parse_markdown_heading_without_colon(self, service):
        """Parse responses using Markdown headings without colon after 解説"""
        raw = """
ご依頼の日本語テキストを、用語集（glossary.csv）を参照し、指定のスタイル・ルールに従って英訳します。

***

### 原文

18.3兆円の補正予算案が審議入り 野党からは財政悪化に懸念も

***

### 訳文: English translation

The 18.3 trillion yen supplementary budget proposal is under review. Opposition parties also voice concerns about worsening fiscal health.

***

### 解説

- システムが回答冒頭に付けた説明を含むフォーマット
"""

        options = service._parse_single_translation_result(raw)

        assert len(options) == 1
        assert "English translation" in options[0].text
        assert "concerns about worsening fiscal health" in options[0].text
        assert "システムが回答冒頭" in options[0].explanation

    def test_parse_avoids_english_translation_metadata_leak(self, service):
        """Ensure 英訳/和訳 in preamble doesn't cause metadata to leak into translation.

        Regression test for bug where "の英訳（簡潔・略語活用...）" in Copilot's
        preamble would match the old 「訳」 pattern, causing metadata to leak.
        """
        raw = """ユーザーの依頼内容：「グーグルの『Gemini 3 Pro』、高度な視覚・空間認識で18世紀の帳簿もデータ化」の英訳（簡潔・略語活用・記号変換・段落保持・用語集参照）

訳文: Google's Gemini 3 Pro digitizes even 18th-century ledgers using advanced visual and spatial recognition.

解説:

この表現を選んだ理由：原文の要点（製品名・技術・対象物）を簡潔にまとめ、冗長な説明を省略しました。"""

        options = service._parse_single_translation_result(raw)

        assert len(options) == 1
        # Translation should NOT contain metadata from preamble
        assert "簡潔" not in options[0].text
        assert "略語活用" not in options[0].text
        assert "記号変換" not in options[0].text
        # Translation should start with the actual content
        assert options[0].text.startswith("Google's Gemini 3 Pro")
        assert "この表現を選んだ理由" in options[0].explanation


class TestParseSingleOptionResult:
    """Tests for TranslationService._parse_single_option_result()"""

    @pytest.fixture
    def service(self):
        return TranslationService(Mock(), AppSettings())

    def test_parse_standard_format(self, service):
        """Parse standard single option format"""
        raw = """訳文: Adjusted text
解説: Explanation of adjustment"""

        option = service._parse_single_option_result(raw)

        assert option is not None
        assert option.text == "Adjusted text"
        assert option.explanation == "Explanation of adjustment"

    def test_parse_fallback_format(self, service):
        """Parse fallback format (whole text as result)"""
        raw = "Just the adjusted text"

        option = service._parse_single_option_result(raw)

        assert option is not None
        assert option.text == "Just the adjusted text"

    def test_parse_empty_result(self, service):
        """Parse empty result returns None"""
        option = service._parse_single_option_result("")
        assert option is None

    def test_parse_whitespace_only(self, service):
        """Parse whitespace only returns None"""
        option = service._parse_single_option_result("   \n\t  ")
        assert option is None


# =============================================================================
# Tests: Batch Size Boundary Conditions (Character Limit Only)
# =============================================================================

class TestBatchSizeBoundaries:
    """Comprehensive tests for batch character limit boundaries"""

    @pytest.fixture
    def batch_translator(self):
        """Create BatchTranslator with explicit limits for boundary testing"""
        mock_copilot = Mock()
        mock_copilot.translate_sync.return_value = []
        prompt_builder = PromptBuilder()  # Use real PromptBuilder
        # Use explicit values for boundary tests (not defaults)
        return BatchTranslator(
            mock_copilot, prompt_builder,
            max_chars_per_batch=10000,  # Explicit value for boundary tests
        )

    # --- MAX_CHARS_PER_BATCH (10000) boundary tests ---

    def test_exactly_9999_chars_total(self, batch_translator):
        """9999 characters (one below limit) stays in one batch"""
        # Create blocks with exactly 9999 total characters
        # 10 blocks with 999 chars each, plus 9 extra in first
        text_999 = "x" * 999
        text_1008 = "x" * 1008  # 999 + 9 = 1008 to get 9999 total
        blocks = [TextBlock(id="0", text=text_1008, location="A0")]
        for i in range(1, 10):
            blocks.append(TextBlock(id=str(i), text=text_999, location=f"A{i}"))

        # Total = 1008 + 9*999 = 1008 + 8991 = 9999
        batches = batch_translator._create_batches(blocks)

        assert len(batches) == 1

    def test_exactly_10000_chars_total(self, batch_translator):
        """10000 characters (at limit) stays in one batch"""
        text_1000 = "x" * 1000
        blocks = [
            TextBlock(id=str(i), text=text_1000, location=f"A{i}")
            for i in range(10)
        ]

        batches = batch_translator._create_batches(blocks)

        # All 10 blocks should fit in one batch
        assert len(batches) == 1
        assert len(batches[0]) == 10

    def test_10001_chars_splits_batch(self, batch_translator):
        """10001 characters causes batch split"""
        text_5001 = "x" * 5001
        blocks = [
            TextBlock(id="0", text=text_5001, location="A0"),
            TextBlock(id="1", text=text_5001, location="A1"),
        ]

        # 5001 + 5001 = 10002 > 10000
        batches = batch_translator._create_batches(blocks)

        assert len(batches) == 2
        assert len(batches[0]) == 1
        assert len(batches[1]) == 1

    def test_large_single_block_gets_own_batch(self, batch_translator):
        """Very large single block (>10000 chars) gets its own batch"""
        large_text = "x" * 15000
        blocks = [
            TextBlock(id="0", text="Short", location="A0"),
            TextBlock(id="1", text=large_text, location="A1"),
            TextBlock(id="2", text="Short", location="A2"),
        ]

        batches = batch_translator._create_batches(blocks)

        # Each block should be in its own batch due to size
        assert len(batches) >= 2

    def test_char_limit_triggers_split(self, batch_translator):
        """Character limit triggers batch split"""
        # 20 blocks with 600 chars each = 12000 chars
        text_600 = "x" * 600
        blocks = [
            TextBlock(id=str(i), text=text_600, location=f"A{i}")
            for i in range(20)
        ]

        batches = batch_translator._create_batches(blocks)

        # Should split due to character limit
        # First batch: blocks 0-15 (16 * 600 = 9600)
        # Actually: 16 blocks fit, 17th would exceed
        assert len(batches) >= 2

    # --- Edge cases ---

    def test_single_block_at_char_limit(self, batch_translator):
        """Single block exactly at character limit"""
        text_10000 = "x" * 10000
        blocks = [TextBlock(id="0", text=text_10000, location="A0")]

        batches = batch_translator._create_batches(blocks)

        assert len(batches) == 1
        assert len(batches[0]) == 1

    def test_single_block_over_char_limit(self, batch_translator):
        """Single block over character limit still creates one batch"""
        text_20000 = "x" * 20000
        blocks = [TextBlock(id="0", text=text_20000, location="A0")]

        batches = batch_translator._create_batches(blocks)

        # Single block should still be in its own batch
        assert len(batches) == 1
        assert len(batches[0]) == 1

    def test_many_tiny_blocks_single_batch(self, batch_translator):
        """Many tiny blocks fit in single batch when under char limit"""
        # 200 blocks with 1 char each = 200 chars (well under 10000)
        blocks = [
            TextBlock(id=str(i), text="a", location=f"A{i}")
            for i in range(200)
        ]

        batches = batch_translator._create_batches(blocks)

        # All blocks should fit in one batch (no block count limit)
        assert len(batches) == 1
        assert len(batches[0]) == 200

    def test_alternating_large_small_blocks(self, batch_translator):
        """Alternating large and small blocks batch correctly"""
        blocks = []
        for i in range(20):
            if i % 2 == 0:
                blocks.append(TextBlock(id=str(i), text="x" * 1000, location=f"A{i}"))
            else:
                blocks.append(TextBlock(id=str(i), text="y", location=f"A{i}"))

        batches = batch_translator._create_batches(blocks)

        # Total chars = 10*1000 + 10*1 = 10010
        # Should split into 2 batches
        assert len(batches) >= 2

    def test_batch_order_preserved(self, batch_translator):
        """Block order is preserved within and across batches"""
        # Use blocks with enough text to force multiple batches
        text_1000 = "x" * 1000
        blocks = [
            TextBlock(id=str(i), text=text_1000, location=f"A{i}")
            for i in range(15)
        ]

        batches = batch_translator._create_batches(blocks)

        # Verify order within batches
        all_ids = []
        for batch in batches:
            for block in batch:
                all_ids.append(int(block.id))

        # IDs should be in order
        assert all_ids == list(range(15))


class TestBatchTranslatorTranslateBlocksAdditional:
    """Additional tests for BatchTranslator.translate_blocks() method"""

    @pytest.fixture
    def mock_copilot(self):
        mock = Mock()
        mock.translate_sync.return_value = ["Trans1", "Trans2", "Trans3"]
        return mock

    @pytest.fixture
    def prompt_builder(self):
        """Use real PromptBuilder"""
        return PromptBuilder()

    def test_translate_returns_dict(self, mock_copilot, prompt_builder):
        """translate_blocks returns dict of id -> translation"""
        translator = BatchTranslator(mock_copilot, prompt_builder)

        blocks = [
            TextBlock(id="a", text="Text1", location="A1"),
            TextBlock(id="b", text="Text2", location="A2"),
            TextBlock(id="c", text="Text3", location="A3"),
        ]

        mock_copilot.translate_sync.return_value = ["Trans1", "Trans2", "Trans3"]

        results = translator.translate_blocks(blocks)

        assert isinstance(results, dict)
        assert results["a"] == "Trans1"
        assert results["b"] == "Trans2"
        assert results["c"] == "Trans3"

    def test_translate_calls_copilot_per_batch(self, mock_copilot, prompt_builder):
        """Copilot is called once per batch (split by char limit)"""
        # Use small char limit to force multiple batches
        translator = BatchTranslator(mock_copilot, prompt_builder, max_chars_per_batch=1000)

        # Create 4 blocks with 400 chars each (2 batches of 2 blocks)
        text_400 = "x" * 400
        blocks = [
            TextBlock(id=str(i), text=text_400, location=f"A{i}")
            for i in range(4)
        ]

        # Return enough translations for each call
        mock_copilot.translate_sync.side_effect = [
            ["Trans0", "Trans1"],
            ["Trans2", "Trans3"],
        ]

        translator.translate_blocks(blocks)

        assert mock_copilot.translate_sync.call_count == 2

    def test_translate_progress_callback(self, mock_copilot, prompt_builder):
        """Progress callback is called for each batch"""
        # Use small char limit to force multiple batches
        translator = BatchTranslator(mock_copilot, prompt_builder, max_chars_per_batch=1000)

        # Create 4 blocks with 400 chars each (2 batches of 2 blocks)
        text_400 = "x" * 400
        blocks = [
            TextBlock(id=str(i), text=text_400, location=f"A{i}")
            for i in range(4)
        ]

        progress_calls = []

        def on_progress(progress):
            progress_calls.append(progress)

        mock_copilot.translate_sync.side_effect = [
            ["Trans0", "Trans1"],
            ["Trans2", "Trans3"],
        ]

        translator.translate_blocks(blocks, on_progress=on_progress)

        assert len(progress_calls) == 2

    def test_translate_with_output_language(self, mock_copilot, prompt_builder):
        """Output language affects prompt content (verified via copilot call)"""
        translator = BatchTranslator(mock_copilot, prompt_builder)

        blocks = [TextBlock(id="1", text="Test", location="A1")]
        mock_copilot.translate_sync.return_value = ["翻訳"]

        translator.translate_blocks(blocks, output_language="jp")

        # Verify copilot was called with prompt containing Japanese translation instructions
        mock_copilot.translate_sync.assert_called_once()
        call_args = mock_copilot.translate_sync.call_args
        prompt_used = call_args[0][1]  # Second positional arg is the prompt
        # Japanese output prompt should contain Japanese language instruction
        assert "日本語" in prompt_used


# --- Tests: _export_glossary_csv() ---

class TestExportGlossaryCsv:
    """Tests for TranslationService._export_glossary_csv()"""

    @pytest.fixture
    def service(self):
        return TranslationService(Mock(), AppSettings())

    def test_export_basic_glossary(self, service, tmp_path):
        """Basic glossary export creates valid CSV"""
        import csv

        blocks = [
            TextBlock(id="1", text="原文テキスト", location="A1"),
            TextBlock(id="2", text="別のテキスト", location="A2"),
        ]
        translations = {
            "1": "Translated text",
            "2": "Another translation",
        }

        output_path = tmp_path / "glossary.csv"
        service._export_glossary_csv(blocks, translations, output_path)

        assert output_path.exists()

        with open(output_path, 'r', encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            rows = list(reader)

        assert rows[0] == ['original', 'translated']
        assert rows[1] == ['原文テキスト', 'Translated text']
        assert rows[2] == ['別のテキスト', 'Another translation']

    def test_export_skips_empty_translations(self, service, tmp_path):
        """Empty translations are skipped"""
        import csv

        blocks = [
            TextBlock(id="1", text="有効なテキスト", location="A1"),
            TextBlock(id="2", text="", location="A2"),  # Empty original
            TextBlock(id="3", text="別のテキスト", location="A3"),
        ]
        translations = {
            "1": "Valid translation",
            "2": "",  # Empty translation
            "3": "Another translation",
        }

        output_path = tmp_path / "glossary.csv"
        service._export_glossary_csv(blocks, translations, output_path)

        with open(output_path, 'r', encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            rows = list(reader)

        # Header + 2 valid rows (skipping empty ones)
        assert len(rows) == 3
        assert rows[1][0] == '有効なテキスト'
        assert rows[2][0] == '別のテキスト'

    def test_export_skips_untranslated_blocks(self, service, tmp_path):
        """Blocks without translations are skipped"""
        import csv

        blocks = [
            TextBlock(id="1", text="翻訳済み", location="A1"),
            TextBlock(id="2", text="未翻訳", location="A2"),
        ]
        translations = {
            "1": "Translated",
            # "2" is missing from translations
        }

        output_path = tmp_path / "glossary.csv"
        service._export_glossary_csv(blocks, translations, output_path)

        with open(output_path, 'r', encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            rows = list(reader)

        # Header + 1 valid row
        assert len(rows) == 2
        assert rows[1][0] == '翻訳済み'

    def test_export_strips_whitespace(self, service, tmp_path):
        """Leading/trailing whitespace is stripped"""
        import csv

        blocks = [
            TextBlock(id="1", text="  原文  ", location="A1"),
        ]
        translations = {
            "1": "  Translated  ",
        }

        output_path = tmp_path / "glossary.csv"
        service._export_glossary_csv(blocks, translations, output_path)

        with open(output_path, 'r', encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            rows = list(reader)

        assert rows[1] == ['原文', 'Translated']

    def test_export_handles_special_characters(self, service, tmp_path):
        """CSV properly escapes special characters"""
        import csv

        blocks = [
            TextBlock(id="1", text='カンマ,を含む', location="A1"),
            TextBlock(id="2", text='改行\nを含む', location="A2"),
            TextBlock(id="3", text='"引用符"を含む', location="A3"),
        ]
        translations = {
            "1": "Contains, comma",
            "2": "Contains\nnewline",
            "3": '"Has quotes"',
        }

        output_path = tmp_path / "glossary.csv"
        service._export_glossary_csv(blocks, translations, output_path)

        with open(output_path, 'r', encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            rows = list(reader)

        assert len(rows) == 4
        assert rows[1][0] == 'カンマ,を含む'
        assert rows[2][0] == '改行\nを含む'
        assert rows[3][0] == '"引用符"を含む'

    def test_export_empty_blocks(self, service, tmp_path):
        """Empty block list creates CSV with only header"""
        import csv

        output_path = tmp_path / "glossary.csv"
        service._export_glossary_csv([], {}, output_path)

        with open(output_path, 'r', encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            rows = list(reader)

        assert len(rows) == 1
        assert rows[0] == ['original', 'translated']


class TestCreateBilingualOutput:
    """Tests for TranslationService._create_bilingual_output()"""

    @pytest.fixture
    def service(self):
        return TranslationService(Mock(), AppSettings())

    def test_create_bilingual_excel(self, service, tmp_path):
        """Creates bilingual Excel workbook"""
        import openpyxl

        # Create original file
        original_path = tmp_path / "original.xlsx"
        wb_orig = openpyxl.Workbook()
        ws = wb_orig.active
        ws["A1"] = "日本語"
        wb_orig.save(original_path)

        # Create translated file
        translated_path = tmp_path / "translated.xlsx"
        wb_trans = openpyxl.Workbook()
        ws = wb_trans.active
        ws["A1"] = "Japanese"
        wb_trans.save(translated_path)

        # Get Excel processor
        processor = service.processors['.xlsx']

        # Create bilingual output
        result = service._create_bilingual_output(
            original_path, translated_path, processor
        )

        assert result is not None
        assert result.exists()
        assert "_bilingual.xlsx" in result.name

    def test_create_bilingual_word(self, service, tmp_path):
        """Creates bilingual Word document"""
        from docx import Document

        # Create original file
        original_path = tmp_path / "original.docx"
        doc_orig = Document()
        doc_orig.add_paragraph("日本語テキスト")
        doc_orig.save(original_path)

        # Create translated file
        translated_path = tmp_path / "translated.docx"
        doc_trans = Document()
        doc_trans.add_paragraph("Japanese text")
        doc_trans.save(translated_path)

        # Get Word processor
        processor = service.processors['.docx']

        # Create bilingual output
        result = service._create_bilingual_output(
            original_path, translated_path, processor
        )

        assert result is not None
        assert result.exists()
        assert "_bilingual.docx" in result.name

    def test_create_bilingual_pptx(self, service, tmp_path):
        """Creates bilingual PowerPoint presentation"""
        from pptx import Presentation
        from pptx.util import Inches

        # Create original file
        original_path = tmp_path / "original.pptx"
        prs_orig = Presentation()
        slide = prs_orig.slides.add_slide(prs_orig.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.text = "日本語"
        prs_orig.save(original_path)

        # Create translated file
        translated_path = tmp_path / "translated.pptx"
        prs_trans = Presentation()
        slide = prs_trans.slides.add_slide(prs_trans.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.text = "Japanese"
        prs_trans.save(translated_path)

        # Get PowerPoint processor
        processor = service.processors['.pptx']

        # Create bilingual output
        result = service._create_bilingual_output(
            original_path, translated_path, processor
        )

        assert result is not None
        assert result.exists()
        assert "_bilingual.pptx" in result.name

    def test_returns_none_for_unsupported_type(self, service, tmp_path):
        """Returns None for unsupported file types"""
        # Create a mock processor without bilingual method
        mock_processor = Mock()
        mock_processor.create_bilingual_workbook = None
        del mock_processor.create_bilingual_workbook

        # Try to create bilingual with unsupported type
        result = service._create_bilingual_output(
            tmp_path / "file.xyz",
            tmp_path / "translated.xyz",
            mock_processor
        )

        assert result is None

    def test_returns_none_on_error(self, service, tmp_path):
        """Returns None when processor raises exception"""
        # Create a mock processor that raises exception
        mock_processor = Mock()
        mock_processor.create_bilingual_workbook = Mock(side_effect=Exception("Test error"))

        original_path = tmp_path / "original.xlsx"
        original_path.touch()
        translated_path = tmp_path / "translated.xlsx"
        translated_path.touch()

        result = service._create_bilingual_output(
            original_path, translated_path, mock_processor
        )

        assert result is None


# =============================================================================
# Tests: TranslationCache
# =============================================================================

class TestTranslationCache:
    """Tests for TranslationCache class"""

    def test_get_returns_none_for_empty_cache(self):
        """Empty cache returns None for any key"""
        cache = TranslationCache()
        assert cache.get("hello") is None

    def test_set_and_get_basic(self):
        """Basic set and get operation"""
        cache = TranslationCache()
        cache.set("hello", "こんにちは")
        assert cache.get("hello") == "こんにちは"

    def test_get_nonexistent_key(self):
        """Get nonexistent key returns None"""
        cache = TranslationCache()
        cache.set("hello", "こんにちは")
        assert cache.get("world") is None

    def test_set_overwrites_existing(self):
        """Set overwrites existing value"""
        cache = TranslationCache()
        cache.set("hello", "こんにちは")
        cache.set("hello", "ハロー")
        assert cache.get("hello") == "ハロー"

    def test_stats_initial(self):
        """Initial stats are all zeros"""
        cache = TranslationCache()
        stats = cache.stats
        assert stats["size"] == 0
        assert stats["hits"] == 0
        assert stats["misses"] == 0
        assert stats["hit_rate"] == "0.0%"

    def test_stats_after_hits(self):
        """Stats correctly count hits"""
        cache = TranslationCache()
        cache.set("hello", "こんにちは")
        cache.get("hello")  # hit
        cache.get("hello")  # hit

        stats = cache.stats
        assert stats["hits"] == 2
        assert stats["misses"] == 0
        assert stats["hit_rate"] == "100.0%"

    def test_stats_after_misses(self):
        """Stats correctly count misses"""
        cache = TranslationCache()
        cache.get("nonexistent1")  # miss
        cache.get("nonexistent2")  # miss

        stats = cache.stats
        assert stats["hits"] == 0
        assert stats["misses"] == 2
        assert stats["hit_rate"] == "0.0%"

    def test_stats_hit_rate_calculation(self):
        """Hit rate is calculated correctly"""
        cache = TranslationCache()
        cache.set("hello", "こんにちは")
        cache.get("hello")  # hit
        cache.get("world")  # miss

        stats = cache.stats
        assert stats["hits"] == 1
        assert stats["misses"] == 1
        assert stats["hit_rate"] == "50.0%"

    def test_clear_removes_all_entries(self):
        """Clear removes all cached entries"""
        cache = TranslationCache()
        cache.set("hello", "こんにちは")
        cache.set("world", "世界")
        cache.clear()

        assert cache.get("hello") is None
        assert cache.get("world") is None

    def test_clear_resets_stats(self):
        """Clear resets hit/miss statistics"""
        cache = TranslationCache()
        cache.set("hello", "こんにちは")
        cache.get("hello")  # hit
        cache.get("world")  # miss
        cache.clear()

        stats = cache.stats
        assert stats["hits"] == 0
        assert stats["misses"] == 0
        assert stats["size"] == 0

    def test_max_size_evicts_lru(self):
        """Cache evicts least recently used entry when max_size is reached"""
        cache = TranslationCache(max_size=4)

        # Add 4 entries (at max)
        cache.set("a", "1")
        cache.set("b", "2")
        cache.set("c", "3")
        cache.set("d", "4")

        # Adding 5th should evict oldest (a)
        cache.set("e", "5")

        stats = cache.stats
        # Should have 4 entries: b, c, d, e (a evicted)
        assert stats["size"] == 4
        assert cache.get("a") is None  # Evicted (LRU)
        assert cache.get("b") == "2"
        assert cache.get("c") == "3"
        assert cache.get("d") == "4"
        assert cache.get("e") == "5"

    def test_lru_access_updates_order(self):
        """Accessing an entry moves it to end (most recently used)"""
        cache = TranslationCache(max_size=3)

        # Add 3 entries
        cache.set("a", "1")
        cache.set("b", "2")
        cache.set("c", "3")

        # Access "a" - moves it to end (most recently used)
        assert cache.get("a") == "1"

        # Add new entry - should evict "b" (now the LRU)
        cache.set("d", "4")

        assert cache.get("b") is None  # Evicted (was LRU after "a" was accessed)
        assert cache.get("a") == "1"   # Still present (was accessed recently)
        assert cache.get("c") == "3"
        assert cache.get("d") == "4"

    def test_thread_safety(self):
        """Cache is thread-safe for concurrent access"""
        import threading

        cache = TranslationCache()
        errors = []

        def writer():
            for i in range(100):
                try:
                    cache.set(f"key{i}", f"value{i}")
                except Exception as e:
                    errors.append(e)

        def reader():
            for i in range(100):
                try:
                    cache.get(f"key{i}")
                except Exception as e:
                    errors.append(e)

        threads = [
            threading.Thread(target=writer),
            threading.Thread(target=reader),
            threading.Thread(target=writer),
            threading.Thread(target=reader),
        ]

        for t in threads:
            t.start()
        for t in threads:
            t.join()

        assert len(errors) == 0


# =============================================================================
# Tests: BatchTranslator Cache Integration
# =============================================================================

class TestBatchTranslatorCache:
    """Tests for BatchTranslator cache integration"""

    @pytest.fixture
    def mock_copilot(self):
        mock = Mock()
        mock.translate_sync.return_value = ["Translation1", "Translation2"]
        return mock

    @pytest.fixture
    def prompt_builder(self):
        return PromptBuilder()

    def test_cache_enabled_by_default(self, mock_copilot, prompt_builder):
        """Cache is enabled by default"""
        translator = BatchTranslator(mock_copilot, prompt_builder)
        assert translator._cache is not None

    def test_cache_can_be_disabled(self, mock_copilot, prompt_builder):
        """Cache can be disabled via enable_cache=False"""
        translator = BatchTranslator(mock_copilot, prompt_builder, enable_cache=False)
        assert translator._cache is None

    def test_cache_hit_skips_copilot(self, mock_copilot, prompt_builder):
        """Cached translations skip Copilot calls"""
        translator = BatchTranslator(mock_copilot, prompt_builder)

        blocks = [
            TextBlock(id="1", text="Hello", location="A1"),
            TextBlock(id="2", text="World", location="A2"),
        ]

        # First translation - calls Copilot
        result1 = translator.translate_blocks_with_result(blocks, output_language="日本語")
        assert mock_copilot.translate_sync.call_count == 1

        # Second translation - should use cache, no Copilot call
        mock_copilot.reset_mock()
        result2 = translator.translate_blocks_with_result(blocks, output_language="日本語")
        assert mock_copilot.translate_sync.call_count == 0

        # Results should be the same
        assert result1.translations == result2.translations

    def test_partial_cache_hit(self, mock_copilot, prompt_builder):
        """Mixed cached and uncached blocks work correctly"""
        translator = BatchTranslator(mock_copilot, prompt_builder)

        # First translation
        blocks1 = [
            TextBlock(id="1", text="Hello", location="A1"),
            TextBlock(id="2", text="World", location="A2"),
        ]
        mock_copilot.translate_sync.return_value = ["Trans1", "Trans2"]
        translator.translate_blocks_with_result(blocks1, output_language="日本語")

        # Second translation with one cached, one new
        mock_copilot.reset_mock()
        mock_copilot.translate_sync.return_value = ["NewTrans"]
        blocks2 = [
            TextBlock(id="1", text="Hello", location="A1"),  # cached
            TextBlock(id="3", text="New text", location="A3"),  # new
        ]
        result = translator.translate_blocks_with_result(blocks2, output_language="日本語")

        # Copilot should only be called once for the new block
        assert mock_copilot.translate_sync.call_count == 1

        # Verify translations
        assert result.translations["1"] == "Trans1"  # from cache
        assert result.translations["3"] == "NewTrans"  # from Copilot

    def test_all_cached_returns_early(self, mock_copilot, prompt_builder):
        """All cached blocks return early without Copilot call"""
        translator = BatchTranslator(mock_copilot, prompt_builder)

        blocks = [TextBlock(id="1", text="Hello", location="A1")]
        mock_copilot.translate_sync.return_value = ["Trans1"]

        # First call populates cache
        translator.translate_blocks_with_result(blocks, output_language="日本語")
        assert mock_copilot.translate_sync.call_count == 1

        # Second call should not call Copilot at all
        mock_copilot.reset_mock()
        result = translator.translate_blocks_with_result(blocks, output_language="日本語")

        assert mock_copilot.translate_sync.call_count == 0
        assert result.translations["1"] == "Trans1"
        assert result.translated_count == 1

    def test_cache_stores_translation_after_success(self, mock_copilot, prompt_builder):
        """Successful translations are stored in cache"""
        translator = BatchTranslator(mock_copilot, prompt_builder)

        blocks = [TextBlock(id="1", text="Hello", location="A1")]
        mock_copilot.translate_sync.return_value = ["こんにちは"]

        translator.translate_blocks_with_result(blocks, output_language="日本語")

        # Verify cache contains the translation
        cached = translator._cache.get("Hello")
        assert cached == "こんにちは"

    def test_cache_stats_after_translation(self, mock_copilot, prompt_builder):
        """Cache stats are updated correctly after translation"""
        translator = BatchTranslator(mock_copilot, prompt_builder)

        blocks = [
            TextBlock(id="1", text="Hello", location="A1"),
            TextBlock(id="2", text="World", location="A2"),
        ]
        mock_copilot.translate_sync.return_value = ["Trans1", "Trans2"]

        # First translation - 2 misses
        translator.translate_blocks_with_result(blocks, output_language="日本語")
        stats1 = translator._cache.stats
        assert stats1["misses"] == 2
        assert stats1["hits"] == 0

        # Second translation - 2 hits
        translator.translate_blocks_with_result(blocks, output_language="日本語")
        stats2 = translator._cache.stats
        assert stats2["misses"] == 2
        assert stats2["hits"] == 2
        assert stats2["hit_rate"] == "50.0%"

    def test_disabled_cache_always_calls_copilot(self, mock_copilot, prompt_builder):
        """With cache disabled, Copilot is always called"""
        translator = BatchTranslator(mock_copilot, prompt_builder, enable_cache=False)

        blocks = [TextBlock(id="1", text="Hello", location="A1")]
        mock_copilot.translate_sync.return_value = ["Trans1"]

        # First call
        translator.translate_blocks_with_result(blocks, output_language="日本語")
        assert mock_copilot.translate_sync.call_count == 1

        # Second call - should still call Copilot (no cache)
        translator.translate_blocks_with_result(blocks, output_language="日本語")
        assert mock_copilot.translate_sync.call_count == 2

    def test_clear_cache_method(self, mock_copilot, prompt_builder):
        """clear_cache() removes all cached entries"""
        translator = BatchTranslator(mock_copilot, prompt_builder)

        blocks = [TextBlock(id="1", text="Hello", location="A1")]
        mock_copilot.translate_sync.return_value = ["Trans1"]

        # Populate cache
        translator.translate_blocks_with_result(blocks, output_language="日本語")
        assert translator._cache.stats["size"] == 1

        # Clear cache
        translator.clear_cache()
        assert translator._cache.stats["size"] == 0

        # Next translation should call Copilot again
        mock_copilot.reset_mock()
        translator.translate_blocks_with_result(blocks, output_language="日本語")
        assert mock_copilot.translate_sync.call_count == 1

    def test_cache_key_exact_match_required(self, mock_copilot, prompt_builder):
        """Cache requires exact text match (not normalized)"""
        translator = BatchTranslator(mock_copilot, prompt_builder)

        # Translate "Hello"
        blocks1 = [TextBlock(id="1", text="Hello", location="A1")]
        mock_copilot.translate_sync.return_value = ["Trans1"]
        translator.translate_blocks_with_result(blocks1, output_language="日本語")

        # Try "Hello " (with trailing space) - should be cache miss
        mock_copilot.reset_mock()
        mock_copilot.translate_sync.return_value = ["Trans2"]
        blocks2 = [TextBlock(id="2", text="Hello ", location="A2")]
        result = translator.translate_blocks_with_result(blocks2, output_language="日本語")

        # Should call Copilot because key doesn't match exactly
        assert mock_copilot.translate_sync.call_count == 1
        assert result.translations["2"] == "Trans2"
