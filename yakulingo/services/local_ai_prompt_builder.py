# yakulingo/services/local_ai_prompt_builder.py
from __future__ import annotations

import csv
import heapq
import io
import json
import os
import logging
import re
import threading
import time
import unicodedata
from decimal import Decimal, InvalidOperation
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Optional, Sequence, cast

from yakulingo.config.settings import AppSettings
from yakulingo.services.prompt_builder import PromptBuilder

logger = logging.getLogger(__name__)

_TIMING_ENABLED = os.environ.get("YAKULINGO_LOCAL_AI_TIMING") == "1"


_SUPPORTED_REFERENCE_EXTENSIONS = {
    ".csv",
    ".txt",
    ".md",
    ".json",
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
}
_BUNDLED_GLOSSARY_FILENAMES = {"glossary.csv", "glossary_old.csv"}
_RE_GLOSSARY_MATCH_SEPARATORS = re.compile(r"[\s_/\\-]+")
_RE_GLOSSARY_ASCII_WORD = re.compile(r"^[a-z0-9]+$")
# NOTE: Do NOT use Unicode \b here; Japanese kana/kanji are treated as \w and break
# matching for short ASCII terms (e.g. "AIを", "GPUを").
_RE_GLOSSARY_TEXT_ASCII_WORD = re.compile(r"[a-z0-9]+")
_ASCII_ALNUM = frozenset("abcdefghijklmnopqrstuvwxyz0123456789")
_RE_JP_YEN_AMOUNT = re.compile(
    r"(?P<sign>[▲+\-−])?\s*(?:(?P<trillion>\d[\d,]*(?:\.\d+)?)兆(?:(?P<oku>\d[\d,]*(?:\.\d+)?)億)?|(?P<oku_only>\d[\d,]*(?:\.\d+)?)億)(?P<yen>円)?"
)
_RE_JP_MAN_SEN_AMOUNT = re.compile(
    r"(?P<sign>[▲+\-−])?\s*(?P<man>\d[\d,]*(?:\.\d+)?)万(?:(?P<sen>\d[\d,]*(?:\.\d+)?)千)?(?P<unit>円|台)?"
)
_RE_JP_SEN_AMOUNT = re.compile(
    r"(?P<sign>[▲+\-−])?\s*(?P<sen>\d[\d,]*(?:\.\d+)?)千(?P<unit>円|台)?"
)
_RE_TO_EN_FORBIDDEN_SYMBOLS = re.compile(r"(?:>=|<=|[><~→↑↓≥≧≤≦])")
_RE_TO_EN_MONTH = re.compile(r"\d{1,2}月")
_RE_TO_EN_MONTH_NUMBER = re.compile(r"(?P<month>\d{1,2})月")
_RE_TO_EN_TRIANGLE_NUMBER_PLAIN = re.compile(
    r"▲\s*(?P<number>\d[\d,]*(?:\.\d+)?)(?!\s*(?:兆|億|万|千))"
)
_RE_TO_EN_COMPARISON_EXAMPLE = re.compile(
    r"(?P<left>[A-Za-z0-9_.%]+)\s*(?P<op>>=|<=|≥|≧|≤|≦|>|<)\s*(?P<right>[A-Za-z0-9_.%]+)"
)
_RE_TO_EN_NUMBERED_LIST_PREFIX = re.compile(
    r"(?m)^[ \t]*(?:\d{1,3}|[０-９]{1,3})[.)．）]"
)
_RE_TO_EN_BULLET_LIST_PREFIX = re.compile(r"(?m)^[ \t]*(?:[-*•]|・)")
_RE_TO_EN_YOY_TERMS = re.compile(
    r"(前年同期比|前期比|前年比|前年度比|YoY|QoQ|CAGR)", re.IGNORECASE
)
_RE_TO_JP_OKU_WORD = re.compile(r"\boku\b", re.IGNORECASE)
_RE_TO_JP_NUMBER_K = re.compile(r"\b\d[\d,]*(?:\.\d+)?\s*k\b", re.IGNORECASE)
_RE_TO_JP_YEN_BN = re.compile(
    r"(?:¥|￥)\s*[\d,]+(?:\.\d+)?\s*(?:billion|bn)\b", re.IGNORECASE
)
_RE_TO_JP_ACCOUNTING_PAREN = re.compile(r"\(\s*[$¥￥]?\s*\d[\d,]*(?:\.\d+)?\s*\)")


@dataclass(frozen=True)
class EmbeddedReference:
    text: str
    warnings: list[str]
    truncated: bool


@dataclass(frozen=True)
class GlossaryIndex:
    pairs: list[tuple[str, str, str, str, str, str]]
    ascii_word_to_indices: dict[str, tuple[int, ...]]
    ascii_prefix4_to_indices: dict[str, tuple[int, ...]]
    compact_prefix3_to_indices: dict[str, tuple[int, ...]]
    compact_prefix2_to_indices: dict[str, tuple[int, ...]]
    compact_char1_to_indices: dict[str, tuple[int, ...]]
    scan_indices: tuple[int, ...]


class LocalPromptBuilder:
    def __init__(
        self,
        prompts_dir: Optional[Path],
        *,
        base_prompt_builder: PromptBuilder,
        settings: AppSettings,
    ) -> None:
        self.prompts_dir = prompts_dir
        self._base = base_prompt_builder
        self._settings = settings
        self._template_cache: dict[str, str] = {}
        self._template_lock = threading.Lock()
        self._rules_lock = threading.Lock()
        self._reference_cache: Optional[
            tuple[
                tuple[tuple[str, int, int], ...],
                Optional[tuple[int, str, str]],
                EmbeddedReference,
            ]
        ] = None
        self._reference_lock = threading.Lock()
        self._reference_file_cache: dict[tuple[str, int, int], tuple[str, bool]] = {}
        self._reference_file_lock = threading.Lock()
        self._glossary_cache: dict[
            tuple[str, int, int], list[tuple[str, str, str, str, str, str]]
        ] = {}
        self._glossary_index_cache: dict[tuple[str, int, int], GlossaryIndex] = {}
        self._glossary_lock = threading.Lock()

    def _get_translation_rules(self, output_language: str) -> str:
        with self._rules_lock:
            rules = self._base.get_translation_rules(output_language)
            return rules.strip()

    def _get_effective_reference_files(
        self,
        reference_files: Optional[Sequence[Path]],
        *,
        input_text: Optional[str],
    ) -> list[Path]:
        files: list[Path] = list(reference_files) if reference_files else []
        if not self._settings.use_bundled_glossary:
            return files
        if not input_text or not input_text.strip():
            return files
        if not self.prompts_dir:
            return files

        glossary_path = self.prompts_dir.parent / "glossary.csv"
        if glossary_path.exists() and glossary_path not in files:
            files.insert(0, glossary_path)
        return files

    @staticmethod
    def _input_fingerprint(text: Optional[str]) -> Optional[tuple[int, str, str]]:
        if not text:
            return None
        normalized = text.strip()
        if not normalized:
            return None
        if len(normalized) <= 128:
            return (len(normalized), normalized, normalized)
        return (len(normalized), normalized[:64], normalized[-64:])

    @staticmethod
    def _file_cache_key(path: Path) -> tuple[str, int, int]:
        try:
            stat = path.stat()
            mtime_ns = getattr(stat, "st_mtime_ns", None)
            mtime_key = (
                int(mtime_ns) if isinstance(mtime_ns, int) else int(stat.st_mtime)
            )
            return (str(path), mtime_key, int(stat.st_size))
        except OSError:
            return (str(path), 0, 0)

    def _load_glossary_pairs(
        self, path: Path, file_key: tuple[str, int, int]
    ) -> list[tuple[str, str, str, str, str, str]]:
        with self._glossary_lock:
            cached = self._glossary_cache.get(file_key)
            if cached is not None:
                return cached

        try:
            raw = path.read_text(encoding="utf-8-sig", errors="replace")
        except OSError:
            pairs: list[tuple[str, str, str, str, str, str]] = []
        else:
            pairs = []
            for row in csv.reader(io.StringIO(raw)):
                if not row:
                    continue
                first = (row[0] or "").strip()
                if not first or first.startswith("#"):
                    continue
                second = (row[1] or "").strip() if len(row) > 1 else ""
                source_folded = self._normalize_for_glossary_match(first)
                target_folded = (
                    self._normalize_for_glossary_match(second) if second else ""
                )
                source_compact = self._compact_for_glossary_match(source_folded)
                target_compact = self._compact_for_glossary_match(target_folded)
                pairs.append(
                    (
                        first,
                        second,
                        source_folded,
                        target_folded,
                        source_compact,
                        target_compact,
                    )
                )

        with self._glossary_lock:
            self._glossary_cache[file_key] = pairs
        return pairs

    def _load_glossary_index(
        self, path: Path, file_key: tuple[str, int, int]
    ) -> GlossaryIndex:
        with self._glossary_lock:
            cached = self._glossary_index_cache.get(file_key)
            if cached is not None:
                return cached

        pairs = self._load_glossary_pairs(path, file_key)
        index = self._build_glossary_index(pairs)

        with self._glossary_lock:
            self._glossary_index_cache[file_key] = index
        return index

    @staticmethod
    def _normalize_for_glossary_match(text: str) -> str:
        if not text:
            return ""
        normalized = unicodedata.normalize("NFKC", text)
        normalized = normalized.replace("\u3000", " ")
        return normalized.casefold()

    @staticmethod
    def _compact_for_glossary_match(text_folded: str) -> str:
        if not text_folded:
            return ""
        return _RE_GLOSSARY_MATCH_SEPARATORS.sub("", text_folded)

    @staticmethod
    def _sample_text_for_glossary_match(text: str, *, max_chars: int) -> str:
        text = (text or "").strip()
        if not text:
            return ""

        max_chars = max(0, int(max_chars))
        if len(text) <= max_chars:
            return text

        segments = 4
        segment_len = max(1, max_chars // segments)
        if segment_len >= len(text):
            return text[:max_chars]

        span = len(text) - segment_len
        starts = [int(span * i / (segments - 1)) for i in range(segments)]
        parts = [text[start : start + segment_len].strip() for start in starts]
        return " | ".join([part for part in parts if part])

    @staticmethod
    def _matches_glossary_term(
        *,
        text_folded: str,
        text_compact: str,
        term_folded: str,
        term_compact: str,
    ) -> bool:
        term_folded = (term_folded or "").strip()
        if not term_folded:
            return False

        if term_folded.isascii() and _RE_GLOSSARY_ASCII_WORD.match(term_folded):
            term_len = len(term_folded)
            start = 0
            while True:
                idx = text_folded.find(term_folded, start)
                if idx < 0:
                    break
                before_ok = idx == 0 or text_folded[idx - 1] not in _ASCII_ALNUM
                after_pos = idx + term_len
                after_ok = (
                    after_pos >= len(text_folded)
                    or text_folded[after_pos] not in _ASCII_ALNUM
                )
                if before_ok and after_ok:
                    return True
                start = idx + term_len
        else:
            if term_folded in text_folded:
                return True

        term_compact = (term_compact or "").strip()
        if not term_compact:
            return False
        if len(term_compact) < 4:
            return False
        return term_compact in text_compact

    @staticmethod
    def _filter_glossary_pairs(
        glossary: GlossaryIndex,
        input_text: str,
        *,
        max_lines: int,
    ) -> tuple[list[tuple[str, str]], bool]:
        text = (input_text or "").strip()
        if not text:
            return [], False

        if len(text) > 12000:
            text = LocalPromptBuilder._sample_text_for_glossary_match(
                text, max_chars=12000
            )

        text_folded = LocalPromptBuilder._normalize_for_glossary_match(text)
        text_compact = LocalPromptBuilder._compact_for_glossary_match(text_folded)

        indices_to_check: set[int] = (
            set(glossary.scan_indices) if glossary.scan_indices else set()
        )

        for match in _RE_GLOSSARY_TEXT_ASCII_WORD.finditer(text_folded):
            word = match.group(0)
            if not word:
                continue
            indices_to_check.update(glossary.ascii_word_to_indices.get(word, ()))

        if len(text_compact) >= 4 and glossary.ascii_prefix4_to_indices:
            ascii_prefix4 = glossary.ascii_prefix4_to_indices
            seen_prefixes: set[str] = set()
            for idx in range(len(text_compact) - 3):
                prefix = text_compact[idx : idx + 4]
                if not prefix.isascii():
                    continue
                indices = ascii_prefix4.get(prefix)
                if not indices:
                    continue
                if prefix in seen_prefixes:
                    continue
                seen_prefixes.add(prefix)
                indices_to_check.update(indices)

        if not text_folded.isascii():
            if glossary.compact_prefix3_to_indices and len(text_compact) >= 3:
                prefix3 = glossary.compact_prefix3_to_indices
                seen_prefixes3: set[str] = set()
                for idx in range(len(text_compact) - 2):
                    prefix = text_compact[idx : idx + 3]
                    indices = prefix3.get(prefix)
                    if not indices:
                        continue
                    if prefix in seen_prefixes3:
                        continue
                    seen_prefixes3.add(prefix)
                    indices_to_check.update(indices)

            if glossary.compact_prefix2_to_indices and len(text_compact) >= 2:
                prefix2 = glossary.compact_prefix2_to_indices
                seen_prefixes2: set[str] = set()
                for idx in range(len(text_compact) - 1):
                    prefix = text_compact[idx : idx + 2]
                    indices = prefix2.get(prefix)
                    if not indices:
                        continue
                    if prefix in seen_prefixes2:
                        continue
                    seen_prefixes2.add(prefix)
                    indices_to_check.update(indices)

            if glossary.compact_char1_to_indices:
                char1 = glossary.compact_char1_to_indices
                seen_chars: set[str] = set()
                for ch in text_compact:
                    if ch in seen_chars:
                        continue
                    indices = char1.get(ch)
                    if not indices:
                        continue
                    seen_chars.add(ch)
                    indices_to_check.update(indices)

        if not indices_to_check:
            return [], False

        seen: set[str] = set()
        heap: list[tuple[int, int, str, str]] = []
        matched_count = 0
        for idx in sorted(indices_to_check):
            (
                source,
                target,
                source_folded,
                target_folded,
                source_compact,
                target_compact,
            ) = glossary.pairs[idx]
            source = (source or "").strip()
            if not source:
                continue
            if source in seen:
                continue

            matched = LocalPromptBuilder._matches_glossary_term(
                text_folded=text_folded,
                text_compact=text_compact,
                term_folded=source_folded,
                term_compact=source_compact,
            )
            if not matched and target:
                matched = LocalPromptBuilder._matches_glossary_term(
                    text_folded=text_folded,
                    text_compact=text_compact,
                    term_folded=target_folded,
                    term_compact=target_compact,
                )
            if not matched:
                continue

            seen.add(source)
            matched_count += 1
            key = max(len(source_folded or source), len(target_folded or target))
            item = (key, -idx, source, target)
            if len(heap) < max_lines:
                heapq.heappush(heap, item)
            else:
                if item > heap[0]:
                    heapq.heapreplace(heap, item)

        if not heap:
            return [], False

        selected = sorted(heap, key=lambda x: (-x[0], -x[1]))
        return [(source, target) for _, _, source, target in selected], (
            matched_count > max_lines
        )

    @staticmethod
    def _build_glossary_index(
        pairs: list[tuple[str, str, str, str, str, str]],
    ) -> GlossaryIndex:
        ascii_word_to_indices: dict[str, list[int]] = {}
        ascii_prefix4_to_indices: dict[str, list[int]] = {}
        compact_prefix3_to_indices: dict[str, list[int]] = {}
        compact_prefix2_to_indices: dict[str, list[int]] = {}
        compact_char1_to_indices: dict[str, list[int]] = {}
        scan_indices: list[int] = []

        for idx, (
            _source,
            target,
            source_folded,
            target_folded,
            source_compact,
            target_compact,
        ) in enumerate(pairs):
            source_is_ascii_word = bool(
                source_folded.isascii() and _RE_GLOSSARY_ASCII_WORD.match(source_folded)
            )
            target_is_ascii_word = bool(
                target_folded.isascii() and _RE_GLOSSARY_ASCII_WORD.match(target_folded)
            )

            if source_is_ascii_word:
                ascii_word_to_indices.setdefault(source_folded, []).append(idx)
            if target_is_ascii_word:
                ascii_word_to_indices.setdefault(target_folded, []).append(idx)

            source_compact_is_alnum = bool(
                source_compact
                and source_compact.isascii()
                and _RE_GLOSSARY_ASCII_WORD.match(source_compact)
            )
            target_compact_is_alnum = bool(
                target_compact
                and target_compact.isascii()
                and _RE_GLOSSARY_ASCII_WORD.match(target_compact)
            )

            if source_compact_is_alnum and len(source_compact) >= 4:
                ascii_prefix4_to_indices.setdefault(source_compact[:4], []).append(idx)
            if target_compact_is_alnum and len(target_compact) >= 4:
                ascii_prefix4_to_indices.setdefault(target_compact[:4], []).append(idx)

            if source_compact and not source_folded.isascii():
                if len(source_compact) >= 3:
                    compact_prefix3_to_indices.setdefault(
                        source_compact[:3], []
                    ).append(idx)
                elif len(source_compact) >= 2:
                    compact_prefix2_to_indices.setdefault(
                        source_compact[:2], []
                    ).append(idx)
                else:
                    compact_char1_to_indices.setdefault(source_compact, []).append(idx)

            if target_compact and target and not target_folded.isascii():
                if len(target_compact) >= 3:
                    compact_prefix3_to_indices.setdefault(
                        target_compact[:3], []
                    ).append(idx)
                elif len(target_compact) >= 2:
                    compact_prefix2_to_indices.setdefault(
                        target_compact[:2], []
                    ).append(idx)
                else:
                    compact_char1_to_indices.setdefault(target_compact, []).append(idx)

            is_indexed = bool(
                source_is_ascii_word
                or target_is_ascii_word
                or (source_compact_is_alnum and len(source_compact) >= 4)
                or (target_compact_is_alnum and len(target_compact) >= 4)
                or (source_compact and not source_folded.isascii())
                or (target_compact and target and not target_folded.isascii())
            )
            if not is_indexed:
                scan_indices.append(idx)

        return GlossaryIndex(
            pairs=pairs,
            ascii_word_to_indices={
                key: tuple(indices) for key, indices in ascii_word_to_indices.items()
            },
            ascii_prefix4_to_indices={
                key: tuple(indices) for key, indices in ascii_prefix4_to_indices.items()
            },
            compact_prefix3_to_indices={
                key: tuple(indices)
                for key, indices in compact_prefix3_to_indices.items()
            },
            compact_prefix2_to_indices={
                key: tuple(indices)
                for key, indices in compact_prefix2_to_indices.items()
            },
            compact_char1_to_indices={
                key: tuple(indices) for key, indices in compact_char1_to_indices.items()
            },
            scan_indices=tuple(scan_indices),
        )

    @staticmethod
    def _join_lines_with_limit(
        lines: Sequence[str], *, max_chars: int
    ) -> tuple[str, bool]:
        if max_chars <= 0:
            return "", True
        kept: list[str] = []
        total = 0
        for line in lines:
            if not line:
                continue
            separator_len = 0 if not kept else 1  # newline
            needed = separator_len + len(line)
            if total + needed > max_chars:
                return "\n".join(kept), True
            kept.append(line)
            total += needed
        return "\n".join(kept), False

    @staticmethod
    def _split_top_level_rule_blocks(section: str) -> list[tuple[str, list[str]]]:
        blocks: list[tuple[str, list[str]]] = []
        head: Optional[str] = None
        body: list[str] = []
        for raw in (section or "").splitlines():
            line = raw.rstrip()
            stripped = line.lstrip()
            indent = len(line) - len(stripped)
            if stripped.startswith("- ") and indent == 0:
                if head is not None:
                    blocks.append((head, body))
                head = line
                body = []
                continue
            if head is None:
                if line:
                    blocks.append((line, []))
                continue
            body.append(line)
        if head is not None:
            blocks.append((head, body))
        return blocks

    @staticmethod
    def _filter_translation_rules_to_en(section: str, text: str) -> str:
        text = text or ""
        has_symbol = bool(_RE_TO_EN_FORBIDDEN_SYMBOLS.search(text))
        has_month = bool(_RE_TO_EN_MONTH.search(text))
        has_chou_oku = ("兆" in text) or ("億" in text)
        has_man = "万" in text
        has_sen = "千" in text
        has_triangle = "▲" in text
        has_yoy = bool(_RE_TO_EN_YOY_TERMS.search(text))
        lowered = text.casefold()
        has_bn_word = any(token in lowered for token in ("bn", "billion", "trillion"))

        lines: list[str] = []
        for head, body in LocalPromptBuilder._split_top_level_rule_blocks(section):
            if "禁止記号" in head:
                lines.append(head)
                if has_symbol:
                    lines.extend(body)
                continue

            if "数値/単位" in head:
                selected: list[str] = []
                for sub in body:
                    item = sub.strip()
                    if not item:
                        continue
                    if "兆/億" in item:
                        if has_chou_oku:
                            selected.append(sub)
                        continue
                    if "万→k" in item:
                        if has_man:
                            selected.append(sub)
                        continue
                    if "千→k" in item:
                        if has_sen:
                            selected.append(sub)
                        continue
                    if "▲→" in item:
                        if has_triangle:
                            selected.append(sub)
                        continue
                    if any(token in item for token in ("YoY", "QoQ", "CAGR")):
                        if has_yoy:
                            selected.append(sub)
                        continue
                    if any(
                        token in item.casefold()
                        for token in ("billion", "trillion", "bn")
                    ):
                        if has_chou_oku or has_bn_word:
                            selected.append(sub)
                        continue
                if selected:
                    lines.append(head)
                    lines.extend(selected)
                continue

            if "月名" in head:
                if has_month:
                    lines.append(head)
                continue

            if "+" in head or "「+」" in head:
                lines.append(head)
                continue

            lines.append(head)
            lines.extend(body)

        return "\n".join(lines).strip()

    @staticmethod
    def _filter_translation_rules_to_jp(section: str, text: str) -> str:
        text = text or ""
        has_oku = bool(_RE_TO_JP_OKU_WORD.search(text))
        has_k = bool(_RE_TO_JP_NUMBER_K.search(text))
        has_yen_bn = bool(_RE_TO_JP_YEN_BN.search(text))
        has_accounting = bool(_RE_TO_JP_ACCOUNTING_PAREN.search(text))

        lines: list[str] = []
        for head, body in LocalPromptBuilder._split_top_level_rule_blocks(section):
            if "数値/単位" in head:
                selected: list[str] = []
                for sub in body:
                    item = sub.strip()
                    if not item:
                        continue
                    if "例:" in item and any(
                        token in item.casefold() for token in ("billion", "bn")
                    ):
                        if has_yen_bn:
                            selected.append(sub)
                        continue
                    if "oku→億" in item:
                        if has_oku:
                            selected.append(sub)
                        continue
                    if "k→" in item:
                        if has_k:
                            selected.append(sub)
                        continue
                    if "¥/￥" in item or "billion/bn" in item:
                        if has_yen_bn:
                            selected.append(sub)
                        continue
                    if "(" in item and "▲" in item:
                        if has_accounting:
                            selected.append(sub)
                        continue
                if selected:
                    lines.append(head)
                    lines.extend(selected)
                continue

            lines.append(head)
            lines.extend(body)

        return "\n".join(lines).strip()

    def _get_translation_rules_for_text(self, output_language: str, text: str) -> str:
        if not text or not text.strip():
            return ""
        if output_language not in ("en", "jp"):
            return self._get_translation_rules(output_language)

        with self._rules_lock:
            base = self._base
            base.get_translation_rules(output_language)
            has_sections = getattr(base, "_translation_rules_has_sections", False)
            if not has_sections:
                return base.get_translation_rules(output_language).strip()

            common = base.get_translation_rules("common").strip()
            sections = getattr(base, "_translation_rules_sections", {})
            specific = ""
            if isinstance(sections, dict):
                specific = str(sections.get(output_language, "") or "").strip()

        if not specific:
            return common

        filtered = (
            self._filter_translation_rules_to_en(specific, text)
            if output_language == "en"
            else self._filter_translation_rules_to_jp(specific, text)
        ).strip()

        if common and filtered:
            return f"{common}\n\n{filtered}"
        return common or filtered

    @staticmethod
    def _format_oku_amount(value: Decimal) -> str:
        sign = "-" if value < 0 else ""
        value = abs(value)
        if value == value.to_integral():
            return f"{sign}{int(value):,}"
        text = format(value.normalize(), "f")
        int_part, _, frac_part = text.partition(".")
        int_part_with_commas = f"{int(int_part):,}"
        frac_part = frac_part.rstrip("0")
        if frac_part:
            return f"{sign}{int_part_with_commas}.{frac_part}"
        return f"{sign}{int_part_with_commas}"

    @staticmethod
    def _parse_decimal(value: str) -> Optional[Decimal]:
        if not value:
            return None
        try:
            return Decimal(value.replace(",", ""))
        except InvalidOperation:
            return None

    def _build_to_en_numeric_hints(self, text: str) -> str:
        if not text:
            return ""
        conversions: list[tuple[str, str]] = []
        max_lines = 12
        seen: set[str] = set()
        for match in _RE_JP_YEN_AMOUNT.finditer(text):
            raw = (match.group(0) or "").strip()
            if not raw or raw in seen:
                continue
            seen.add(raw)
            sign_marker = match.group("sign") or ""
            sign = -1 if sign_marker in {"▲", "-", "−"} else 1
            trillion = self._parse_decimal(match.group("trillion") or "")
            oku_part = self._parse_decimal(match.group("oku") or "")
            oku_only = self._parse_decimal(match.group("oku_only") or "")
            if trillion is not None:
                oku_value = trillion * Decimal("10000")
                if oku_part is not None:
                    oku_value += oku_part
            elif oku_only is not None:
                oku_value = oku_only
            else:
                continue
            oku_value *= sign
            formatted = self._format_oku_amount(oku_value)
            if sign < 0:
                formatted = f"({formatted.lstrip('-')})"
            unit = "oku yen" if match.group("yen") else "oku"
            conversions.append((raw, f"{formatted} {unit}".strip()))
            if len(conversions) >= max_lines:
                break

        man_sen_spans: list[tuple[int, int]] = []

        def unit_suffix(unit: str) -> str:
            if unit == "円":
                return " yen"
            if unit == "台":
                return " units"
            return ""

        for match in _RE_JP_MAN_SEN_AMOUNT.finditer(text):
            if len(conversions) >= max_lines:
                break
            raw = (match.group(0) or "").strip()
            if not raw or raw in seen:
                continue
            seen.add(raw)
            man_sen_spans.append(match.span())

            sign_marker = match.group("sign") or ""
            sign = -1 if sign_marker in {"▲", "-", "−"} else 1
            man = self._parse_decimal(match.group("man") or "")
            if man is None:
                continue
            sen = self._parse_decimal(match.group("sen") or "") or Decimal(0)
            k_value = (man * Decimal("10") + sen) * sign
            formatted = self._format_oku_amount(k_value)
            if sign < 0:
                formatted = f"({formatted.lstrip('-')})"
            suffix = unit_suffix((match.group("unit") or "").strip())
            conversions.append((raw, f"{formatted}k{suffix}"))

        def is_inside_man_sen(match_start: int) -> bool:
            return any(start <= match_start < end for start, end in man_sen_spans)

        for match in _RE_JP_SEN_AMOUNT.finditer(text):
            if len(conversions) >= max_lines:
                break
            if is_inside_man_sen(match.start()):
                continue
            raw = (match.group(0) or "").strip()
            if not raw or raw in seen:
                continue
            seen.add(raw)

            sign_marker = match.group("sign") or ""
            sign = -1 if sign_marker in {"▲", "-", "−"} else 1
            sen = self._parse_decimal(match.group("sen") or "")
            if sen is None:
                continue
            k_value = sen * sign
            formatted = self._format_oku_amount(k_value)
            if sign < 0:
                formatted = f"({formatted.lstrip('-')})"
            suffix = unit_suffix((match.group("unit") or "").strip())
            conversions.append((raw, f"{formatted}k{suffix}"))

        if not conversions:
            return ""

        lines = ["### 数値変換ヒント（必ず使用）"]
        for raw, converted in conversions:
            lines.append(f"- {raw} -> {converted}")
        return "\n".join(lines) + "\n"

    def _build_to_en_rule_hints(self, text: str) -> str:
        text = (text or "").strip()
        if not text:
            return ""

        lines: list[str] = []

        seen: set[str] = set()
        for match in _RE_TO_EN_TRIANGLE_NUMBER_PLAIN.finditer(text):
            raw = (match.group(0) or "").strip()
            number = (match.group("number") or "").strip()
            if not raw or not number:
                continue
            if raw in seen:
                continue
            seen.add(raw)
            normalized_number = number.replace(" ", "")
            lines.append(f"- ▲{normalized_number} -> ({normalized_number})")

        months: set[int] = set()
        for match in _RE_TO_EN_MONTH_NUMBER.finditer(text):
            try:
                month = int(match.group("month") or "")
            except ValueError:
                continue
            if 1 <= month <= 12:
                months.add(month)
        if months:
            month_map = {
                1: "Jan.",
                2: "Feb.",
                3: "Mar.",
                4: "Apr.",
                5: "May",
                6: "Jun.",
                7: "Jul.",
                8: "Aug.",
                9: "Sep.",
                10: "Oct.",
                11: "Nov.",
                12: "Dec.",
            }
            for month in sorted(months):
                abbrev = month_map.get(month)
                if abbrev:
                    lines.append(f"- {month}月 -> {abbrev}")

        if _RE_TO_EN_FORBIDDEN_SYMBOLS.search(text):
            op_map = {
                ">": "more than",
                "<": "less than",
                ">=": "at least",
                "≤": "at most",
                "<=": "at most",
                "≥": "at least",
                "≧": "at least",
                "≦": "at most",
            }
            for match in _RE_TO_EN_COMPARISON_EXAMPLE.finditer(text):
                raw = (match.group(0) or "").strip()
                left = (match.group("left") or "").strip()
                right = (match.group("right") or "").strip()
                op = (match.group("op") or "").strip()
                phrase = op_map.get(op)
                if not raw or not left or not right or not phrase:
                    continue
                key = f"cmp:{raw}"
                if key in seen:
                    continue
                seen.add(key)
                lines.append(f"- {raw} -> {left} {phrase} {right}")

            symbol_map = [
                ("~", "about"),
                ("→", "leads to"),
                ("↑", "up"),
                ("↓", "down"),
            ]
            for symbol, phrase in symbol_map:
                if symbol in text:
                    key = f"sym:{symbol}"
                    if key in seen:
                        continue
                    seen.add(key)
                    lines.append(f"- {symbol} -> {phrase}")

        if not lines:
            return ""

        header = "### ルール適用ヒント（必ず使用）"
        return "\n".join([header, *lines]) + "\n"

    def _build_to_en_structure_hints(self, text: str, *, include_item_ids: bool) -> str:
        text = (text or "").strip()
        if not text and not include_item_ids:
            return ""

        lines: list[str] = []
        if include_item_ids:
            lines.append("- 各訳文の先頭に [[ID:n]] をそのまま残す（削除/改変しない）")
        if "\n" in text or "\r" in text:
            lines.append("- 入力の改行/インデントを維持し、行結合/行分割しない")
        if _RE_TO_EN_NUMBERED_LIST_PREFIX.search(text):
            lines.append("- 番号付きリストの番号・記号・順序を保持し、再番号付けしない")
        if _RE_TO_EN_BULLET_LIST_PREFIX.search(text):
            lines.append("- 箇条書きの先頭記号（・/-/* 等）を保持する")

        if not lines:
            return ""
        header = "### 形式維持ヒント（必ず使用）"
        return "\n".join([header, *lines]) + "\n"

    def _get_cached_reference_text(
        self,
        path: Path,
        *,
        file_key: tuple[str, int, int],
        max_chars: int,
    ) -> tuple[Optional[str], bool]:
        with self._reference_file_lock:
            cached = self._reference_file_cache.get(file_key)
        if cached is not None:
            return cached

        try:
            content = path.read_text(encoding="utf-8-sig", errors="replace")
        except OSError:
            return None, False

        content = content.strip()
        truncated = False
        if content and len(content) > max_chars:
            content = content[:max_chars]
            truncated = True

        with self._reference_file_lock:
            self._reference_file_cache[file_key] = (content, truncated)
        return content, truncated

    def _get_cached_binary_reference_text(
        self,
        path: Path,
        *,
        file_key: tuple[str, int, int],
        suffix: str,
        max_chars: int,
    ) -> tuple[Optional[str], bool]:
        with self._reference_file_lock:
            cached = self._reference_file_cache.get(file_key)
        if cached is not None:
            return cached

        content = self._extract_binary_reference_text(
            path,
            suffix=suffix,
            max_chars=max_chars,
        )
        if content is None:
            return None, False

        truncated = False
        with self._reference_file_lock:
            self._reference_file_cache[file_key] = (content, truncated)
        return content, truncated

    def _load_template(self, filename: str) -> str:
        with self._template_lock:
            cached = self._template_cache.get(filename)
            if cached is not None:
                return cached
            if not self.prompts_dir:
                raise FileNotFoundError(
                    f"prompts_dir is not set (missing template: {filename})"
                )
            path = self.prompts_dir / filename
            if not path.exists():
                raise FileNotFoundError(f"Missing local AI prompt template: {path}")
            text = path.read_text(encoding="utf-8")
            self._template_cache[filename] = text
            return text

    def preload_startup_templates(self) -> None:
        """Preload templates needed for the first local translation (best-effort)."""
        filenames = (
            "local_text_translate_to_en_single_json.txt",
            "local_text_translate_to_en_3style_json.txt",
            "local_text_translate_to_en_missing_styles_json.txt",
            "local_text_translate_to_jp_json.txt",
            "local_batch_translate_to_en_json.txt",
            "local_batch_translate_to_jp_json.txt",
        )
        for filename in filenames:
            try:
                self._load_template(filename)
            except Exception as e:
                logger.debug(
                    "Local prompt template preload skipped (%s): %s",
                    filename,
                    e,
                )

    @staticmethod
    def _append_limited_text(
        chunks: list[str],
        text: str,
        *,
        max_chars: int,
        total: int,
    ) -> tuple[int, bool]:
        if not text:
            return total, False
        remaining = max_chars - total
        if remaining <= 0:
            return total, True
        if len(text) > remaining:
            chunks.append(text[:remaining])
            return max_chars, True
        chunks.append(text)
        return total + len(text), False

    @staticmethod
    def _extract_binary_reference_text(
        path: Path,
        *,
        suffix: str,
        max_chars: int,
    ) -> Optional[str]:
        def add_chunk(chunks: list[str], raw: str, total: int) -> tuple[int, bool]:
            text = (raw or "").strip()
            if not text:
                return total, False
            return LocalPromptBuilder._append_limited_text(
                chunks,
                f"{text}\n",
                max_chars=max_chars,
                total=total,
            )

        try:
            if suffix == ".pdf":
                import fitz

                doc = fitz.open(path)
                try:
                    chunks: list[str] = []
                    total = 0
                    for page in doc:
                        total, truncated = add_chunk(
                            chunks, cast(str, page.get_text("text")), total
                        )
                        if truncated:
                            break
                    return "".join(chunks).strip()
                finally:
                    doc.close()

            if suffix == ".docx":
                from docx import Document

                doc = Document(str(path))
                chunks = []
                total = 0
                for para in doc.paragraphs:
                    total, truncated = add_chunk(chunks, para.text, total)
                    if truncated:
                        break
                if total < max_chars:
                    for table in doc.tables:
                        if total >= max_chars:
                            break
                        for row in table.rows:
                            if total >= max_chars:
                                break
                            cells = [
                                cell.text
                                for cell in row.cells
                                if cell.text and cell.text.strip()
                            ]
                            if not cells:
                                continue
                            total, truncated = add_chunk(
                                chunks, "\t".join(cells), total
                            )
                            if truncated:
                                break
                return "".join(chunks).strip()

            if suffix == ".xlsx":
                import openpyxl

                wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
                try:
                    chunks = []
                    total = 0
                    for sheet in wb.worksheets:
                        if total >= max_chars:
                            break
                        total, truncated = add_chunk(
                            chunks, f"[Sheet] {sheet.title}", total
                        )
                        if truncated:
                            break
                        for row in sheet.iter_rows(values_only=True):
                            if total >= max_chars:
                                break
                            row_values = [
                                str(value).strip()
                                for value in row
                                if value is not None and str(value).strip()
                            ]
                            if not row_values:
                                continue
                            total, truncated = add_chunk(
                                chunks, "\t".join(row_values), total
                            )
                            if truncated:
                                break
                    return "".join(chunks).strip()
                finally:
                    wb.close()

            if suffix == ".pptx":
                from pptx import Presentation

                pres = Presentation(str(path))
                chunks = []
                total = 0
                for slide in pres.slides:
                    if total >= max_chars:
                        break
                    for shape in slide.shapes:
                        if total >= max_chars:
                            break
                        if getattr(shape, "has_text_frame", False):
                            total, truncated = add_chunk(
                                chunks, cast(Any, shape).text, total
                            )
                            if truncated:
                                break
                return "".join(chunks).strip()
        except Exception:
            return None
        return None

    def build_reference_embed(
        self,
        reference_files: Optional[Sequence[Path]],
        *,
        input_text: Optional[str] = None,
    ) -> EmbeddedReference:
        timing_enabled = _TIMING_ENABLED and logger.isEnabledFor(logging.DEBUG)
        t0 = time.perf_counter() if timing_enabled else 0.0

        reference_files = self._get_effective_reference_files(
            reference_files, input_text=input_text
        )
        if not reference_files:
            if timing_enabled:
                logger.debug(
                    "[TIMING] LocalPromptBuilder.build_reference_embed: %.4fs (files=0 cache=miss input_chars=%d embedded_chars=0 truncated=False warnings=0)",
                    time.perf_counter() - t0,
                    len((input_text or "").strip()),
                )
            return EmbeddedReference(text="", warnings=[], truncated=False)

        key_items: list[tuple[str, int, int]] = []
        file_keys: dict[Path, tuple[str, int, int]] = {}
        for path in reference_files:
            file_key = self._file_cache_key(path)
            file_keys[path] = file_key
            key_items.append(file_key)
        cache_key = tuple(key_items)
        text_key = self._input_fingerprint(input_text)

        with self._reference_lock:
            if (
                self._reference_cache
                and self._reference_cache[0] == cache_key
                and self._reference_cache[1] == text_key
            ):
                if timing_enabled:
                    embedded = self._reference_cache[2]
                    logger.debug(
                        "[TIMING] LocalPromptBuilder.build_reference_embed: %.4fs (files=%d cache=hit input_chars=%d embedded_chars=%d truncated=%s warnings=%d)",
                        time.perf_counter() - t0,
                        len(reference_files),
                        len((input_text or "").strip()),
                        len(embedded.text or ""),
                        bool(embedded.truncated),
                        len(embedded.warnings or []),
                    )
                return self._reference_cache[2]

        max_total_chars = 4000
        max_file_chars = 2000
        glossary_max_lines = 80

        warnings: list[str] = []
        truncated = False
        total = 0
        parts: list[str] = []

        for path in reference_files:
            file_key = file_keys.get(path) or self._file_cache_key(path)
            suffix = path.suffix.lower()
            if suffix not in _SUPPORTED_REFERENCE_EXTENSIONS:
                warnings.append(f"未対応の参照ファイルをスキップしました: {path.name}")
                continue

            is_bundled_glossary = (
                suffix == ".csv" and path.name.casefold() in _BUNDLED_GLOSSARY_FILENAMES
            )
            if is_bundled_glossary:
                if not input_text or not input_text.strip():
                    continue
                glossary = self._load_glossary_index(path, file_key)
                matched, glossary_truncated = self._filter_glossary_pairs(
                    glossary, input_text or "", max_lines=glossary_max_lines
                )
                if not matched:
                    continue
                lines = [
                    f"{source} 翻译成 {target}" for source, target in matched if target
                ]
                if not lines:
                    continue
                remaining_total = max_total_chars - total
                max_glossary_chars = min(max_file_chars, max(0, remaining_total))
                content, was_truncated = self._join_lines_with_limit(
                    lines, max_chars=max_glossary_chars
                )
                was_truncated = was_truncated or glossary_truncated
                if not content:
                    continue
            elif suffix in {".txt", ".md", ".json", ".csv"}:
                content, was_truncated = self._get_cached_reference_text(
                    path,
                    file_key=file_key,
                    max_chars=max_file_chars,
                )
                if content is None:
                    warnings.append(f"参照ファイルを読み込めませんでした: {path.name}")
                    continue
                if not content:
                    continue
            else:
                content, was_truncated = self._get_cached_binary_reference_text(
                    path,
                    suffix=suffix,
                    file_key=file_key,
                    max_chars=max_file_chars,
                )
                if not content:
                    warnings.append(f"参照ファイルを読み込めませんでした: {path.name}")
                    continue

            if was_truncated or len(content) > max_file_chars:
                if len(content) > max_file_chars:
                    content = content[:max_file_chars]
                truncated = True
                if not is_bundled_glossary:
                    warnings.append(
                        f"参照ファイルを一部省略しました（上限 {max_file_chars} 文字）: {path.name}"
                    )

            remaining = max_total_chars - total
            if remaining <= 0:
                truncated = True
                if not is_bundled_glossary:
                    warnings.append(
                        f"参照ファイルを一部省略しました（合計上限 {max_total_chars} 文字）"
                    )
                break

            if len(content) > remaining:
                content = content[:remaining]
                truncated = True
                if not is_bundled_glossary:
                    warnings.append(
                        f"参照ファイルを一部省略しました（合計上限 {max_total_chars} 文字）"
                    )

            total += len(content)
            parts.append(f"[REFERENCE:file={path.name}]\n{content}\n[/REFERENCE]")

        if not parts:
            embedded = EmbeddedReference(
                text="", warnings=warnings, truncated=truncated
            )
            with self._reference_lock:
                self._reference_cache = (cache_key, text_key, embedded)
            if timing_enabled:
                logger.debug(
                    "[TIMING] LocalPromptBuilder.build_reference_embed: %.4fs (files=%d cache=miss input_chars=%d embedded_chars=0 truncated=%s warnings=%d)",
                    time.perf_counter() - t0,
                    len(reference_files),
                    len((input_text or "").strip()),
                    bool(truncated),
                    len(warnings),
                )
            return embedded

        header = (
            "### 参照（埋め込み）\n"
            "以下の参照を優先して翻訳してください（必要に応じて省略されています）。\n"
        )
        header = (
            "### Reference files (critical)\n"
            "- GLOSSARY (mandatory): apply glossary terms everywhere they appear, including inside longer sentences.\n"
            "- Use the glossary's preferred wording verbatim; do not paraphrase glossary terms.\n"
            "\n"
        )
        embedded_text = header + "\n\n".join(parts)
        embedded = EmbeddedReference(
            text=embedded_text, warnings=warnings, truncated=truncated
        )
        with self._reference_lock:
            self._reference_cache = (cache_key, text_key, embedded)
        if timing_enabled:
            logger.debug(
                "[TIMING] LocalPromptBuilder.build_reference_embed: %.4fs (files=%d cache=miss input_chars=%d embedded_chars=%d truncated=%s warnings=%d)",
                time.perf_counter() - t0,
                len(reference_files),
                len((input_text or "").strip()),
                len(embedded.text or ""),
                bool(embedded.truncated),
                len(embedded.warnings or []),
            )
        return embedded

    def build_batch(
        self,
        texts: list[str],
        has_reference_files: bool = False,
        output_language: str = "en",
        translation_style: str = "concise",
        include_item_ids: bool = False,
        reference_files: Optional[Sequence[Path]] = None,
    ) -> str:
        timing_enabled = _TIMING_ENABLED and logger.isEnabledFor(logging.DEBUG)
        t0 = time.perf_counter() if timing_enabled else 0.0

        if output_language not in ("en", "jp"):
            output_language = "en"

        def build_rule_context_text(*, max_chars: int) -> str:
            if output_language != "en":
                return ""
            if not texts:
                return ""

            def slice_around(value: str, *, pos: int, budget: int) -> str:
                if budget <= 0 or not value:
                    return ""
                if len(value) <= budget:
                    return value
                half = budget // 2
                start = max(0, pos - half)
                end = start + budget
                if end > len(value):
                    end = len(value)
                    start = max(0, end - budget)
                return value[start:end]

            def first_char_pos(value: str, chars: tuple[str, ...]) -> int | None:
                positions = [value.find(ch) for ch in chars]
                positions = [p for p in positions if p != -1]
                return min(positions) if positions else None

            parts: list[str] = []
            total = 0
            seen: set[str] = set()
            per_snippet_max = min(800, max_chars)

            def add_snippet(item: str, *, pos: int) -> None:
                nonlocal total
                separator_len = 0 if not parts else 1
                remaining = max_chars - total - separator_len
                if remaining <= 0:
                    return
                budget = min(per_snippet_max, remaining)
                snippet = slice_around(item, pos=pos, budget=budget).strip()
                if not snippet or snippet in seen:
                    return
                if separator_len:
                    total += 1
                parts.append(snippet)
                seen.add(snippet)
                total += len(snippet)

            for item in texts:
                if not item:
                    continue
                pos = first_char_pos(item, ("兆", "億"))
                if pos is not None:
                    add_snippet(item, pos=pos)
                    break

            for item in texts:
                if not item:
                    continue
                pos = item.find("▲")
                if pos != -1:
                    add_snippet(item, pos=pos)
                    break

            for item in texts:
                if not item:
                    continue
                if total >= max_chars:
                    break
                pos = first_char_pos(item, ("万", "千"))
                if pos is not None:
                    add_snippet(item, pos=pos)
                    continue
                for pattern in (
                    _RE_TO_EN_FORBIDDEN_SYMBOLS,
                    _RE_TO_EN_MONTH,
                    _RE_TO_EN_YOY_TERMS,
                ):
                    match = pattern.search(item)
                    if match:
                        add_snippet(item, pos=match.start())
                        break

            return "\n".join(parts).strip()

        filename = (
            "local_batch_translate_to_en_json.txt"
            if output_language == "en"
            else "local_batch_translate_to_jp_json.txt"
        )
        template = self._load_template(filename)

        max_context_chars = 3000
        context_parts: list[str] = []
        total_chars = 0
        for item in texts:
            if not item:
                continue
            if total_chars >= max_context_chars:
                break
            remaining = max_context_chars - total_chars
            if len(item) > remaining:
                context_parts.append(item[:remaining])
                total_chars = max_context_chars
                break
            context_parts.append(item)
            total_chars += len(item) + 1
        context_text = "\n".join(context_parts)
        rule_context_text = build_rule_context_text(max_chars=max_context_chars)
        if not rule_context_text:
            rule_context_text = context_text
        translation_rules = (
            self._get_translation_rules_for_text(output_language, rule_context_text)
            if rule_context_text.strip()
            else self._get_translation_rules(output_language)
        )
        embedded_ref = self.build_reference_embed(
            reference_files, input_text=context_text
        )
        reference_section = embedded_ref.text if embedded_ref.text else ""
        numeric_hints = (
            self._build_to_en_numeric_hints(rule_context_text)
            if output_language == "en"
            else ""
        )
        rule_hints = (
            self._build_to_en_rule_hints(rule_context_text)
            if output_language == "en"
            else ""
        )
        structure_hints = (
            self._build_to_en_structure_hints(
                context_text, include_item_ids=include_item_ids
            )
            if output_language == "en"
            else ""
        )
        hint_parts = [
            part.strip()
            for part in (numeric_hints, rule_hints, structure_hints)
            if part and part.strip()
        ]
        merged_hints = "\n\n".join(hint_parts).strip()
        numeric_hints = f"{merged_hints}\n" if merged_hints else ""

        items = [
            {
                "id": i + 1,
                "text": f"[[ID:{i + 1}]] {text}" if include_item_ids else text,
            }
            for i, text in enumerate(texts)
        ]
        items_json = json.dumps(
            {"items": items}, ensure_ascii=False, separators=(",", ":")
        )

        prompt = template.replace("{translation_rules}", translation_rules)
        prompt = prompt.replace("{numeric_hints}", numeric_hints)
        prompt = prompt.replace("{reference_section}", reference_section)
        prompt = prompt.replace("{style}", translation_style)
        prompt = prompt.replace("{items_json}", items_json)
        prompt = prompt.replace("{output_language}", output_language)
        prompt = prompt.replace("{n_items}", str(len(items)))
        if timing_enabled:
            logger.debug(
                "[TIMING] LocalPromptBuilder.build_batch: %.4fs (items=%d output=%s style=%s prompt_chars=%d ref_chars=%d)",
                time.perf_counter() - t0,
                len(items),
                output_language,
                translation_style,
                len(prompt),
                len(reference_section or ""),
            )
        return prompt

    def build_text_to_en_3style(
        self,
        text: str,
        *,
        reference_files: Optional[Sequence[Path]] = None,
        detected_language: str = "日本語",
        extra_instruction: str | None = None,
    ) -> str:
        timing_enabled = _TIMING_ENABLED and logger.isEnabledFor(logging.DEBUG)
        t0 = time.perf_counter() if timing_enabled else 0.0

        template = self._load_template("local_text_translate_to_en_3style_json.txt")
        timing_enabled = _TIMING_ENABLED and logger.isEnabledFor(logging.DEBUG)
        t0 = time.perf_counter() if timing_enabled else 0.0

        embedded_ref = self.build_reference_embed(reference_files, input_text=text)
        translation_rules = self._get_translation_rules_for_text("en", text)
        numeric_hints = self._build_to_en_numeric_hints(text)
        rule_hints = self._build_to_en_rule_hints(text)
        reference_section = embedded_ref.text if embedded_ref.text else ""
        prompt_input_text = self._base.normalize_input_text(text, "en")
        extra_instruction = extra_instruction.strip() if extra_instruction else ""
        extra_parts = [part for part in (rule_hints.strip(), extra_instruction) if part]
        merged_extra = "\n\n".join(extra_parts).strip()
        extra_instruction = f"{merged_extra}\n" if merged_extra else ""
        prompt = template.replace("{translation_rules}", translation_rules)
        prompt = prompt.replace("{numeric_hints}", numeric_hints)
        prompt = prompt.replace("{reference_section}", reference_section)
        prompt = prompt.replace("{extra_instruction}", extra_instruction)
        prompt = prompt.replace("{input_text}", prompt_input_text)
        prompt = prompt.replace("{detected_language}", detected_language)
        if timing_enabled:
            logger.debug(
                "[TIMING] LocalPromptBuilder.build_text_to_en_3style: %.4fs (input_chars=%d prompt_chars=%d ref_chars=%d)",
                time.perf_counter() - t0,
                len(text or ""),
                len(prompt),
                len(reference_section or ""),
            )
        return prompt

    def build_text_to_en_missing_styles(
        self,
        text: str,
        *,
        styles: Sequence[str],
        reference_files: Optional[Sequence[Path]] = None,
        detected_language: str = "日本語",
        extra_instruction: str | None = None,
    ) -> str:
        timing_enabled = _TIMING_ENABLED and logger.isEnabledFor(logging.DEBUG)
        t0 = time.perf_counter() if timing_enabled else 0.0

        template = self._load_template(
            "local_text_translate_to_en_missing_styles_json.txt"
        )
        embedded_ref = self.build_reference_embed(reference_files, input_text=text)
        translation_rules = self._get_translation_rules_for_text("en", text)
        numeric_hints = self._build_to_en_numeric_hints(text)
        rule_hints = self._build_to_en_rule_hints(text)
        reference_section = embedded_ref.text if embedded_ref.text else ""
        prompt_input_text = self._base.normalize_input_text(text, "en")
        extra_instruction = extra_instruction.strip() if extra_instruction else ""
        extra_parts = [part for part in (rule_hints.strip(), extra_instruction) if part]
        merged_extra = "\n\n".join(extra_parts).strip()
        extra_instruction = f"{merged_extra}\n" if merged_extra else ""
        style_list: list[str] = []
        seen: set[str] = set()
        for style in styles:
            if not style or style in seen:
                continue
            seen.add(style)
            style_list.append(style)
        styles_json = json.dumps(style_list, ensure_ascii=False, separators=(",", ":"))
        prompt = template.replace("{translation_rules}", translation_rules)
        prompt = prompt.replace("{numeric_hints}", numeric_hints)
        prompt = prompt.replace("{reference_section}", reference_section)
        prompt = prompt.replace("{extra_instruction}", extra_instruction)
        prompt = prompt.replace("{input_text}", prompt_input_text)
        prompt = prompt.replace("{detected_language}", detected_language)
        prompt = prompt.replace("{styles_json}", styles_json)
        prompt = prompt.replace("{n_styles}", str(len(style_list)))
        if timing_enabled:
            logger.debug(
                "[TIMING] LocalPromptBuilder.build_text_to_en_missing_styles: %.4fs (input_chars=%d styles=%d prompt_chars=%d ref_chars=%d)",
                time.perf_counter() - t0,
                len(text or ""),
                len(style_list),
                len(prompt),
                len(reference_section or ""),
            )
        return prompt

    def build_text_to_en_single(
        self,
        text: str,
        *,
        style: str,
        reference_files: Optional[Sequence[Path]] = None,
        detected_language: str = "日本語",
        extra_instruction: str | None = None,
    ) -> str:
        timing_enabled = _TIMING_ENABLED and logger.isEnabledFor(logging.DEBUG)
        t0 = time.perf_counter() if timing_enabled else 0.0

        template = self._load_template("local_text_translate_to_en_single_json.txt")
        embedded_ref = self.build_reference_embed(reference_files, input_text=text)
        translation_rules = self._get_translation_rules_for_text("en", text)
        numeric_hints = self._build_to_en_numeric_hints(text)
        rule_hints = self._build_to_en_rule_hints(text)
        reference_section = embedded_ref.text if embedded_ref.text else ""
        prompt_input_text = self._base.normalize_input_text(text, "en")
        extra_instruction = extra_instruction.strip() if extra_instruction else ""
        extra_parts = [part for part in (rule_hints.strip(), extra_instruction) if part]
        merged_extra = "\n\n".join(extra_parts).strip()
        extra_instruction = f"{merged_extra}\n" if merged_extra else ""
        prompt = template.replace("{translation_rules}", translation_rules)
        prompt = prompt.replace("{numeric_hints}", numeric_hints)
        prompt = prompt.replace("{extra_instruction}", extra_instruction)
        prompt = prompt.replace("{reference_section}", reference_section)
        prompt = prompt.replace("{input_text}", prompt_input_text)
        prompt = prompt.replace("{style}", style)
        prompt = prompt.replace("{detected_language}", detected_language)
        if timing_enabled:
            logger.debug(
                "[TIMING] LocalPromptBuilder.build_text_to_en_single: %.4fs (input_chars=%d style=%s prompt_chars=%d ref_chars=%d)",
                time.perf_counter() - t0,
                len(text or ""),
                style,
                len(prompt),
                len(reference_section or ""),
            )
        return prompt

    def build_text_to_jp(
        self,
        text: str,
        *,
        reference_files: Optional[Sequence[Path]] = None,
        detected_language: str = "英語",
    ) -> str:
        timing_enabled = _TIMING_ENABLED and logger.isEnabledFor(logging.DEBUG)
        t0 = time.perf_counter() if timing_enabled else 0.0

        template = self._load_template("local_text_translate_to_jp_json.txt")
        embedded_ref = self.build_reference_embed(reference_files, input_text=text)
        translation_rules = self._get_translation_rules_for_text("jp", text)
        reference_section = embedded_ref.text if embedded_ref.text else ""
        prompt_input_text = self._base.normalize_input_text(text, "jp")
        prompt = template.replace("{translation_rules}", translation_rules)
        prompt = prompt.replace("{reference_section}", reference_section)
        prompt = prompt.replace("{input_text}", prompt_input_text)
        prompt = prompt.replace("{detected_language}", detected_language)
        if timing_enabled:
            logger.debug(
                "[TIMING] LocalPromptBuilder.build_text_to_jp: %.4fs (input_chars=%d prompt_chars=%d ref_chars=%d)",
                time.perf_counter() - t0,
                len(text or ""),
                len(prompt),
                len(reference_section or ""),
            )
        return prompt
