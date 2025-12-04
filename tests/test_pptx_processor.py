# tests/test_pptx_processor.py
"""Tests for yakulingo.processors.pptx_processor"""

import pytest
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

from yakulingo.processors.pptx_processor import PptxProcessor
from yakulingo.models.types import FileType


# --- Fixtures ---

@pytest.fixture
def processor():
    """PptxProcessor instance"""
    return PptxProcessor()


@pytest.fixture
def sample_pptx(tmp_path):
    """Create a sample PowerPoint with text"""
    file_path = tmp_path / "sample.pptx"
    prs = Presentation()

    # Slide 1
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)

    # Add text box
    left = Inches(1)
    top = Inches(1)
    width = Inches(5)
    height = Inches(1)
    txBox = slide1.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "スライド1のテキスト"

    # Add another shape with number (should skip)
    txBox2 = slide1.shapes.add_textbox(left, Inches(2), width, height)
    txBox2.text_frame.text = "12345"

    prs.save(file_path)
    return file_path


@pytest.fixture
def pptx_with_table(tmp_path):
    """Create PowerPoint with table"""
    file_path = tmp_path / "table.pptx"
    prs = Presentation()

    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Add title text
    left = Inches(0.5)
    top = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, Inches(8), Inches(0.5))
    txBox.text_frame.text = "Table Slide Title"

    # Add table
    rows, cols = 2, 2
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(6), Inches(2)).table

    table.cell(0, 0).text = "ヘッダー1"
    table.cell(0, 1).text = "Header 2"
    table.cell(1, 0).text = "12345"  # Should skip
    table.cell(1, 1).text = "データ"

    prs.save(file_path)
    return file_path


@pytest.fixture
def pptx_with_multiple_slides(tmp_path):
    """Create PowerPoint with multiple slides"""
    file_path = tmp_path / "multi.pptx"
    prs = Presentation()

    slide_layout = prs.slide_layouts[5]

    # Slide 1 - Japanese text
    slide1 = prs.slides.add_slide(slide_layout)
    txBox1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    txBox1.text_frame.text = "スライド1のテキスト"

    # Slide 2 - Japanese text
    slide2 = prs.slides.add_slide(slide_layout)
    txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    txBox2.text_frame.text = "スライド2のテキスト"

    prs.save(file_path)
    return file_path


@pytest.fixture
def empty_pptx(tmp_path):
    """Create empty PowerPoint"""
    file_path = tmp_path / "empty.pptx"
    prs = Presentation()
    prs.save(file_path)
    return file_path


# --- Tests: Properties ---

class TestPptxProcessorProperties:
    """Test PptxProcessor properties"""

    def test_file_type(self, processor):
        assert processor.file_type == FileType.POWERPOINT

    def test_supported_extensions(self, processor):
        extensions = processor.supported_extensions
        assert ".pptx" in extensions
        # .ppt (legacy format) is not supported by python-pptx
        assert ".ppt" not in extensions


# --- Tests: get_file_info ---

class TestPptxProcessorGetFileInfo:
    """Test PptxProcessor.get_file_info()"""

    def test_file_info_basic(self, processor, sample_pptx):
        """Basic file info retrieval"""
        info = processor.get_file_info(sample_pptx)

        assert info.path == sample_pptx
        assert info.file_type == FileType.POWERPOINT
        assert info.size_bytes > 0
        assert info.slide_count == 1

    def test_file_info_multiple_slides(self, processor, pptx_with_multiple_slides):
        """File info counts slides"""
        info = processor.get_file_info(pptx_with_multiple_slides)
        assert info.slide_count == 2

    def test_file_info_with_table(self, processor, pptx_with_table):
        """File info includes slide count"""
        info = processor.get_file_info(pptx_with_table)
        assert info.slide_count == 1

    def test_file_info_empty(self, processor, empty_pptx):
        """File info for empty presentation"""
        info = processor.get_file_info(empty_pptx)
        assert info.slide_count == 0


# --- Tests: extract_text_blocks ---

class TestPptxProcessorExtractTextBlocks:
    """Test PptxProcessor.extract_text_blocks()"""

    def test_extracts_shape_text(self, processor, sample_pptx):
        """Extracts translatable shape text"""
        blocks = list(processor.extract_text_blocks(sample_pptx))

        # Should have 1 block (number skipped)
        assert len(blocks) == 1
        assert blocks[0].text == "スライド1のテキスト"

    def test_skips_numbers(self, processor, sample_pptx):
        """Skips non-translatable content"""
        blocks = list(processor.extract_text_blocks(sample_pptx))
        texts = [b.text for b in blocks]

        assert "12345" not in texts

    def test_extracts_table_cells(self, processor, pptx_with_table):
        """Extracts table cells (Japanese only)"""
        blocks = list(processor.extract_text_blocks(pptx_with_table))

        # Filter table cells - only Japanese cells are extracted
        table_blocks = [b for b in blocks if b.metadata.get("type") == "table_cell"]
        # "ヘッダー1" and "データ" are Japanese, "Header 2" and "12345" are skipped
        assert len(table_blocks) == 2

        table_texts = [b.text for b in table_blocks]
        assert "ヘッダー1" in table_texts
        # Note: "Header 2" is skipped as English-only
        assert "データ" in table_texts
        assert "12345" not in table_texts

    def test_extracts_from_multiple_slides(self, processor, pptx_with_multiple_slides):
        """Extracts from all slides"""
        blocks = list(processor.extract_text_blocks(pptx_with_multiple_slides))

        assert len(blocks) == 2

        slides = {b.metadata["slide"] for b in blocks}
        assert 0 in slides
        assert 1 in slides

    def test_block_metadata_shape(self, processor, sample_pptx):
        """Shape blocks have correct metadata"""
        blocks = list(processor.extract_text_blocks(sample_pptx))

        shape_block = blocks[0]
        assert shape_block.metadata["type"] == "shape"
        assert "slide" in shape_block.metadata
        assert "shape" in shape_block.metadata

    def test_block_metadata_table(self, processor, pptx_with_table):
        """Table cell blocks have correct metadata"""
        blocks = list(processor.extract_text_blocks(pptx_with_table))

        table_block = next(b for b in blocks if b.metadata.get("type") == "table_cell")
        assert "slide" in table_block.metadata
        assert "table" in table_block.metadata
        assert "row" in table_block.metadata
        assert "col" in table_block.metadata

    def test_block_ids_unique(self, processor, pptx_with_table):
        """All block IDs are unique"""
        blocks = list(processor.extract_text_blocks(pptx_with_table))
        ids = [b.id for b in blocks]
        assert len(ids) == len(set(ids))

    def test_block_location_readable(self, processor, sample_pptx):
        """Blocks have human-readable location"""
        blocks = list(processor.extract_text_blocks(sample_pptx))

        assert "Slide 1" in blocks[0].location

    def test_empty_file_returns_no_blocks(self, processor, empty_pptx):
        """Empty file yields no blocks"""
        blocks = list(processor.extract_text_blocks(empty_pptx))
        assert blocks == []


# --- Tests: apply_translations ---

class TestPptxProcessorApplyTranslations:
    """Test PptxProcessor.apply_translations()"""

    def test_applies_translations_to_shapes(self, processor, sample_pptx, tmp_path):
        """Applies translations to shapes"""
        output_path = tmp_path / "output.pptx"

        blocks = list(processor.extract_text_blocks(sample_pptx))
        block_id = blocks[0].id

        translations = {
            block_id: "Slide 1 Text Translated",
        }

        processor.apply_translations(
            sample_pptx, output_path, translations, "jp_to_en"
        )

        # Verify output
        prs = Presentation(output_path)
        slide = prs.slides[0]

        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text:
                        texts.append(para.text)

        assert "Slide 1 Text Translated" in texts

    def test_applies_translations_to_tables(self, processor, pptx_with_table, tmp_path):
        """Applies translations to table cells"""
        output_path = tmp_path / "output.pptx"

        translations = {
            "s0_tbl0_r0_c0": "Header 1 EN",
            "s0_tbl0_r1_c1": "Data EN",
        }

        processor.apply_translations(
            pptx_with_table, output_path, translations, "jp_to_en"
        )

        # Verify output
        prs = Presentation(output_path)
        slide = prs.slides[0]

        table_shape = next(s for s in slide.shapes if s.has_table)
        table = table_shape.table

        assert table.cell(0, 0).text == "Header 1 EN"
        assert table.cell(1, 1).text == "Data EN"
        # Unchanged
        assert table.cell(0, 1).text == "Header 2"

    def test_preserves_untranslated_content(self, processor, pptx_with_multiple_slides, tmp_path):
        """Content not in translations dict is unchanged"""
        output_path = tmp_path / "output.pptx"

        blocks = list(processor.extract_text_blocks(pptx_with_multiple_slides))
        block_0_id = blocks[0].id

        translations = {
            block_0_id: "Translated Slide 1",
            # Slide 2 not translated
        }

        processor.apply_translations(
            pptx_with_multiple_slides, output_path, translations, "jp_to_en"
        )

        prs = Presentation(output_path)

        # Get texts from all slides
        all_texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text:
                            all_texts.append(para.text)

        assert "Translated Slide 1" in all_texts
        assert "スライド2のテキスト" in all_texts  # unchanged (Japanese)

    def test_creates_output_file(self, processor, sample_pptx, tmp_path):
        """Output file is created"""
        output_path = tmp_path / "output.pptx"

        processor.apply_translations(
            sample_pptx, output_path, {}, "jp_to_en"
        )

        assert output_path.exists()


# --- Tests: Edge cases ---

class TestPptxProcessorEdgeCases:
    """Test edge cases"""

    def test_multiple_paragraphs_in_shape(self, processor, tmp_path):
        """Handles shapes with multiple paragraphs"""
        file_path = tmp_path / "multi_para.pptx"
        prs = Presentation()

        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
        tf = txBox.text_frame
        tf.text = "最初の段落"  # Japanese
        p = tf.add_paragraph()
        p.text = "二番目の段落"  # Japanese

        prs.save(file_path)

        blocks = list(processor.extract_text_blocks(file_path))
        assert len(blocks) == 2

    def test_empty_shapes_skipped(self, processor, tmp_path):
        """Empty shapes don't create blocks"""
        file_path = tmp_path / "empty_shape.pptx"
        prs = Presentation()

        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # Empty text box
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.text = ""

        # Shape with whitespace only
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
        txBox2.text_frame.text = "   "

        prs.save(file_path)

        blocks = list(processor.extract_text_blocks(file_path))
        assert len(blocks) == 0

    def test_unicode_content(self, processor, tmp_path):
        """Handles various Unicode characters"""
        file_path = tmp_path / "unicode.pptx"
        prs = Presentation()

        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # Only Japanese/CJK texts are extracted
        texts = ["日本語テスト", "中文测试", "한국어 테스트", "Emoji 😀🎉"]
        for i, text in enumerate(texts):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(5), Inches(0.5))
            txBox.text_frame.text = text

        prs.save(file_path)

        blocks = list(processor.extract_text_blocks(file_path))
        # Japanese and Chinese (CJK) are extracted, Korean and English-only are skipped
        assert len(blocks) == 2

    def test_many_slides(self, processor, tmp_path):
        """Handles presentations with many slides"""
        file_path = tmp_path / "many_slides.pptx"
        prs = Presentation()

        slide_layout = prs.slide_layouts[5]

        for i in range(10):
            slide = prs.slides.add_slide(slide_layout)
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            txBox.text_frame.text = f"スライド{i+1}の内容"  # Japanese text

        prs.save(file_path)

        blocks = list(processor.extract_text_blocks(file_path))
        assert len(blocks) == 10

        info = processor.get_file_info(file_path)
        assert info.slide_count == 10


# --- Tests: create_bilingual_presentation ---

class TestPptxProcessorCreateBilingualPresentation:
    """Test PptxProcessor.create_bilingual_presentation()"""

    def test_creates_bilingual_presentation(self, processor, tmp_path):
        """Creates presentation with interleaved original and translated slides"""
        # Create original presentation
        original_path = tmp_path / "original.pptx"
        prs_orig = Presentation()
        slide_layout = prs_orig.slide_layouts[5]

        slide1 = prs_orig.slides.add_slide(slide_layout)
        txBox1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox1.text_frame.text = "日本語スライド1"

        slide2 = prs_orig.slides.add_slide(slide_layout)
        txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox2.text_frame.text = "日本語スライド2"

        prs_orig.save(original_path)

        # Create translated presentation
        translated_path = tmp_path / "translated.pptx"
        prs_trans = Presentation()
        slide_layout = prs_trans.slide_layouts[5]

        slide1_t = prs_trans.slides.add_slide(slide_layout)
        txBox1_t = slide1_t.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox1_t.text_frame.text = "Japanese Slide 1"

        slide2_t = prs_trans.slides.add_slide(slide_layout)
        txBox2_t = slide2_t.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox2_t.text_frame.text = "Japanese Slide 2"

        prs_trans.save(translated_path)

        # Create bilingual presentation
        output_path = tmp_path / "bilingual.pptx"
        result = processor.create_bilingual_presentation(
            original_path, translated_path, output_path
        )

        # Verify result
        assert result["original_slides"] == 2
        assert result["translated_slides"] == 2
        # Total should include merged slides
        assert result["total_slides"] >= 2

        # Verify output file exists and can be loaded
        assert output_path.exists()
        prs_out = Presentation(output_path)
        assert len(prs_out.slides) >= 2

    def test_handles_single_slide(self, processor, sample_pptx, tmp_path):
        """Works with single-slide presentations"""
        # Create translated version
        translated_path = tmp_path / "translated.pptx"
        prs_trans = Presentation()
        slide_layout = prs_trans.slide_layouts[5]
        slide = prs_trans.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.text = "Slide 1 Text Translated"
        prs_trans.save(translated_path)

        output_path = tmp_path / "bilingual.pptx"
        result = processor.create_bilingual_presentation(
            sample_pptx, translated_path, output_path
        )

        assert result["original_slides"] == 1
        assert result["translated_slides"] == 1
        assert output_path.exists()

    def test_handles_empty_presentations(self, processor, empty_pptx, tmp_path):
        """Handles empty presentations gracefully"""
        # Create empty translated
        translated_path = tmp_path / "translated.pptx"
        prs_trans = Presentation()
        prs_trans.save(translated_path)

        output_path = tmp_path / "bilingual.pptx"
        result = processor.create_bilingual_presentation(
            empty_pptx, translated_path, output_path
        )

        assert result["original_slides"] == 0
        assert result["translated_slides"] == 0
        assert output_path.exists()

    def test_returns_slide_counts(self, processor, pptx_with_multiple_slides, tmp_path):
        """Returns correct slide counts"""
        # Create translated with same number of slides
        translated_path = tmp_path / "translated.pptx"
        prs_trans = Presentation()
        slide_layout = prs_trans.slide_layouts[5]

        for i in range(2):
            slide = prs_trans.slides.add_slide(slide_layout)
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            txBox.text_frame.text = f"Translated Slide {i+1}"

        prs_trans.save(translated_path)

        output_path = tmp_path / "bilingual.pptx"
        result = processor.create_bilingual_presentation(
            pptx_with_multiple_slides, translated_path, output_path
        )

        assert "original_slides" in result
        assert "translated_slides" in result
        assert "total_slides" in result
        assert result["original_slides"] == 2
        assert result["translated_slides"] == 2
